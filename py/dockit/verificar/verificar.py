#!/usr/bin/env python3
"""Comprueba una entrega contra su BRIEF.md. Responde "¿esto está listo?".

    python3 _taller/bin/verificar.py <carpeta-del-trabajo> [--rapido]

Todo lo que se puede comprobar solo, se comprueba solo (existencia, que el
archivo abra, páginas/diapositivas mínimas, fuentes obligatorias, rastro de
citas). Lo que solo puede confirmar una persona se lista aparte, sin darlo
por hecho.

Sale con código 1 si algo automático falla.
"""
from __future__ import annotations

import argparse
import re
import subprocess
import sys
import tempfile
from pathlib import Path

import yaml

OK, MAL, DUDA = "  ok  ", " FALTA", "  ?   "


# ── lectura del brief ─────────────────────────────────────────────────────

def leer_brief(carpeta: Path) -> dict:
    brief = carpeta / "BRIEF.md"
    if not brief.exists():
        sys.exit(f"no hay BRIEF.md en {carpeta} — créalo con _taller/bin/nuevo.py")
    texto = brief.read_text(encoding="utf-8")
    m = re.match(r"^---\n(.*?)\n---\n", texto, re.S)
    if not m:
        sys.exit(f"{brief} no tiene cabecera YAML entre --- y ---")
    datos = yaml.safe_load(m.group(1)) or {}
    datos["_cuerpo"] = texto[m.end():]
    return datos


# ── extracción de texto por tipo ──────────────────────────────────────────

def texto_docx(ruta: Path) -> str:
    from docx import Document
    doc = Document(str(ruta))
    trozos = [p.text for p in doc.paragraphs]
    for t in doc.tables:
        for fila in t.rows:
            trozos += [c.text for c in fila.cells]
    return "\n".join(trozos)


def texto_pptx(ruta: Path) -> tuple[str, int]:
    from pptx import Presentation
    pres = Presentation(str(ruta))
    trozos = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                trozos.append(shape.text_frame.text)
    return "\n".join(trozos), len(pres.slides)


def texto_pdf(ruta: Path) -> tuple[str, int]:
    from pypdf import PdfReader
    lector = PdfReader(str(ruta))
    trozos = [(pg.extract_text() or "") for pg in lector.pages]
    return "\n".join(trozos), len(lector.pages)


def paginas_via_libreoffice(ruta: Path) -> int | None:
    """Cuenta páginas reales de un .docx convirtiéndolo. Lento pero exacto."""
    with tempfile.TemporaryDirectory() as tmp:
        try:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", tmp, str(ruta)],
                check=True, capture_output=True, timeout=180,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired,
                FileNotFoundError):
            return None
        pdfs = list(Path(tmp).glob("*.pdf"))
        if not pdfs:
            return None
        from pypdf import PdfReader
        return len(PdfReader(str(pdfs[0])).pages)


# ── comprobaciones ────────────────────────────────────────────────────────

class Informe:
    def __init__(self) -> None:
        self.fallos = 0
        self.lineas: list[str] = []

    def di(self, estado: str, msg: str) -> None:
        if estado is MAL:
            self.fallos += 1
        self.lineas.append(f"[{estado}] {msg}")

    def titulo(self, t: str) -> None:
        self.lineas.append(f"\n{t}")


def revisar_entregables(carpeta: Path, brief: dict, inf: Informe,
                        rapido: bool) -> str:
    """Comprueba cada entregable y devuelve todo su texto concatenado."""
    inf.titulo("Entregables")
    entregables = brief.get("entregables") or []
    if not entregables:
        inf.di(DUDA, "el brief no declara ningún entregable")
        return ""

    todo_el_texto = []
    for ent in entregables:
        rel = ent.get("archivo")
        if not rel:
            continue
        ruta = carpeta / rel
        tipo = (ent.get("tipo") or ruta.suffix.lstrip(".")).lower()

        if not ruta.exists():
            inf.di(MAL, f"{rel} — no existe")
            continue
        if ruta.stat().st_size == 0:
            inf.di(MAL, f"{rel} — está vacío")
            continue

        try:
            if tipo == "pptx":
                texto, n = texto_pptx(ruta)
                unidad, minimo = "diapositivas", ent.get("minimo_diapositivas")
            elif tipo == "pdf":
                texto, n = texto_pdf(ruta)
                unidad, minimo = "páginas", ent.get("minimo_paginas")
            elif tipo == "docx":
                texto = texto_docx(ruta)
                minimo = ent.get("minimo_paginas")
                unidad = "páginas"
                n = None
                if minimo and not rapido:
                    n = paginas_via_libreoffice(ruta)
            else:
                texto = ruta.read_text(encoding="utf-8", errors="replace")
                n, unidad, minimo = None, "", None
        except Exception as e:                     # archivo corrupto
            inf.di(MAL, f"{rel} — no abre: {type(e).__name__}: {e}")
            continue

        todo_el_texto.append(texto)
        detalle = f"{n} {unidad}" if n is not None else f"{len(texto.split())} palabras"
        if minimo and n is not None and n < int(minimo):
            inf.di(MAL, f"{rel} — {detalle}, el brief pide {minimo}")
        elif minimo and n is None:
            inf.di(DUDA, f"{rel} — {detalle}; no pude contar {unidad} (usa sin --rapido)")
        else:
            inf.di(OK, f"{rel} — {detalle}")

    return "\n".join(todo_el_texto)


def revisar_fuentes(carpeta: Path, brief: dict, inf: Informe) -> None:
    fuentes = brief.get("fuentes") or {}
    obligatorias = fuentes.get("obligatorias") or []
    if not obligatorias:
        return
    inf.titulo("Fuentes obligatorias")
    for rel in obligatorias:
        ruta = carpeta / rel
        inf.di(OK if ruta.exists() else MAL,
               f"{rel}{'' if ruta.exists() else ' — no está en la carpeta'}")


# Reconoce las citas en el texto tal y como las escribe citeproc:
#   (Nath et al., 2024) · (Ćosić & Bača, 2010) · (Pérez-Gómez, 2019a)
# Ojo con dos detalles que costaron un fallo silencioso: «et al.» lleva espacio
# duro (U+00A0) y APA usa «&» para dos autores. Sin contemplarlos, el
# verificador contaba cero citas en documentos que sí las tenían.
RE_CITA = re.compile(
    r"\("
    r"[A-ZÁÉÍÓÚÑÀÈÌÒÙÄËÏÖÜÑĆČŠŽ]"      # el apellido empieza en mayúscula
    r"[^()]{1,80}?"                     # autores, sin paréntesis anidados
    r",\s*"
    r"(?:\d{4}[a-z]?|s\.\s*f\.)"        # año, o «s. f.» si no lo tiene
    r"\)"
)


def normalizar_espacios(s: str) -> str:
    """Sustituye los espacios duros que mete citeproc por espacios normales."""
    return (s or "").replace("\u00a0", " ").replace("\u202f", " ")
RE_BIBLIO = re.compile(r"^\s*(bibliograf|referencias|fuentes consultadas)",
                       re.I | re.M)


def revisar_citas(texto: str, brief: dict, inf: Informe) -> None:
    modo = ((brief.get("fuentes") or {}).get("citas") or "no").lower()
    if modo == "no" or not texto:
        return
    inf.titulo("Citas")
    plano = normalizar_espacios(texto)
    citas = RE_CITA.findall(plano)
    tiene_biblio = bool(RE_BIBLIO.search(plano))
    estado = OK if citas else (MAL if modo == "obligatorias" else DUDA)
    inf.di(estado, f"{len(citas)} citas con formato (Autor, año)")
    inf.di(OK if tiene_biblio else (MAL if modo == "obligatorias" else DUDA),
           "sección de bibliografía/referencias"
           + ("" if tiene_biblio else " — no la encuentro"))


def revisar_prohibido(texto: str, brief: dict, inf: Informe) -> None:
    """No puede comprobarse solo: se recuerda para la revisión humana."""
    prohibido = brief.get("prohibido") or []
    imgs = (brief.get("imagenes") or {}).get("prohibido") or []
    if not prohibido and not imgs:
        return
    inf.titulo("Líneas rojas (compruébalas tú)")
    for regla in list(prohibido) + [f"en imágenes: {r}" for r in imgs]:
        inf.di(DUDA, str(regla))

    nota = (brief.get("imagenes") or {}).get("nota_obligatoria")
    if nota and texto:
        presente = nota.lower() in texto.lower()
        inf.di(OK if presente else MAL,
               f"leyenda obligatoria «{nota}»"
               + ("" if presente else " — no aparece en el texto"))


def revisar_investigacion(carpeta: Path, brief: dict, inf: Informe) -> None:
    """Biblioteca y afirmaciones: lo que hace que el contenido sea verificable."""
    biblioteca = carpeta / "fuentes" / "biblioteca.json"
    afirmaciones = carpeta / "afirmaciones.json"
    exige_citas = ((brief.get("fuentes") or {}).get("citas") or "").lower() == "obligatorias"
    if not biblioteca.exists() and not afirmaciones.exists() and not exige_citas:
        return

    inf.titulo("Investigación")

    if biblioteca.exists():
        import json as _json
        refs = _json.loads(biblioteca.read_text())
        con_doi = sum(1 for e in refs if e.get("DOI"))
        inf.di(OK if refs else MAL,
               f"biblioteca: {len(refs)} referencias, {con_doi} con DOI")
        if refs and con_doi < len(refs):
            inf.di(DUDA, f"{len(refs) - con_doi} sin DOI — verifícalas a mano "
                         f"(bibliografia.py --verificar)")
    elif exige_citas:
        inf.di(MAL, "el brief exige citas y no hay fuentes/biblioteca.json")

    if afirmaciones.exists():
        r = subprocess.run(
            [sys.executable, str(Path(__file__).with_name("afirmaciones.py")),
             "--carpeta", str(carpeta)],
            capture_output=True, text=True, timeout=600)
        resumen = [l for l in r.stdout.splitlines() if "sin respaldo" in l]
        detalle = resumen[0].strip() if resumen else "no pude resumir"
        inf.di(OK if r.returncode == 0 else MAL, f"afirmaciones: {detalle}")
        if r.returncode != 0:
            for linea in r.stdout.splitlines():
                if linea.startswith("[") and "ok" not in linea[:14]:
                    inf.lineas.append(f"         {linea}")
    elif exige_citas:
        inf.di(DUDA, "no hay afirmaciones.json — sin él, las cifras del texto "
                     "no están respaldadas por nada comprobable")


def revisar_terminado(brief: dict, inf: Informe) -> None:
    items = brief.get("terminado") or []
    if not items:
        return
    inf.titulo("Criterio de terminado (confirmación tuya)")
    for it in items:
        inf.di(DUDA, str(it))


def main() -> int:
    p = argparse.ArgumentParser()
    p.add_argument("carpeta", type=Path)
    p.add_argument("--rapido", action="store_true",
                   help="no convierte docx a PDF para contar páginas")
    args = p.parse_args()

    carpeta = args.carpeta.resolve()
    brief = leer_brief(carpeta)
    inf = Informe()

    print(f"\n{brief.get('titulo') or carpeta.name}")
    if brief.get("entrega"):
        print(f"entrega: {brief['entrega']}")

    texto = revisar_entregables(carpeta, brief, inf, args.rapido)
    revisar_fuentes(carpeta, brief, inf)
    revisar_citas(texto, brief, inf)
    revisar_prohibido(texto, brief, inf)
    revisar_investigacion(carpeta, brief, inf)
    revisar_terminado(brief, inf)

    print("\n".join(inf.lineas))
    pendientes = sum(1 for l in inf.lineas if DUDA in l)
    print(f"\n{inf.fallos} fallo(s) automático(s), {pendientes} por confirmar")
    if inf.fallos:
        print("NO está listo.")
    return 1 if inf.fallos else 0


if __name__ == "__main__":
    raise SystemExit(main())
