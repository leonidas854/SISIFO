#!/usr/bin/env python3
"""Extrae las unidades visuales que deben recibir una fotografía.

Una unidad visual es el título de la diapositiva, un subtítulo o un apartado
destacado. La detección combina el título semántico ya auditado con el formato
real del PPTX; el resultado se conserva en JSON para poder revisarlo antes de
generar ninguna imagen.
"""

from __future__ import annotations

import argparse
import json
import re
import unicodedata
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]


def normalizar(texto: str) -> str:
    texto = texto.replace("\x0b", " ").replace("\n", " ")
    return re.sub(r"\s+", " ", texto).strip(" \t\r\n-–•")


def clave(texto: str) -> str:
    texto = unicodedata.normalize("NFKD", texto.casefold())
    texto = "".join(c for c in texto if not unicodedata.combining(c))
    return re.sub(r"[^a-z0-9]+", " ", texto).strip()


def palabras(texto: str) -> int:
    return len(re.findall(r"\w+", texto, flags=re.UNICODE))


def es_rotulo_corto(texto: str, tamanos: list[float]) -> bool:
    """Reconoce rótulos visuales que heredan la negrita desde el tema."""
    if not texto or palabras(texto) > 12 or len(texto) > 105:
        return False
    if texto.endswith((".", ";", ",")):
        return False
    if not tamanos or max(tamanos) < 18:
        return False
    # Las frases verbales largas suelen ser cuerpo aunque no terminen en punto.
    inicio = clave(texto).split(" ", 1)[0]
    if palabras(texto) >= 8 and inicio in {
        "aplicar", "brindar", "coordinar", "determinar", "establecer",
        "evitar", "garantizar", "hacer", "identificar", "implementar",
        "incrementar", "llevar", "mantener", "permitir", "promover",
        "realizar", "reconocer", "reducir", "revisar", "se", "uso",
        "utilizar", "verificar",
    }:
        return False
    return True


def candidatos_diapositiva(slide, titulo: str) -> list[dict]:
    unidades: list[dict] = [
        {"rotulo": normalizar(titulo), "tipo": "titulo", "fuente": "plan"}
    ]
    vistos = {clave(titulo)}

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for par in shape.text_frame.paragraphs:
            runs = [run for run in par.runs if normalizar(run.text)]
            texto = normalizar("".join(run.text for run in par.runs))
            if not texto:
                continue
            tamanos = [run.font.size.pt for run in runs if run.font.size]
            runs_negrita = [normalizar(run.text) for run in runs if run.font.bold is True]
            todo_negrita = bool(runs) and all(run.font.bold is True for run in runs)

            piezas: list[tuple[str, str, str]] = []
            if todo_negrita:
                # Dos runs contiguos pueden ser título y subtítulo distintos.
                if len(runs_negrita) > 1 and all(palabras(x) >= 2 for x in runs_negrita):
                    piezas.extend((x, "destacado", "negrita_run") for x in runs_negrita)
                else:
                    piezas.append((texto, "destacado", "negrita"))
            elif runs_negrita:
                piezas.extend((x, "destacado", "negrita_run") for x in runs_negrita)
            elif es_rotulo_corto(texto, tamanos):
                piezas.append((texto, "apartado", "rotulo_visual"))

            for pieza, tipo, fuente in piezas:
                pieza = normalizar(pieza)
                k = clave(pieza)
                if not k or k in vistos:
                    continue
                # Evita duplicar el título cuando una versión contiene a la otra.
                kt = clave(titulo)
                if (k in kt or kt in k) and min(len(k), len(kt)) >= 12:
                    continue
                vistos.add(k)
                unidades.append({"rotulo": pieza, "tipo": tipo, "fuente": fuente})
    return unidades


def extraer_tema(tema: int, pptx: Path, plan: Path) -> dict:
    data = json.loads(plan.read_text(encoding="utf-8"))
    titulos = {}
    for trabajo in data["trabajos"]:
        titulos.setdefault(int(trabajo["diapositiva"]), trabajo["titulo"])

    prs = Presentation(pptx)
    slides = list(prs.slides)
    diapositivas = []
    for numero in sorted(titulos):
        unidades = candidatos_diapositiva(slides[numero - 1], titulos[numero])
        diapositivas.append(
            {"diapositiva": numero, "titulo": titulos[numero], "unidades": unidades}
        )
    return {
        "tema": tema,
        "titulo_tema": data.get("titulo_tema", ""),
        "archivo_original": str(pptx),
        "diapositivas": diapositivas,
        "total_unidades": sum(len(d["unidades"]) for d in diapositivas),
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--originales", type=Path, default=ROOT.parent / "diapos_original")
    parser.add_argument("--planes", type=Path, default=ROOT / "prompts")
    parser.add_argument("--salida", type=Path, default=ROOT / "paneles" / "inventario_paneles.json")
    args = parser.parse_args()

    temas = []
    for tema in range(3, 15):
        candidatos = sorted(args.originales.glob(f"{tema}_*.pptx"))
        if len(candidatos) != 1:
            raise SystemExit(f"Tema {tema}: se esperaba un PPTX y se encontraron {len(candidatos)}")
        temas.append(
            extraer_tema(
                tema,
                candidatos[0],
                args.planes / f"tema{tema:02d}_expandido.json",
            )
        )

    payload = {
        "criterio": "titulo + subtitulos/apartados destacados en el PPTX",
        "temas": temas,
        "total_unidades": sum(t["total_unidades"] for t in temas),
    }
    args.salida.parent.mkdir(parents=True, exist_ok=True)
    args.salida.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps({"salida": str(args.salida), "total": payload["total_unidades"],
                      "por_tema": {t["tema"]: t["total_unidades"] for t in temas}},
                     ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
