#!/usr/bin/env python3
"""Inventaría texto, notas y recursos visuales de presentaciones PPTX.

El script es deliberadamente de solo lectura sobre los PPTX.  Produce un informe
Markdown pensado para revisar cada diapositiva sin perder el orden visual básico
(de arriba hacia abajo y de izquierda a derecha).

Uso:
    python3 extraer_texto_pptx.py diapos_original --output inventario.md
    python3 extraer_texto_pptx.py archivo.pptx
"""

from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Iterator

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Emu


EMU_PER_CM = 360_000
DIAGRAM_TYPES = {
    MSO_SHAPE_TYPE.CHART,
    MSO_SHAPE_TYPE.DIAGRAM,
    MSO_SHAPE_TYPE.IGX_GRAPHIC,
}
PICTURE_TYPES = {
    MSO_SHAPE_TYPE.PICTURE,
    MSO_SHAPE_TYPE.LINKED_PICTURE,
}


@dataclass(frozen=True)
class TextItem:
    top: int
    left: int
    z_order: int
    kind: str
    name: str
    text: str


@dataclass(frozen=True)
class VisualItem:
    kind: str
    name: str
    left: int
    top: int
    width: int
    height: int
    detail: str = ""


def natural_key(path: Path) -> tuple:
    """Orden natural: 3 antes que 10."""
    return tuple(
        int(part) if part.isdigit() else part.casefold()
        for part in re.split(r"(\d+)", path.name)
    )


def cm(value: int | Emu) -> float:
    return int(value) / EMU_PER_CM


def clean_text(value: str | None) -> str:
    if not value:
        return ""
    lines = []
    for line in value.replace("\x0b", "\n").splitlines():
        line = re.sub(r"[ \t]+", " ", line).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def shape_kind(shape) -> str:
    if getattr(shape, "is_placeholder", False):
        try:
            placeholder_name = shape.placeholder_format.type.name
        except (AttributeError, ValueError):
            placeholder_name = "PLACEHOLDER"
        return f"marcador:{placeholder_name.lower()}"
    try:
        return shape.shape_type.name.lower()
    except AttributeError:
        return str(shape.shape_type)


def shape_text(shape) -> str:
    """Extrae el texto visible de una forma sin duplicar tablas/gráficos."""
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            cells = [clean_text(cell.text).replace("\n", " / ") for cell in row.cells]
            rows.append(" | ".join(cells))
        return clean_text("\n".join(rows))

    if getattr(shape, "has_chart", False):
        chart = shape.chart
        if getattr(chart, "has_title", False):
            return clean_text(chart.chart_title.text_frame.text)

    if getattr(shape, "has_text_frame", False):
        return clean_text(shape.text_frame.text)

    # SmartArt y algunos objetos gráficos no están plenamente expuestos por
    # python-pptx. Sus nodos a:t sí conservan las cadenas mostradas.
    try:
        xml_strings = [clean_text(node.text) for node in shape.element.xpath(".//a:t")]
    except (AttributeError, TypeError):
        xml_strings = []
    return clean_text("\n".join(text for text in xml_strings if text))


def iter_shapes(shapes, prefix: str = "") -> Iterator:
    """Recorre formas y subformas agrupadas manteniendo su orden Z."""
    for index, shape in enumerate(shapes, start=1):
        path = f"{prefix}{index}" if not prefix else f"{prefix}.{index}"
        yield path, index, shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes, path)


def text_items(slide) -> list[TextItem]:
    items: list[TextItem] = []
    for path, z_order, shape in iter_shapes(slide.shapes):
        # El contenedor de grupo suele repetir los textos de sus hijos.
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            continue
        text = shape_text(shape)
        if not text:
            continue
        items.append(
            TextItem(
                top=int(shape.top),
                left=int(shape.left),
                z_order=z_order,
                kind=shape_kind(shape),
                name=f"{path}: {shape.name}",
                text=text,
            )
        )
    # PowerPoint almacena formas en orden Z, que no siempre coincide con el
    # orden de lectura. El informe usa posición visual y conserva Z como
    # desempate para que el resultado sea estable.
    return sorted(items, key=lambda item: (item.top, item.left, item.z_order))


def title_text(slide, items: list[TextItem]) -> tuple[str, bool]:
    """Devuelve título y si procede de un marcador de título real."""
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            placeholder_type = shape.placeholder_format.type
        except (AttributeError, ValueError):
            continue
        if placeholder_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            title = clean_text(getattr(shape, "text", ""))
            if title:
                return title, True
    if items:
        return items[0].text.splitlines()[0], False
    return "(sin título)", False


def notes_text(slide) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    try:
        return clean_text(slide.notes_slide.notes_text_frame.text)
    except (AttributeError, ValueError):
        return ""


def picture_detail(shape) -> str:
    details: list[str] = []
    try:
        image = shape.image
        details.append(f"archivo `{image.filename}`")
        details.append(f"{image.size[0]}×{image.size[1]} px")
        details.append(str(image.ext).upper())
    except (AttributeError, ValueError, KeyError):
        pass
    return ", ".join(details)


def visual_items(slide) -> list[VisualItem]:
    visuals: list[VisualItem] = []
    for path, _z_order, shape in iter_shapes(slide.shapes):
        shape_type = shape.shape_type
        name = f"{path}: {shape.name}"
        if shape_type in PICTURE_TYPES:
            visuals.append(
                VisualItem(
                    "imagen incrustada" if shape_type == MSO_SHAPE_TYPE.PICTURE else "imagen vinculada",
                    name,
                    int(shape.left),
                    int(shape.top),
                    int(shape.width),
                    int(shape.height),
                    picture_detail(shape),
                )
            )
        elif shape_type in DIAGRAM_TYPES:
            visuals.append(
                VisualItem(
                    "gráfico/diagrama",
                    name,
                    int(shape.left),
                    int(shape.top),
                    int(shape.width),
                    int(shape.height),
                    shape_type.name,
                )
            )
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            # El grupo solo se considera diagrama si está compuesto por varias
            # piezas visuales nativas. Los hijos siguen inventariados aparte.
            non_text_children = sum(
                1
                for child in shape.shapes
                if child.shape_type
                in {
                    MSO_SHAPE_TYPE.AUTO_SHAPE,
                    MSO_SHAPE_TYPE.FREEFORM,
                    MSO_SHAPE_TYPE.LINE,
                    MSO_SHAPE_TYPE.CANVAS,
                }
            )
            if non_text_children >= 2:
                visuals.append(
                    VisualItem(
                        "composición vectorial agrupada",
                        name,
                        int(shape.left),
                        int(shape.top),
                        int(shape.width),
                        int(shape.height),
                        f"{non_text_children} piezas visuales",
                    )
                )
    return visuals


def size_description(item: VisualItem, slide_width: int, slide_height: int) -> str:
    area_ratio = (item.width * item.height) / (slide_width * slide_height)
    if area_ratio >= 0.45:
        label = "muy grande"
    elif area_ratio >= 0.20:
        label = "grande"
    elif area_ratio >= 0.06:
        label = "mediana"
    else:
        label = "pequeña"
    return (
        f"{cm(item.width):.1f}×{cm(item.height):.1f} cm; "
        f"≈{area_ratio * 100:.0f}% del área, {label}"
    )


def visual_status(
    visuals: list[VisualItem], slide_width: int, slide_height: int, all_text: str
) -> str:
    pictures = [item for item in visuals if item.kind.startswith("imagen")]
    diagrams = [item for item in visuals if "diagrama" in item.kind or "vectorial" in item.kind]
    explicit_request = bool(
        re.search(
            r"\b(imagen|foto|fotograf[ií]a|ilustraci[oó]n|gr[aá]fico)\b",
            all_text,
            flags=re.IGNORECASE,
        )
    )
    largest_picture = max(
        ((item.width * item.height) / (slide_width * slide_height) for item in pictures),
        default=0.0,
    )
    if largest_picture >= 0.20:
        return "Ya contiene al menos una imagen principal grande; revisar como posible ejemplo resuelto."
    if diagrams:
        return "Ya contiene un gráfico/diagrama nativo o una composición vectorial relevante."
    if pictures:
        return "Solo contiene imagen(es) pequeña(s); probablemente aún necesita una imagen principal."
    if explicit_request:
        return "No contiene imágenes insertadas y el texto menciona un recurso visual: requiere imagen."
    return "No contiene fotos ni diagramas detectables en la diapositiva: requiere imagen o revisión manual."


def markdown_escape(value: str) -> str:
    return value.replace("|", "\\|")


def render_report(paths: list[Path]) -> str:
    lines = [
        "# Inventario de las presentaciones originales",
        "",
        "Informe generado a partir de los archivos PPTX, sin modificarlos. El texto se ordena "
        "por posición visual (arriba→abajo, izquierda→derecha). Las imágenes heredadas del "
        "patrón/tema pueden no aparecer como formas propias de la diapositiva.",
        "",
        "## Resumen",
        "",
        "| Tema/archivo | Diapositivas | Con imagen grande | Sin foto/diagrama |",
        "|---|---:|---:|---:|",
    ]
    presentations = []
    for path in paths:
        prs = Presentation(path)
        slide_rows = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            items = text_items(slide)
            visuals = visual_items(slide)
            combined_text = "\n".join(item.text for item in items)
            slide_rows.append((slide_index, slide, items, visuals, combined_text))
        with_large = sum(
            1
            for _, _, _, visuals, _ in slide_rows
            if any(
                item.kind.startswith("imagen")
                and (item.width * item.height) / (prs.slide_width * prs.slide_height) >= 0.20
                for item in visuals
            )
        )
        without_visual = sum(1 for _, _, _, visuals, _ in slide_rows if not visuals)
        lines.append(
            f"| `{markdown_escape(path.name)}` | {len(prs.slides)} | {with_large} | {without_visual} |"
        )
        presentations.append((path, prs, slide_rows))

    for path, prs, slide_rows in presentations:
        lines.extend(
            [
                "",
                f"## {path.name}",
                "",
                f"- Total: **{len(prs.slides)} diapositiva(s)**.",
                f"- Formato: **{cm(prs.slide_width):.1f}×{cm(prs.slide_height):.1f} cm**.",
            ]
        )
        for slide_index, slide, items, visuals, combined_text in slide_rows:
            title, true_title = title_text(slide, items)
            notes = notes_text(slide)
            pictures = [item for item in visuals if item.kind.startswith("imagen")]
            diagrams = [
                item for item in visuals if "diagrama" in item.kind or "vectorial" in item.kind
            ]
            lines.extend(
                [
                    "",
                    f"### Diapositiva {slide_index}: {title}",
                    "",
                    f"- Título: **{title}**"
                    + (" (marcador de título)" if true_title else " (inferido por posición)"),
                    f"- Fotos/imágenes insertadas: **{len(pictures)}**.",
                    f"- Diagramas/gráficos detectados: **{len(diagrams)}**.",
                    f"- Diagnóstico automático: {visual_status(visuals, prs.slide_width, prs.slide_height, combined_text)}",
                    "",
                    "#### Texto completo en orden visual",
                    "",
                ]
            )
            if items:
                for order, item in enumerate(items, start=1):
                    position = f"x={cm(item.left):.1f}, y={cm(item.top):.1f} cm"
                    lines.append(
                        f"{order}. `{item.kind}` — {item.name} ({position})"
                    )
                    for text_line in item.text.splitlines():
                        lines.append(f"   - {text_line}")
            else:
                lines.append("(No se detectó texto en formas de la diapositiva.)")

            lines.extend(["", "#### Notas del presentador", ""])
            if notes:
                for note_line in notes.splitlines():
                    lines.append(f"- {note_line}")
            else:
                lines.append("(Sin notas del presentador.)")

            lines.extend(["", "#### Recursos visuales existentes", ""])
            if visuals:
                for item in visuals:
                    position = f"x={cm(item.left):.1f}, y={cm(item.top):.1f} cm"
                    details = f"; {item.detail}" if item.detail else ""
                    lines.append(
                        f"- **{item.kind}** — {item.name}: "
                        f"{size_description(item, prs.slide_width, prs.slide_height)}; {position}{details}."
                    )
            else:
                lines.append("- Ninguna foto, gráfico o diagrama insertado detectado.")

    lines.extend(
        [
            "",
            "## Criterio técnico",
            "",
            "- Una imagen ocupa el área de su marco, no necesariamente el área visible después de recortes.",
            "- “Diagrama” incluye gráficos/SmartArt y grupos vectoriales con varias piezas; no cuenta cada rectángulo decorativo aislado.",
            "- La clasificación automática es una ayuda. La identificación final de ejemplos se debe confirmar visualmente.",
            "",
        ]
    )
    return "\n".join(lines)


def collect_paths(inputs: Iterable[str]) -> list[Path]:
    paths: list[Path] = []
    for raw in inputs:
        path = Path(raw).expanduser()
        if path.is_dir():
            paths.extend(path.glob("*.pptx"))
        elif path.is_file() and path.suffix.casefold() == ".pptx":
            paths.append(path)
        else:
            raise FileNotFoundError(f"No es un PPTX ni un directorio válido: {path}")
    unique_paths = {path.resolve(): path.resolve() for path in paths}
    return sorted(unique_paths.values(), key=natural_key)


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("inputs", nargs="+", help="Archivo(s) PPTX o directorio(s)")
    parser.add_argument("-o", "--output", type=Path, help="Guardar informe Markdown")
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = parse_args(argv or sys.argv[1:])
    try:
        paths = collect_paths(args.inputs)
    except FileNotFoundError as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 2
    if not paths:
        print("Error: no se encontraron archivos .pptx", file=sys.stderr)
        return 2
    report = render_report(paths)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(report, encoding="utf-8")
        print(f"Informe escrito en {args.output} ({len(paths)} PPTX).")
    else:
        print(report)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
