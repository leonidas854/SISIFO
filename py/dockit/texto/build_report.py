from __future__ import annotations

import os

import json
from pathlib import Path

from pptx import Presentation


def _base_proyecto() -> str:
    """Carpeta del proyecto con los .pptx fuente.

    Se fija con TALLER_PROYECTO; si no, se busca subiendo desde donde se
    ejecuta. Nunca una ruta absoluta escrita a mano: el motor no puede depender
    de que exista una carpeta de trabajo concreta, porque se borran.
    """
    if _v := os.environ.get("TALLER_PROYECTO"):
        return str(Path(_v).resolve())
    _aqui = Path.cwd().resolve()
    for _c in [_aqui, *_aqui.parents]:
        if any(_c.glob("*DIAPOSITIVA*.pptx")) or (_c / "diapos_original").is_dir():
            return str(_c)
    return str(_aqui)


BASE = Path(_base_proyecto()) / "diapos_original"
OUT = Path(__file__).resolve().parent

FILES = {
    "11": "11_ DIAPOSITIVA.pptx",
    "12": "12_ DIAPOSITIVA.pptx",
    "13": "13_ DIAPOSITIVA.pptx",
    "14": "14_ DIAPOSITIVA -.pptx",
}

# Índices de forma en orden de lectura humano, validados contra todas las formas
# que contienen texto nativo en cada diapositiva.
ORDER = {
    "11": {
        1: [9, 3, 2, 10, 4],
        2: [1, 4, 2, 7],
        3: [1, 4, 2, 7],
        4: [1, 4, 2, 7, 11, 12, 9],
        5: [1, 2, 5, 4],
        6: [1, 6, 2, 7, 4, 8, 5],
        7: [1, 2],
        8: [1, 5, 2, 6, 4],
        9: [1, 2, 4],
        10: [1, 2, 4],
        11: [1, 2, 4],
        12: [1, 2, 4],
        13: [1, 2],
        14: [1, 5, 2, 4],
        15: [1, 5, 2, 4, 6],
        16: [1, 5, 2, 4],
        17: [1, 5, 2, 4, 6],
        18: [1, 5, 2, 4],
        19: [1, 4, 2],
        20: [7, 2],
    },
    "12": {
        1: [9, 3, 2, 10, 4],
        2: [1, 4, 2, 5, 6],
        3: [1, 5, 4, 8, 7, 2],
        4: [1, 4, 2, 5, 6, 9, 7],
        5: [1, 4, 2, 7],
        6: [1, 4, 2, 5, 8, 12, 9, 13, 10, 14, 11, 15],
        7: [1, 4, 2, 7],
        8: [7, 2, 1, 5],
        9: [1, 4, 6, 2, 7, 8, 11, 10],
        10: [1, 4, 6, 2, 7, 8, 10],
        11: [1, 3, 4, 5, 6, 9, 7, 8, 10],
        12: [1, 4, 2, 5],
        13: [1, 4, 2, 8, 5],
        14: [1, 4, 2, 8, 5, 9],
        15: [1, 4, 2, 7, 6],
        16: [1, 4, 2, 8, 7],
        17: [1, 4, 2, 6],
        18: [1, 4, 2, 7, 9],
        19: [7, 2],
    },
    "13": {
        1: [9, 3, 2, 10, 4],
        2: [1, 4, 5, 2],
        3: [1, 4, 2, 7],
        4: [3, 5, 1, 4],
        5: [1, 3, 4, 5, 6, 7, 8],
        6: [1, 4, 2],
        7: [1, 2, 4, 5],
        8: [1, 2, 4, 5],
        9: [1, 2, 4, 5],
        10: [1, 2, 4, 5],
        11: [1, 2, 4, 8],
        12: [1, 2, 4, 5],
        13: [1, 2, 4, 5],
        14: [1, 2, 4, 5, 9],
        15: [1, 2, 4],
        16: [1, 4, 6],
        17: [1, 3, 5, 4, 6, 8, 7, 9],
        18: [1, 3],
        19: [1, 3, 5, 4],
        20: [1, 3, 5, 4, 6, 7],
        21: [1, 3, 4, 7, 8, 9],
        22: [1, 3, 4, 7, 8, 9],
        23: [1, 3, 4, 7, 8, 9],
        24: [1, 3, 4, 7, 8, 9],
        25: [7, 2],
    },
    "14": {
        1: [8, 10, 2, 9, 3],
        2: [1, 2],
        3: [1, 4, 2, 5, 6, 7, 8, 9],
        4: [1, 4, 2],
        5: [1, 4, 5, 2, 6, 7],
        6: [1, 4, 5, 2, 6, 7],
        7: [1, 4, 5, 2, 6, 7],
        8: [1, 4, 5, 2, 6, 7],
        9: [1, 5, 2, 4],
        10: [1, 5, 2, 4],
        11: [1, 5, 2, 4],
        12: [1, 5, 2, 4],
        13: [1, 5, 2, 4],
        14: [1, 5, 2, 4],
        15: [1, 5, 2, 4],
        16: [1, 5, 2, 4],
        17: [1, 5, 2, 4],
        18: [1, 5, 2, 4],
        19: [1, 5, 2, 4],
        20: [1, 5, 2, 4],
        21: [1, 5, 2, 6, 4],
        22: [1, 2, 4],
        23: [7, 2],
    },
}


LOGO_OCR = {
    "description": "Logotipo institucional bicentenario, repetido en la esquina inferior derecha.",
    "text": ["200", "POLICIA BOLIVIANA", "1826 - 2026"],
    "uncertainty": "El lema diminuto del listón dorado inferior no es resoluble con fidelidad en la imagen fuente de 167 x 103 píxeles.",
}


SPECIAL_IMAGE_13_16 = {
    "description": (
        "Ilustración de un oficial de la Policía Boliviana señalando un mapa digital de despliegue, "
        "con bandera boliviana, otros agentes y vehículos policiales."
    ),
    "text": [
        "PLANIFRULLAJE - Policía Boliviana",
        (
            "Consiste en diseñar y estructurar el despliegue del personal policial en un territorio "
            "determinado y durante un tiempo específico. Su fin principal es disuadir la comisión de "
            "delitos, garantizar la seguridad ciudadana y mantener el orden público."
        ),
        "Safety Shield",
        "SHIFT TIMELINES",
        "1, 2, 3, 4",
        "TACTICAL DISTRICO",
    ],
    "uncertainty": (
        "La ilustración parece generada por IA y contiene grafías anómalas. Se lee con alta confianza "
        "'PLANIFRULLAJE' y 'TACTICAL DISTRICO', aunque probablemente pretendían decir 'plan de patrullaje' "
        "y 'distrito táctico'. La placa del uniforme parece decir 'ESERGRAVER' y el quinto marcador de la "
        "escala es un glifo similar a G/5; ambos son demasiado ambiguos para una transcripción segura."
    ),
}


def normalize_text(text: str) -> str:
    return text.replace("\v", "\n").replace("\r", "\n").strip()


def image_entries(deck: str, slide_no: int, total: int) -> list[dict]:
    entries = [dict(LOGO_OCR)]
    if slide_no in (1, total):
        entries.append(
            {
                "description": "Escudo/insignia de la Policía Boliviana en la portada o cierre.",
                "text": [],
                "uncertainty": "No contiene texto legible adicional.",
            }
        )
    if deck == "13" and slide_no == 15:
        entries.append(
            {
                "description": "Dos iconos idénticos de flecha descendente.",
                "text": [],
                "uncertainty": "Sin texto.",
            }
        )
    if deck == "13" and slide_no == 16:
        entries.append(dict(SPECIAL_IMAGE_13_16))
    return entries


def slide_uncertainties(deck: str, slide_no: int) -> list[str]:
    notes: list[str] = []
    if deck == "12" and slide_no == 10:
        notes.append(
            "El rótulo nativo dice literalmente 'Rechazo de la Denuncia, Querella, actuaciones policiales o su archivo', "
            "pero el texto asociado describe una imputación formal. Se conserva sin corregir."
        )
    if deck == "13" and slide_no == 3:
        notes.append(
            "La primera frase de la definición empieza literalmente con 'Las leyes y los reglamentos institucionales...' "
            "y parece gramaticalmente incompleta. Se conserva sin corregir."
        )
    if deck == "13" and slide_no == 19:
        notes.append(
            "Hay dos fragmentos visibles anómalos en el archivo: una 'T' aislada y el prefijo 'A•' antes de "
            "'Tribunal Disciplinario Superior'. Se conservan literalmente."
        )
    return notes


def main() -> None:
    result = {
        "scope": "Auditoría de texto de los PowerPoint 11, 12, 13 y 14; no se modificaron los originales.",
        "method": (
            "Texto nativo extraído de todas las formas OOXML y ordenado mediante inspección visual. "
            "Se verificó que todos los elementos <a:t> están representados. Se inspeccionaron las imágenes "
            "únicas y los renders completos."
        ),
        "global_findings": {
            "tables": 0,
            "smartart": 0,
            "charts": 0,
            "speaker_notes_with_text": 0,
            "repeated_background": "Fondo geométrico blanco, verde oscuro y dorado; sin texto.",
            "repeated_logo_ocr": LOGO_OCR,
        },
        "decks": [],
    }

    for deck, filename in FILES.items():
        prs = Presentation(BASE / filename)
        deck_record = {
            "deck": deck,
            "filename": filename,
            "slide_count": len(prs.slides),
            "slides": [],
        }
        assert len(ORDER[deck]) == len(prs.slides)

        for slide_no, slide in enumerate(prs.slides, 1):
            text_by_index = {
                index: normalize_text(shape.text)
                for index, shape in enumerate(slide.shapes)
                if getattr(shape, "has_text_frame", False) and normalize_text(shape.text)
            }
            expected = ORDER[deck][slide_no]
            assert set(expected) == set(text_by_index), (
                deck,
                slide_no,
                expected,
                sorted(text_by_index),
            )
            segments = [text_by_index[index] for index in expected]
            deck_record["slides"].append(
                {
                    "slide": slide_no,
                    "native_text_in_reading_order": segments,
                    "image_content": image_entries(deck, slide_no, len(prs.slides)),
                    "uncertainties": slide_uncertainties(deck, slide_no),
                }
            )
        result["decks"].append(deck_record)

    (OUT / "audit_11_14.json").write_text(
        json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )

    md: list[str] = [
        "# Auditoría de texto — diapositivas 11 a 14",
        "",
        "No se modificaron los PPTX originales. El orden de lectura se determinó mediante inspección visual de cada diapositiva.",
        "",
        "Hallazgos generales:",
        "",
        "- 87 diapositivas en total: 20 + 19 + 25 + 23.",
        "- No hay tablas, SmartArt, gráficos nativos ni notas del orador con contenido.",
        "- El fondo geométrico repetido no contiene texto.",
        "- Todas las diapositivas llevan un logotipo de imagen con texto legible: «200 / POLICIA BOLIVIANA / 1826 - 2026». El lema diminuto del listón inferior no puede resolverse con fidelidad.",
        "- Para audio, el logotipo se registra por diapositiva por literalidad, pero al ser decorativo y repetitivo conviene omitirlo de la locución salvo instrucción contraria.",
        "",
    ]
    for deck_record in result["decks"]:
        md.extend(
            [
                f"## {deck_record['filename']}",
                "",
                f"Total: {deck_record['slide_count']} diapositivas.",
                "",
            ]
        )
        for slide in deck_record["slides"]:
            md.extend([f"### Diapositiva {slide['slide']}", "", "Texto nativo, en orden de lectura:", ""])
            for segment in slide["native_text_in_reading_order"]:
                lines = segment.splitlines() or [segment]
                md.append(f"- {lines[0]}")
                md.extend(f"  {line}" for line in lines[1:])
            md.extend(["", "Texto/elementos en imágenes:", ""])
            for entry in slide["image_content"]:
                md.append(f"- {entry['description']}")
                if entry["text"]:
                    md.append("  Transcripción: " + " | ".join(entry["text"]))
                if entry["uncertainty"]:
                    md.append("  Incertidumbre: " + entry["uncertainty"])
            if slide["uncertainties"]:
                md.extend(["", "Observaciones de fidelidad:", ""])
                md.extend(f"- {note}" for note in slide["uncertainties"])
            md.append("")

    (OUT / "audit_11_14.md").write_text("\n".join(md), encoding="utf-8")


if __name__ == "__main__":
    main()
