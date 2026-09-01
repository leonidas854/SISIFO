#!/usr/bin/env python3
"""Build and verify the semantic reading-order audit for themes 1 through 6.

The source PPTX files are opened read-only.  Each output text block corresponds
to one native PowerPoint shape.  Paragraphs and formatted runs inside a shape
are recomposed without rewriting their literal ``<a:t>`` values.  The only
inserted text is PowerPoint automatic numbering, which is visible on the slide
but is stored as paragraph formatting rather than as ``<a:t>`` content.
"""

from __future__ import annotations

import hashlib
import json
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


ROOT = Path(__file__).resolve().parents[1]
SOURCE_DIR = ROOT / "diapos_original"
OUTPUT_PATH = Path(__file__).with_name("audit_01_06.json")


# One-based PowerPoint shape positions, reordered for natural narration.
# Every text-bearing shape on every slide is listed exactly once.  Titles and
# introductory statements come first; repeated content units are read as
# heading + explanation from left to right and then from top to bottom.
READING_ORDER: dict[int, dict[int, list[int]]] = {
    1: {
        1: [10, 4, 3, 11, 5],
        2: [2, 3, 5],
        3: [3, 4],
        4: [2, 8, 3, 9, 5, 10, 6, 11, 7],
        5: [2, 6, 3, 5, 7],
        6: [4, 7, 5],
        7: [4, 7, 5],
        8: [4, 7, 5],
        9: [2, 7, 3, 8, 5, 9, 6],
        10: [5, 2, 6, 3, 8, 7],
        11: [4, 7, 6, 5, 2],
        12: [2, 5, 3, 6, 9, 8, 7],
        13: [2, 4, 5, 8, 7, 6],
        14: [2, 4, 5, 8, 7, 6],
        15: [2, 4, 5, 8, 7, 6],
        16: [2, 4, 7, 6, 5, 8],
        17: [7, 6, 4],
        18: [4, 7, 5],
        19: [7, 6, 4],
        20: [8, 3],
    },
    2: {
        1: [10, 4, 3, 11, 5],
        2: [6, 3, 4],
        3: [6, 3, 4],
        4: [2, 7, 3, 8, 5, 9, 6],
        5: [2, 6, 7, 5, 10, 4, 11, 9, 12, 8],
        6: [4, 5, 2],
        7: [2, 11, 6, 3, 5, 7, 13, 12],
        8: [4, 5],
        9: [2, 7, 3, 8, 5, 9, 6],
        10: [2, 10, 7, 3, 8, 5, 9, 6],
        11: [2, 5, 6, 7, 3, 8, 9],
        12: [2, 3, 5],
        13: [2, 3, 5],
        14: [2, 3],
        15: [4, 7, 6, 5, 2],
        16: [2, 4, 7, 6, 5],
        17: [2, 7, 3, 8, 5, 9, 6],
        18: [8, 3],
    },
    3: {
        1: [10, 4, 3, 11, 5],
        2: [2, 5, 3, 6],
        3: [3, 4, 2, 8, 6, 9, 7, 10],
        4: [2, 8, 5, 10, 13, 11, 9, 3, 14, 12],
        5: [4, 5],
        6: [2, 8, 5, 10, 12, 11, 9, 3],
        7: [2, 6, 3, 5],
        8: [4, 5, 2, 7, 6],
        9: [2, 8, 5, 10, 12, 11, 9, 3],
        10: [2, 3, 5],
        11: [4, 5],
        12: [2, 4, 5, 8, 7, 6],
        13: [4, 7, 6, 5, 2],
        14: [4, 5],
        15: [4, 5],
        16: [2, 4, 5, 6, 7],
        17: [2, 4, 9, 5, 8, 7, 6],
        18: [8, 3],
    },
    4: {
        1: [10, 4, 3, 11, 5],
        2: [2, 3, 5],
        3: [3, 5, 4, 6, 7, 8, 9],
        4: [3, 5, 4, 7, 8, 6],
        5: [3, 5, 4, 6, 7, 8, 9],
        6: [4, 7, 5],
        7: [7, 3, 4, 5, 6],
        8: [7, 3, 4, 5, 6],
        9: [7, 3, 6, 4, 5],
        10: [2, 3, 5],
        11: [2, 3, 5],
        12: [5, 3, 4, 6],
        13: [7, 3, 6, 5, 4, 9],
        14: [2, 3, 5],
        15: [2, 3, 5],
        16: [2, 3, 5],
        17: [5, 3, 4, 6],
        18: [8, 3],
    },
    5: {
        1: [10, 4, 3, 11, 5],
        2: [2, 5, 3],
        # The semantic sequence is the visible automatic numbering 1, 2, 3, 4.
        3: [2, 4, 6, 5, 8, 7],
        4: [3, 4],
        5: [2, 5, 3, 7, 6],
        6: [2, 5, 3],
        7: [2, 3],
        8: [2, 3],
        9: [2, 3, 5, 6],
        10: [4, 7, 5],
        11: [4, 7, 5],
        12: [6, 2, 3, 5],
        13: [6, 2, 3, 5],
        14: [4, 5, 2],
        15: [4, 5, 2, 6],
        16: [4, 5, 2],
        17: [4, 5, 2],
        18: [8, 3],
    },
    6: {
        1: [10, 4, 3, 11, 5],
        2: [3, 4],
        3: [3, 4],
        4: [3, 4],
        5: [3, 4],
        6: [3, 4],
        7: [3, 4],
        8: [3, 4],
        9: [3, 4],
        10: [2, 5, 6, 3],
        11: [2, 5, 6, 3],
        12: [2, 5, 6, 3],
        13: [2, 5, 6, 3],
        14: [2, 5, 6, 3],
        15: [2, 5, 6, 3],
        16: [2, 5, 6, 3],
        17: [2, 5, 6, 3],
        18: [2, 5, 6, 3],
        19: [2, 5, 6, 3],
        20: [2, 4, 6, 5, 7, 8],
        21: [2, 4, 5, 7, 6],
        22: [2, 4, 5, 7, 6],
        23: [8, 3],
    },
}


def source_path(theme_number: int) -> Path:
    return SOURCE_DIR / f"{theme_number}_ DIAPOSITIVA.pptx"


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def text_shapes(slide) -> dict[int, object]:
    """Return top-level text-bearing shapes keyed by one-based shape position."""
    result: dict[int, object] = {}
    for index, shape in enumerate(slide.shapes, start=1):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            raise AssertionError(
                "Unexpected grouped text shape; reading-order map needs an explicit nested path"
            )
        if shape.element.xpath(".//a:t"):
            result[index] = shape
    return result


def paragraph_plain_text(paragraph_element) -> str:
    """Recompose a paragraph from its literal text runs in document order."""
    return "".join((node.text or "") for node in paragraph_element.xpath(".//a:t"))


def automatic_number_prefix(paragraph_element) -> str | None:
    properties = paragraph_element.find(qn("a:pPr"))
    if properties is None:
        return None
    auto_number = properties.find(qn("a:buAutoNum"))
    if auto_number is None:
        return None
    start_at = auto_number.get("startAt")
    return f"{int(start_at) if start_at else 1}. "


def shape_text_block(shape) -> tuple[str, list[str], list[str]]:
    """Return text block, represented XML node paths, and inserted number labels."""
    paragraphs: list[str] = []
    node_paths: list[str] = []
    inserted_number_labels: list[str] = []
    tree = shape.element.getroottree()

    for paragraph in shape.element.xpath(".//a:p"):
        nodes = paragraph.xpath(".//a:t")
        if not nodes:
            continue
        text = paragraph_plain_text(paragraph)
        number_prefix = automatic_number_prefix(paragraph)
        if number_prefix:
            text = number_prefix + text
            inserted_number_labels.append(number_prefix.rstrip())
        paragraphs.append(text)
        node_paths.extend(tree.getpath(node) for node in nodes)

    if not paragraphs:
        raise AssertionError(f"Text-bearing shape produced no paragraphs: {shape.name}")
    return "\n".join(paragraphs), node_paths, inserted_number_labels


def assert_exact_node_coverage(slide, represented_paths: Iterable[str]) -> int:
    tree = slide.element.getroottree()
    source_paths = [tree.getpath(node) for node in slide.element.xpath(".//a:t")]
    represented = list(represented_paths)
    if len(represented) != len(set(represented)):
        raise AssertionError("At least one <a:t> node was represented more than once")
    if set(source_paths) != set(represented):
        missing = sorted(set(source_paths) - set(represented))
        unexpected = sorted(set(represented) - set(source_paths))
        raise AssertionError(
            f"<a:t> coverage mismatch; missing={missing}, unexpected={unexpected}"
        )
    return len(source_paths)


def build() -> dict:
    themes: list[dict] = []
    total_slides = 0
    total_a_t_nodes = 0
    total_blocks = 0
    total_auto_numbers = 0

    for theme_number in range(1, 7):
        path = source_path(theme_number)
        presentation = Presentation(path)
        expected_slide_orders = READING_ORDER[theme_number]
        if set(expected_slide_orders) != set(range(1, len(presentation.slides) + 1)):
            raise AssertionError(f"Theme {theme_number}: slide order map is incomplete")

        slide_records: list[dict] = []
        theme_a_t_nodes = 0
        theme_blocks = 0
        theme_auto_numbers = 0

        for slide_number, slide in enumerate(presentation.slides, start=1):
            shapes = text_shapes(slide)
            order = expected_slide_orders[slide_number]
            if len(order) != len(set(order)):
                raise AssertionError(
                    f"Theme {theme_number}, slide {slide_number}: duplicate shape in order"
                )
            if set(order) != set(shapes):
                missing = sorted(set(shapes) - set(order))
                unexpected = sorted(set(order) - set(shapes))
                raise AssertionError(
                    f"Theme {theme_number}, slide {slide_number}: "
                    f"shape coverage mismatch; missing={missing}, unexpected={unexpected}"
                )

            blocks: list[str] = []
            represented_paths: list[str] = []
            automatic_number_labels: list[str] = []
            for shape_index in order:
                block, node_paths, number_labels = shape_text_block(shapes[shape_index])
                blocks.append(block)
                represented_paths.extend(node_paths)
                automatic_number_labels.extend(number_labels)

            node_count = assert_exact_node_coverage(slide, represented_paths)
            slide_records.append(
                {
                    "slide_number": slide_number,
                    "native_text_in_reading_order": blocks,
                    "verification": {
                        "native_text_block_count": len(blocks),
                        "a_t_nodes_in_source_slide": node_count,
                        "a_t_nodes_represented": len(represented_paths),
                        "all_a_t_nodes_represented_exactly_once": True,
                        "visible_automatic_number_labels_added": automatic_number_labels,
                    },
                }
            )
            theme_a_t_nodes += node_count
            theme_blocks += len(blocks)
            theme_auto_numbers += len(automatic_number_labels)

        themes.append(
            {
                "theme_number": theme_number,
                "source_file": path.name,
                "source_sha256": sha256(path),
                "slide_count": len(presentation.slides),
                "slides": slide_records,
                "verification": {
                    "a_t_nodes_in_source": theme_a_t_nodes,
                    "a_t_nodes_represented": theme_a_t_nodes,
                    "all_a_t_nodes_represented_exactly_once": True,
                    "native_text_block_count": theme_blocks,
                    "visible_automatic_number_labels_added": theme_auto_numbers,
                },
            }
        )
        total_slides += len(presentation.slides)
        total_a_t_nodes += theme_a_t_nodes
        total_blocks += theme_blocks
        total_auto_numbers += theme_auto_numbers

    return {
        "schema_version": 1,
        "scope": "Original PPTX themes 1 through 6",
        "source_directory": "diapos_original",
        "reading_order_policy": (
            "Titles and introductory statements first; content units follow in semantic "
            "heading-plus-explanation order, from left to right and then top to bottom."
        ),
        "text_policy": (
            "Each block preserves the literal native <a:t> text of one PowerPoint shape. "
            "Rasterized logo text is excluded. Visible automatic numbering is added from "
            "PowerPoint paragraph formatting."
        ),
        "themes": themes,
        "verification": {
            "theme_count": len(themes),
            "slide_count": total_slides,
            "native_text_block_count": total_blocks,
            "a_t_nodes_in_sources": total_a_t_nodes,
            "a_t_nodes_represented": total_a_t_nodes,
            "all_a_t_nodes_represented_exactly_once": True,
            "visible_automatic_number_labels_added": total_auto_numbers,
            "source_pptx_modified": False,
        },
    }


def main() -> None:
    audit = build()
    OUTPUT_PATH.write_text(
        json.dumps(audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    verification = audit["verification"]
    print(
        f"Wrote {OUTPUT_PATH}: {verification['theme_count']} themes, "
        f"{verification['slide_count']} slides, "
        f"{verification['native_text_block_count']} blocks, "
        f"{verification['a_t_nodes_represented']} <a:t> nodes."
    )


if __name__ == "__main__":
    main()
