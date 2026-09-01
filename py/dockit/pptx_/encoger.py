# -*- coding: utf-8 -*-
"""Ajusta la altura de cada caja de texto a su contenido real.
Las cajas del deck son mucho más altas que el texto que llevan: ese aire sobrante es el
que impedía colocar imágenes grandes. Mide con Carlito (métricamente igual a Calibri)."""
from pptx import Presentation
from pptx.util import Emu, Pt
from PIL import ImageFont

EMU_IN  = 914400.0
CARLITO = '/usr/share/fonts/carlito/Carlito-Regular.ttf'
CARLBD  = '/usr/share/fonts/carlito/Carlito-Bold.ttf'
_f = {}

def fuente(pt, bold):
    k = (round(pt*2), bold)
    if k not in _f:
        try:    _f[k] = ImageFont.truetype(CARLBD if bold else CARLITO, int(pt*4))
        except Exception: _f[k] = ImageFont.load_default()
    return _f[k]

def ancho_pt(txt, pt, bold):
    return fuente(pt, bold).getlength(txt) / 4.0

def n_lineas(txt, pt, bold, ancho_disp_pt):
    if not txt.strip(): return 1
    lineas, cur = 0, ''
    for pal in txt.split():
        t = (cur+' '+pal).strip()
        if ancho_pt(t, pt, bold) <= ancho_disp_pt or not cur: cur = t
        else: lineas += 1; cur = pal
    return lineas + (1 if cur else 0)


A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

def _sz_de(el, lvl):
    """Busca sz en un lstStyle/txStyles para el nivel de párrafo dado."""
    if el is None: return None
    for tag in (f'{A}lvl{lvl+1}pPr', f'{A}lvl1pPr'):
        p = el.find(tag)
        if p is not None:
            d = p.find(f'{A}defRPr')
            if d is not None and d.get('sz'): return int(d.get('sz'))/100.0
    return None

def tam_efectivo(sh, par, base=18.0):
    """Resuelve el tamaño de letra real: run -> párrafo -> forma -> diseño -> patrón."""
    for r in par.runs:
        if r.font.size: return r.font.size.pt
    lvl = par.level or 0
    pPr = par._p.find(f'{A}pPr')
    if pPr is not None:
        d = pPr.find(f'{A}defRPr')
        if d is not None and d.get('sz'): return int(d.get('sz'))/100.0
    v = _sz_de(sh.text_frame._txBody.find(f'{A}lstStyle'), lvl)
    if v: return v
    try:
        ph = sh.element.find(f'.//{{http://schemas.openxmlformats.org/presentationml/2006/main}}ph')
        idx = ph.get('idx') if ph is not None else None
        tipo = ph.get('type') if ph is not None else None
        for origen in (sh.part.slide.slide_layout, sh.part.slide.slide_layout.slide_master):
            for o in origen.placeholders:
                oph = o.element.find(f'.//{{http://schemas.openxmlformats.org/presentationml/2006/main}}ph')
                if oph is None: continue
                if (oph.get('idx') == idx) or (tipo and oph.get('type') == tipo):
                    v = _sz_de(o.text_frame._txBody.find(f'{A}lstStyle'), lvl)
                    if v: return v
        ts = sh.part.slide.slide_layout.slide_master.element.find(
                f'{{http://schemas.openxmlformats.org/presentationml/2006/main}}txStyles')
        if ts is not None:
            v = _sz_de(ts.find(f'{{http://schemas.openxmlformats.org/presentationml/2006/main}}bodyStyle'), lvl)
            if v: return v
    except Exception:
        pass
    return base

def alto_texto_in(sh, base=18.0):
    tf = sh.text_frame
    ml = (tf.margin_left  or Emu(91440))/EMU_IN
    mr = (tf.margin_right or Emu(91440))/EMU_IN
    mt = (tf.margin_top   or Emu(45720))/EMU_IN
    mb = (tf.margin_bottom or Emu(45720))/EMU_IN
    ancho_pt_disp = max(24.0, ((sh.width or 0)/EMU_IN - ml - mr) * 72)
    total = 0.0
    for par in tf.paragraphs:
        pt   = tam_efectivo(sh, par, base)
        bold = any(r.font.bold for r in par.runs)
        sangria = 22 if par.level or _tiene_vinieta(par) else 0
        nl = n_lineas(par.text, pt, bold, ancho_pt_disp - sangria)
        total += nl * pt * 1.22
        total += (par.space_before.pt if par.space_before else 0)
        total += (par.space_after.pt  if par.space_after  else 0)
    return total/72 + mt + mb

def _tiene_vinieta(par):
    pPr = par._p.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    if pPr is None: return False
    return any(c.tag.endswith(('buChar','buAutoNum')) for c in pPr)

def encoger(entrada, salida, holgura=0.10, min_alto=0.30, solo_reducir=False):
    prs = Presentation(entrada); rep = []
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame or sh.width is None: continue
            if not sh.text_frame.text.strip(): continue
            if sh.name.startswith(('Título','Titulo')): continue
            viejo = (sh.height or 0)/EMU_IN
            medido = alto_texto_in(sh)
            nuevo  = max(min_alto, medido*1.10 + holgura)
            if solo_reducir and nuevo >= viejo: continue
            if nuevo > viejo*0.88: continue          # solo si sobra bastante aire
            if abs(nuevo - viejo) < 0.15: continue
            sh.height = Emu(int(nuevo*EMU_IN))
            rep.append((i, sh.name, round(viejo,2), round(nuevo,2)))
    prs.save(salida)
    return rep

if __name__ == '__main__':
    import sys
    r = encoger(sys.argv[1], sys.argv[2], solo_reducir=(len(sys.argv)>3))
    print(f'{len(r)} cajas ajustadas')
    for i,n,a,b in r[:14]: print(f'  D{i:02d} {n[:26]:28s} {a:5.2f} -> {b:5.2f} in')
