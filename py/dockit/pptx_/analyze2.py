# -*- coding: utf-8 -*-
"""Calcula los huecos libres de cada diapositiva a partir del RENDER real.
Así se detecta el texto que se desborda de su caja, la banda del fondo y el logo."""
import subprocess, glob, os, json, sys, tempfile
import numpy as np
from PIL import Image
from pptx import Presentation

CELL = 0.05          # pulgadas por celda
PAD  = 3             # celdas de margen alrededor de la tinta (0.15 in)
INK  = 26            # desviación de blanco para considerar "tinta"

def render(pptx, outdir, dpi=100):
    os.makedirs(outdir, exist_ok=True)
    subprocess.run(['libreoffice','--headless','--convert-to','pdf','--outdir',outdir,pptx],
                   check=True, capture_output=True, timeout=600)
    pdf = glob.glob(os.path.join(outdir,'*.pdf'))[0]
    subprocess.run(['pdftoppm','-png','-r',str(dpi),pdf,os.path.join(outdir,'p')],
                   check=True, capture_output=True)
    return sorted(glob.glob(os.path.join(outdir,'p-*.png')))

def dilate(g, k):
    out = g.copy()
    for _ in range(k):
        d = out.copy()
        d[1:,:]  |= out[:-1,:]; d[:-1,:] |= out[1:,:]
        d[:,1:]  |= out[:,:-1]; d[:,:-1] |= out[:,1:]
        out = d
    return out

def busy_from_png(png, nx, ny):
    a = np.asarray(Image.open(png).convert('RGB').resize((nx, ny), Image.BOX)).astype(int)
    return (255 - a).max(axis=2) > INK

def usable(w, h, lo=1.35, hi=2.2):
    """Área realmente aprovechable: el mayor rectángulo de proporción fotográfica
    que cabe dentro. Evita que una franja larga y flaca gane por área bruta."""
    if h <= 0: return 0.0
    r = w/h
    if r > hi:  return (h*hi)*h
    if r < lo:  return w*(w/lo)
    return w*h

def rects(g, topn=4, minw=1.9, minh=1.05):
    rows, cols = g.shape
    found, heights = [], [0]*cols
    for r in range(rows):
        for c in range(cols):
            heights[c] = 0 if g[r,c] else heights[c]+1
        st = []
        for c in range(cols+1):
            hcur = heights[c] if c < cols else 0
            start = c
            while st and st[-1][1] >= hcur:
                s0,h0 = st.pop()
                found.append((s0, r-h0+1, c-s0, h0)); start = s0
            st.append((start,hcur))
    cands = []
    for x,y,w,h in found:
        wi, hi_ = w*CELL, h*CELL
        if wi < minw or hi_ < minh: continue
        cands.append((usable(wi,hi_), x, y, w, h))
    cands.sort(reverse=True)
    out = []
    for u,x,y,w,h in cands:
        if any(not (x+w<=ox or ox+ow<=x or y+h<=oy or oy+oh<=y) for _,ox,oy,ow,oh in out): continue
        out.append((u,x,y,w,h))
        if len(out) >= topn: break
    return [dict(util=round(u,2), area=round(w*h*CELL*CELL,2), left=round(x*CELL,2),
                 top=round(y*CELL,2), w=round(w*CELL,2), h=round(h*CELL,2))
            for u,x,y,w,h in out]

def titulo(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.name.startswith(('Título','Titulo')) and sh.text_frame.text.strip():
            return sh.text_frame.text.strip()
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip()
    return ''

def analizar(pptx, outjson, workdir=None):
    workdir = workdir or tempfile.mkdtemp()
    pngs = render(pptx, workdir)
    prs  = Presentation(pptx)
    W, H = prs.slide_width/914400, prs.slide_height/914400
    nx, ny = int(round(W/CELL)), int(round(H/CELL))
    data = []
    for i, (png, s) in enumerate(zip(pngs, prs.slides), 1):
        g = dilate(busy_from_png(png, nx, ny), PAD)
        data.append(dict(n=i, titulo=titulo(s), ocupacion=round(float(g.mean())*100,1),
                         huecos=rects(g)))
    json.dump(data, open(outjson,'w'), ensure_ascii=False, indent=1)
    return data

if __name__ == '__main__':
    d = analizar(sys.argv[1], sys.argv[2], sys.argv[3] if len(sys.argv)>3 else None)
    for x in d:
        hs = ' ; '.join(f"{h['w']}x{h['h']}@({h['left']},{h['top']})" for h in x['huecos'][:2]) or 'SIN HUECO'
        print(f"D{x['n']:2d} ocup={x['ocupacion']:4.1f}%  {x['titulo'][:38]:40s} -> {hs}")
