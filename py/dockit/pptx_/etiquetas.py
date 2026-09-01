from pathlib import Path
import os
# -*- coding: utf-8 -*-
"""Extrae los 'conceptos' de una diapositiva: los rótulos que encabezan cada bloque.
En estos decks los encabezados van en negrita y en color, así que ese es el criterio
principal; el largo y la forma solo sirven para descartar."""
import re, unicodedata
from pptx import Presentation

GENERICAS = {
 'ejemplos','ejemplo','definicion','concepto','caracteristicas','interpretacion policial',
 'aplicacion','aplicacion policial','enfoque de accion','se subdividen en tres categorias',
 'gracias','gestion 2026','nota','importante','clasificacion','tipos','introduccion',
 'objetivo','objetivos','finalidad','contenido','desarrollo','en que consiste',
}

def norm(s):
    s = unicodedata.normalize('NFD', s.lower())
    return ''.join(c for c in s if unicodedata.category(c) != 'Mn').strip()

def limpio(s):
    s = re.sub(r'^\s*[\dA-Za-z][\.\)\-•]\s+', '', s)
    s = re.sub(r'\s+', ' ', s).strip(' :.-–—•;')
    return s

def _es_encabezado(par):
    runs = [r for r in par.runs if r.text.strip()]
    if not runs: return False
    neg = sum(1 for r in runs if r.font.bold)
    col = 0
    for r in runs:
        try:
            if r.font.color is not None and r.font.color.type is not None: col += 1
        except Exception: pass
    return neg >= max(1, len(runs)//2) or col >= max(1, len(runs)//2)

def conceptos(slide, max_n=4):
    fuertes, debiles = [], []
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        if sh.name.startswith(('Título','Titulo')): continue
        for par in sh.text_frame.paragraphs:
            t = limpio(par.text)
            if not t or len(t) > 46 or len(t) < 3: continue
            if len(t.split()) > 6: continue
            if t.endswith((';', '.')) and len(t.split()) > 3: continue
            n = norm(t)
            if n in GENERICAS: continue
            if any(n == norm(o) for o in fuertes + debiles): continue
            (fuertes if _es_encabezado(par) else debiles).append(t)
    return (fuertes + debiles)[:max_n]

def resumen_deck(path):
    prs = Presentation(path)
    return {i: conceptos(s) for i, s in enumerate(prs.slides, 1)}

if __name__ == '__main__':
    import os
def _base_proyecto() -> str:
    """Carpeta del proyecto con los .pptx fuente.

    Se fija con TALLER_PROYECTO; si no, se busca subiendo desde donde se
    ejecuta. Nunca una ruta absoluta escrita a mano: el motor no puede depender
    de que exista una carpeta de trabajo concreta, porque se borran.
    """
    if _v := os.environ.get("TALLER_PROYECTO"):
        return str(Path(_v).resolve())
    _aqui = Path.cwd().resolve()
    for _c in [_aqui, *_aqui.parents]:
        if any(_c.glob("*DIAPOSITIVA*.pptx")) or (_c / "diapos_original").is_dir():
            return str(_c)
    return str(_aqui)


    BASE=_base_proyecto()
    for f in ['13_ DIAPOSITIVA.pptx','6_ DIAPOSITIVA.pptx']:
        print('==', f)
        r = resumen_deck(os.path.join(BASE,f))
        for n in sorted(r):
            if r[n]: print(f'  {n:2d}: {" | ".join(r[n])}')
