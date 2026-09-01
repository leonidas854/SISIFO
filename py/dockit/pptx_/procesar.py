# -*- coding: utf-8 -*-
"""Procesa un tema completo: libera espacio, elige la imagen de cada diapositiva
(foto web, gráfico propio o generada por IA) y la coloca en el hueco calculado."""
import json, os, glob, shutil, subprocess
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

import analyze2 as A, encoger as EN, reducir as RE, compactar as CO, reacomodar as RA
import photo2 as P2, compose as C, grafgen as GG, graficos as G

BASE = '/home/leonidas/SSD500/Develoment/tareas/proyectos_policias'
MIN_FOTO = (3.0, 1.45)      # ancho, alto mínimos para que una foto aporte
MIN_GRAF = (4.2, 1.00)      # por debajo de esto el gráfico no se lee de lejos

def sirve(h, minimo):
    return h['w'] >= minimo[0] and h['h'] >= minimo[1]

BANDA_TITULO = 1.25       # la franja superior es del título: no se invade

def mejor_hueco(d, tipo):
    minimo = MIN_GRAF if tipo == 'grafico' else MIN_FOTO
    ok = [h for h in d['huecos'] if sirve(h, minimo)]
    bajo = [h for h in ok if h['top'] >= BANDA_TITULO or h['h'] >= 3.0]
    return (bajo or ok or [None])[0]

def _sin_sitio(datos, filas):
    """Diapositivas que aún no tienen un hueco utilizable para lo que les toca."""
    tipos = {f['n']: f['tipo'] for f in filas}
    fuera = []
    for d in datos:
        if d['n'] not in tipos: continue
        if mejor_hueco(d, tipos[d['n']]) is None:
            fuera.append(d['n'])
    return fuera

def preparar(src, tmp, tag, filas, rondas=3):
    """Libera sitio de forma progresiva hasta que toda diapositiva tenga su hueco."""
    p0 = f'{tmp}/{tag}_e.pptx'
    EN.encoger(src, p0, solo_reducir=True)                  # cajas al alto de su texto
    p = f'{tmp}/{tag}_0.pptx'
    RA.reacomodar(p0, p, [f['n'] for f in filas], frac=0.60)  # abre columna lateral
    d = A.analizar(p, f'{tmp}/{tag}_0.json', f'{tmp}/r0_{tag}')
    for i in range(1, rondas+1):
        faltan = _sin_sitio(d, filas)
        if not faltan:
            break
        fr = [0.86, 0.78, 0.70][i-1]
        kc = [0.84, 0.76, 0.68][i-1]
        pa = f'{tmp}/{tag}_{i}a.pptx'; pb = f'{tmp}/{tag}_{i}.pptx'
        RE.reducir(p, pa, {n: fr for n in faltan}, min_chars=60)
        CO.compactar_pptx(pa, pb, {n: kc for n in faltan})
        p = pb
        d = A.analizar(p, f'{tmp}/{tag}_{i}.json', f'{tmp}/r{i}_{tag}')
    return p, d

def _foto(tag, webdir='web2'):
    ops = sorted(glob.glob(f'{webdir}/{tag}_*.jpg'))
    buenas = [o for o in ops if not P2.es_afiche(o)]
    return (buenas or ops or [None])[0]

def colocar(src_pptx, datos, filas, salida, imgdir, webdir='web2', iadir='ia'):
    os.makedirs(imgdir, exist_ok=True)
    prs = Presentation(src_pptx)
    rep = []
    for f in filas:
        n = f['n']
        tipo = f['tipo']
        free = mejor_hueco(datos[n-1], tipo)
        if free is None and tipo != 'grafico':
            free = mejor_hueco(datos[n-1], 'grafico')
        if free is None:
            rep.append((n, 'sin sitio', f['titulo'][:40])); continue
        src  = None
        if tipo == 'foto':
            src = _foto(f['q'].replace(' ','_')[:44], webdir)
            if not src: tipo = 'grafico'
        if tipo == 'ia':
            cand = f'{iadir}/t{f["tema"]}_d{n:02d}.jpg'
            if os.path.exists(cand): src = cand
            else: tipo = 'foto' if _foto(f['q'].replace(' ','_')[:44], webdir) else 'grafico'
            if tipo == 'foto': src = _foto(f['q'].replace(' ','_')[:44], webdir)
        if tipo in ('foto','ia') and src:
            claro, _ = P2.borde_claro(Image.open(src))
            x,y,w,h = P2.encuadre(free, src, align='left' if free['w'] > 5 else 'center',
                                  valign='top', expandir=1.10 if claro else 1.0)
            dest, _m = P2.preparar(src, f'{imgdir}/diapo{n:02d}.jpg', w, h)
        else:
            m = 0.10
            w = round(free['w'] - 2*m, 2)
            h = round(min(free['h'] - 2*m, max(1.0, w*0.30)), 2)
            svg = GG.grafico(f.get('conceptos') or [f['titulo']], w, h,
                             estilo=f.get('estilo','tarjetas'))
            if svg is None:
                rep.append((n, 'sin concepto', f['titulo'][:40])); continue
            x = round(free['left'] + (free['w']-w)/2, 2)
            y = round(free['top']  + (free['h']-h)/2, 2)
            dest = C.save(svg, f'{imgdir}/diapo{n:02d}.png', w)
            tipo = 'grafico'
        prs.slides[n-1].shapes.add_picture(dest, Inches(x), Inches(y), Inches(w), Inches(h))
        rep.append((n, tipo, f'{w}x{h}in  {f["titulo"][:38]}'))
    prs.save(salida)
    return rep

def procesar(tema, plan, tmp='tmp', outdir='salida'):
    os.makedirs(tmp, exist_ok=True); os.makedirs(outdir, exist_ok=True)
    temas = json.load(open('todos_temas.json'))
    src = os.path.join(BASE, temas[tema]['archivo'])
    filas = [dict(fi, tema=tema) for fi in plan[tema]]
    p, datos = preparar(src, tmp, tema, filas)
    salida = f'{outdir}/TEMA{tema}_con_imagenes.pptx'
    rep = colocar(p, datos, filas, salida, f'{outdir}/imagenes_tema{tema}')
    return salida, rep

if __name__ == '__main__':
    import sys
    plan = json.load(open('plan_full.json'))
    for t in sys.argv[1:]:
        s, rep = procesar(t, plan)
        ok = sum(1 for _,x,_ in rep if x in ('foto','grafico','ia'))
        print(f'TEMA {t}: {ok}/{len(rep)} colocadas -> {s}', flush=True)
        for n,x,d in rep:
            if x not in ('foto','grafico','ia'): print(f'   D{n:02d} {x}: {d}')
