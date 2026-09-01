# -*- coding: utf-8 -*-
"""Compacta verticalmente el contenido de una diapositiva para liberar una franja abajo.
Escala posiciones y alturas respecto de un ancla (el pie del título), conservando la
distribución relativa: las cajas y sus flechas siguen alineadas entre sí."""
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0

def es_fondo(sh, W, H):
    return (sh.width or 0) >= W*0.95*EMU_IN and (sh.height or 0) >= H*0.95*EMU_IN

def es_logo(sh, W, H):
    if sh.left is None: return False
    return (sh.left/EMU_IN > W*0.72 and sh.top/EMU_IN > H*0.75)

def compactar(slide, W, H, k=0.82, ancla=None):
    titulos = [s for s in slide.shapes if s.has_text_frame and s.name.startswith(('Título','Titulo'))]
    if ancla is None:
        ancla = max(((t.top or 0)+(t.height or 0))/EMU_IN for t in titulos) if titulos else 1.5
    movidas = 0
    for sh in slide.shapes:
        if sh.left is None or sh.top is None: continue
        if es_fondo(sh, W, H) or es_logo(sh, W, H) or sh in titulos: continue
        t = sh.top/EMU_IN
        h = (sh.height or 0)/EMU_IN
        if t + h <= ancla: continue                    # está por encima del ancla
        nt = ancla + (t - ancla)*k
        nh = h*k
        sh.top    = Emu(int(nt*EMU_IN))
        sh.height = Emu(int(max(nh, 0.18)*EMU_IN))
        movidas += 1
    return movidas

def compactar_pptx(entrada, salida, plan):
    """plan: {n_diapo: k}"""
    prs = Presentation(entrada)
    W, H = prs.slide_width/EMU_IN, prs.slide_height/EMU_IN
    rep = []
    for n, k in plan.items():
        rep.append((n, compactar(prs.slides[n-1], W, H, k)))
    prs.save(salida)
    return rep

if __name__ == '__main__':
    import sys
    r = compactar_pptx(sys.argv[1], sys.argv[2], {4:0.80, 15:0.78})
    for n,m in r: print(f'D{n:02d}: {m} formas compactadas')
