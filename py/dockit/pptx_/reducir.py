# -*- coding: utf-8 -*-
"""Reduce el tamaño de letra de los bloques de texto largos para liberar sitio a la imagen.
No cambia ni una palabra: solo el cuerpo. Los títulos y las etiquetas cortas quedan intactos."""
from pptx import Presentation
from pptx.util import Pt

def reducir_forma(sh, factor, minimo=13.0, base=18.0):
    tocado = False
    for par in sh.text_frame.paragraphs:
        for r in par.runs:
            actual = r.font.size.pt if r.font.size else base
            nuevo  = max(minimo, round(actual*factor, 1))
            if abs(nuevo - actual) > 0.05:
                r.font.size = Pt(nuevo); tocado = True
    return tocado

def reducir(entrada, salida, plan, min_chars=90):
    """plan: {n_diapo: factor}. Se aplica a cuerpos de texto, nunca a los títulos."""
    prs = Presentation(entrada); rep = []
    for n, f in plan.items():
        for sh in prs.slides[n-1].shapes:
            if not sh.has_text_frame: continue
            t = sh.text_frame.text.strip()
            if len(t) < min_chars: continue
            if sh.name.startswith(('Título','Titulo')): continue
            if reducir_forma(sh, f):
                rep.append((n, len(t), f))
    prs.save(salida)
    return rep

if __name__ == '__main__':
    import sys
    plan = {4:0.78, 5:0.80, 6:0.84, 12:0.84, 15:0.78}
    r = reducir(sys.argv[1], sys.argv[2], plan)
    for n,l,f in r: print(f'D{n:02d} cuerpo de {l} caracteres -> x{f}')
