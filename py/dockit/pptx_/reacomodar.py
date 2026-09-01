# -*- coding: utf-8 -*-
"""Angosta la caja de texto que ocupa casi todo el ancho para liberar una columna
lateral donde la imagen entra grande. Es la solución del ejemplo del usuario:
texto a la izquierda, foto a la derecha."""
from pptx import Presentation
from pptx.util import Emu
import encoger as E

EMU_IN = 914400.0

def cuerpos(slide, min_chars=80):
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame or sh.width is None: continue
        if sh.name.startswith(('Título','Titulo')): continue
        if len(sh.text_frame.text.strip()) < min_chars: continue
        out.append(sh)
    return out

def _despejada(slide, sh):
    """La diapositiva debe tener un solo bloque de texto y nada más que estorbe:
    si hay flechas, etiquetas o cajas alrededor, angostar provoca choques."""
    otros = 0
    for o in slide.shapes:
        if o is sh: continue
        if o.left is None: continue
        if o.has_text_frame and o.name.startswith(('Título','Titulo')): continue
        w = (o.width or 0)/EMU_IN; h = (o.height or 0)/EMU_IN
        if w >= 12.6 and h >= 7.0: continue           # fondo
        if o.left/EMU_IN > 9.5 and o.top/EMU_IN > 5.5: continue   # logo
        otros += 1
    return otros == 0

def estrechar(slide, frac=0.60, lado='izquierda', slide_w=13.33, min_ancho=4.2):
    """Reduce el ancho del cuerpo dominante y recalcula su alto para el nuevo ancho."""
    cs = cuerpos(slide)
    if len(cs) != 1: return None                     # solo si hay un cuerpo dominante
    sh = cs[0]
    if not _despejada(slide, sh): return None
    w  = sh.width/EMU_IN
    if w < slide_w*0.62: return None                 # ya es angosta
    nw = max(min_ancho, round(slide_w*frac, 2))
    if lado == 'derecha':
        sh.left = Emu(int((slide_w - nw - (sh.left/EMU_IN))*EMU_IN))
    sh.width = Emu(int(nw*EMU_IN))
    nh = E.alto_texto_in(sh)*1.08 + 0.12             # al angostar, el texto crece a lo alto
    sh.height = Emu(int(max(nh, 0.4)*EMU_IN))
    return round(w,2), nw, round(nh,2)

def reacomodar(entrada, salida, diapos, frac=0.58):
    prs = Presentation(entrada)
    W = prs.slide_width/EMU_IN
    rep = []
    for n in diapos:
        r = estrechar(prs.slides[n-1], frac=frac, slide_w=W)
        rep.append((n, r))
    prs.save(salida)
    return rep

if __name__ == '__main__':
    import sys
    for n, r in reacomodar(sys.argv[1], sys.argv[2], [2,3]):
        print(f'D{n:02d}:', 'sin cambio' if r is None else f'ancho {r[0]} -> {r[1]} in, alto {r[2]} in')
