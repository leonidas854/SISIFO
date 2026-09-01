"""Adaptador de lectura y auditoría de presentaciones PowerPoint."""

from __future__ import annotations

import hashlib
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from .dominio import Hallazgo, PlanVisual, RE_MARCADOR, Visual, normalizar
from .semantica import PuertoSemantico, SemanticaLexica


@dataclass(frozen=True)
class ImagenPPTX:
    nombre: str
    alt: str


@dataclass(frozen=True)
class DiapositivaPPTX:
    numero: int
    titulo: str
    texto: str
    imagenes: tuple[ImagenPPTX, ...]
    visuales: int
    notas: str


def _texto_forma(forma) -> str:
    if getattr(forma, "has_text_frame", False):
        return (forma.text_frame.text or "").strip()
    return ""


def _tamano_maximo(forma) -> float:
    maximo = 0.0
    if not getattr(forma, "has_text_frame", False):
        return maximo
    for parrafo in forma.text_frame.paragraphs:
        for run in parrafo.runs:
            if run.font.size:
                maximo = max(maximo, float(run.font.size.pt))
    return maximo


def _titulo(slide) -> str:
    from pptx.enum.shapes import PP_PLACEHOLDER

    candidatos = []
    for forma in slide.shapes:
        texto = _texto_forma(forma)
        if not texto:
            continue
        if getattr(forma, "is_placeholder", False):
            try:
                if forma.placeholder_format.type in {
                    PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
                }:
                    return texto
            except (ValueError, KeyError):
                pass
        candidatos.append((int(getattr(forma, "top", 0)), -_tamano_maximo(forma), texto))
    # En láminas en blanco generadas por código, el título suele ser el texto
    # más alto. El tamaño desempata cajas que comienzan en la misma coordenada.
    return min(candidatos)[2] if candidatos else ""


def _alt(forma) -> str:
    try:
        nodos = forma._element.xpath(".//p:cNvPr")
        if nodos:
            nodo = nodos[0]
            return (nodo.get("descr") or nodo.get("title") or "").strip()
    except Exception:
        pass
    return ""


def _notas(slide) -> str:
    try:
        frame = slide.notes_slide.notes_text_frame
        return (frame.text or "").strip() if frame else ""
    except (AttributeError, ValueError, KeyError):
        return ""


def leer_presentacion(ruta: str | Path) -> list[DiapositivaPPTX]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pres = Presentation(str(ruta))
    salida = []
    for numero, slide in enumerate(pres.slides, 1):
        textos = [_texto_forma(f) for f in slide.shapes]
        imagenes = []
        visuales = 0
        for forma in slide.shapes:
            tipo = forma.shape_type
            if tipo == MSO_SHAPE_TYPE.PICTURE:
                imagenes.append(ImagenPPTX(getattr(forma, "name", "imagen"), _alt(forma)))
                visuales += 1
            elif tipo in {
                MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TABLE,
                MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.DIAGRAM,
            }:
                visuales += 1
            elif tipo == MSO_SHAPE_TYPE.AUTO_SHAPE and not _texto_forma(forma):
                # Una forma puramente decorativa o un elemento de diagrama.
                visuales += 1
        salida.append(DiapositivaPPTX(
            numero=numero,
            titulo=_titulo(slide),
            texto="\n".join(t for t in textos if t),
            imagenes=tuple(imagenes),
            visuales=visuales,
            notas=_notas(slide),
        ))
    return salida


def _hallazgo(
    codigo: str, severidad: str, mensaje: str, accion: str,
    *, diapositiva: int | None = None, archivo: str = "",
) -> Hallazgo:
    return Hallazgo(
        codigo=codigo, severidad=severidad, mensaje=mensaje, accion=accion,
        diapositiva=diapositiva, archivo=archivo,
    )


def _parecido_titulo(a: str, b: str) -> float:
    sa, sb = set(normalizar(a).split()), set(normalizar(b).split())
    if not sa or not sb:
        return 0.0
    return len(sa & sb) / len(sa | sb)


def _auditar_paquete(ruta: Path) -> list[Hallazgo]:
    hallazgos: list[Hallazgo] = []
    try:
        with zipfile.ZipFile(ruta) as paquete:
            malos = paquete.testzip()
            if malos:
                hallazgos.append(_hallazgo(
                    "PPT-001", "error", f"el paquete ZIP está dañado en {malos}",
                    "regenera la presentación; no repares el XML empaquetado a mano",
                    archivo=str(ruta),
                ))
            medios: dict[str, list[str]] = {}
            for info in paquete.infolist():
                if not info.filename.startswith("ppt/media/"):
                    continue
                if info.file_size == 0:
                    hallazgos.append(_hallazgo(
                        "PPT-002", "error", f"recurso visual vacío: {info.filename}",
                        "reemplaza el recurso y vuelve a generar",
                        archivo=str(ruta),
                    ))
                    continue
                huella = hashlib.sha256(paquete.read(info.filename)).hexdigest()
                medios.setdefault(huella, []).append(info.filename)
            for nombres in medios.values():
                if len(nombres) > 1:
                    hallazgos.append(_hallazgo(
                        "PPT-003", "aviso",
                        f"el mismo recurso visual está duplicado {len(nombres)} veces",
                        "reutiliza una relación o confirma que la repetición sea intencional",
                        archivo=str(ruta),
                    ))
    except zipfile.BadZipFile:
        hallazgos.append(_hallazgo(
            "PPT-004", "error", "el archivo no es un paquete PPTX válido",
            "regenera la presentación desde el plan", archivo=str(ruta),
        ))
    return hallazgos


def auditar_presentacion(
    ruta: str | Path,
    plan: PlanVisual | None = None,
    *,
    semantica: PuertoSemantico | None = None,
) -> list[Hallazgo]:
    ruta = Path(ruta)
    hallazgos = _auditar_paquete(ruta)
    if any(h.severidad == "error" and h.codigo == "PPT-004" for h in hallazgos):
        return hallazgos
    try:
        slides = leer_presentacion(ruta)
    except Exception as exc:
        hallazgos.append(_hallazgo(
            "PPT-005", "error", f"PowerPoint no abre mediante python-pptx: {type(exc).__name__}: {exc}",
            "corrige las relaciones/OOXML o regenera el archivo",
            archivo=str(ruta),
        ))
        return hallazgos

    primarios: dict[int, Visual] = {}
    if plan:
        for numero, opciones in plan.por_diapositiva().items():
            primarios[numero] = sorted(opciones, key=lambda v: v.opcion)[0]
        if len(slides) != len(primarios):
            hallazgos.append(_hallazgo(
                "PPT-010", "error",
                f"el PPTX tiene {len(slides)} diapositivas y el plan cubre {len(primarios)}",
                "sincroniza el plan y la presentación; no aceptes láminas huérfanas",
                archivo=str(ruta),
            ))

    comparaciones: list[tuple[str, str]] = []
    comparaciones_meta: list[tuple[int, str]] = []
    for slide in slides:
        if not slide.titulo:
            hallazgos.append(_hallazgo(
                "PPT-011", "error", "la diapositiva no tiene título detectable",
                "añade un título nativo, no horneado dentro de una imagen",
                diapositiva=slide.numero, archivo=str(ruta),
            ))
        elif RE_MARCADOR.search(slide.titulo):
            hallazgos.append(_hallazgo(
                "PPT-012", "error", f"título de marcador: «{slide.titulo}»",
                "sustituye el texto de plantilla por el definitivo",
                diapositiva=slide.numero, archivo=str(ruta),
            ))

        esperado = primarios.get(slide.numero)
        requiere_visual = slide.numero > 1 and (esperado is None or esperado.tipo != "ninguno")
        if requiere_visual and slide.visuales == 0:
            hallazgos.append(_hallazgo(
                "PPT-013", "error", "diapositiva de contenido sin imagen, gráfico, tabla o diagrama",
                "incorpora una visual semánticamente justificada o marca `tipo: ninguno` con motivo",
                diapositiva=slide.numero, archivo=str(ruta),
            ))
        for imagen in slide.imagenes:
            alt = imagen.alt.strip()
            if not alt or re.fullmatch(r"(?:picture|image|imagen)\s*\d*", alt, re.I):
                hallazgos.append(_hallazgo(
                    "PPT-014", "error", f"{imagen.nombre} no tiene texto alternativo significativo",
                    "escribe qué muestra y cómo apoya la idea; registra también la fuente en el plan",
                    diapositiva=slide.numero, archivo=str(ruta),
                ))
            elif esperado:
                comparaciones.append((esperado.proposito, alt))
                comparaciones_meta.append((slide.numero, imagen.nombre))

        if slide.numero > 1 and not slide.notas:
            hallazgos.append(_hallazgo(
                "PPT-015", "aviso", "no hay notas del expositor",
                "añade notas para exposición, accesibilidad y generación posterior de audio",
                diapositiva=slide.numero, archivo=str(ruta),
            ))

        if esperado:
            if _parecido_titulo(slide.titulo, esperado.titulo) < 0.45:
                hallazgos.append(_hallazgo(
                    "PPT-020", "error",
                    f"el título «{slide.titulo}» no corresponde al plan «{esperado.titulo}»",
                    "corrige el orden o el contenido de la diapositiva",
                    diapositiva=slide.numero, archivo=str(ruta),
                ))
            plano = normalizar(slide.texto)
            for requerido in esperado.texto_visible:
                if normalizar(requerido) not in plano:
                    hallazgos.append(_hallazgo(
                        "PPT-021", "error", f"falta el texto visible obligatorio: «{requerido}»",
                        "colócalo como texto nativo editable; no dentro de una imagen generativa",
                        diapositiva=slide.numero, archivo=str(ruta),
                    ))
            if esperado.es_normativa and not esperado.procedencia.declarada():
                # Se repite deliberadamente al auditar el producto: un plan
                # inválido nunca debe quedar oculto por un PPTX que abre.
                hallazgos.append(_hallazgo(
                    "PPT-022", "error", "lámina normativa sin procedencia oficial/bibliográfica",
                    "añade y muestra una fuente normativa verificada",
                    diapositiva=slide.numero, archivo=str(ruta),
                ))

    if comparaciones:
        proveedor = semantica or SemanticaLexica()
        valores = proveedor.comparar(comparaciones)
        for (numero, nombre), valor in zip(comparaciones_meta, valores):
            if valor < 0.35:
                hallazgos.append(_hallazgo(
                    "PPT-030", "error",
                    f"{nombre} parece ajena al propósito de la diapositiva ({valor:.0%}, {proveedor.nombre})",
                    "reemplaza la imagen o mejora su texto alternativo si la relación sí es válida",
                    diapositiva=numero, archivo=str(ruta),
                ))
    return hallazgos
