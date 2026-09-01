"""Genera diapositivas a partir del mismo guion que produce el documento.

Criterio: una diapositiva por título, con lo que cuelgue de él resumido. Poco
texto por lámina — la diapositiva acompaña a quien expone, no lo sustituye.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt

from . import guion as G

ANCHO, ALTO = Cm(33.87), Cm(19.05)   # 16:9
MAX_VINETAS = 6
TINTA = RGBColor(0x1A, 0x1A, 0x1A)
SUAVE = RGBColor(0x5A, 0x5A, 0x5A)


def _lamina(pres, titulo: str):
    s = pres.slides.add_slide(pres.slide_layouts[6])   # en blanco
    caja = s.shapes.add_textbox(Cm(1.8), Cm(1.2), ANCHO - Cm(3.6), Cm(2.4))
    p = caja.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = titulo
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TINTA
    return s


def _vinetas(slide, lineas: list[str]) -> None:
    if not lineas:
        return
    caja = slide.shapes.add_textbox(Cm(1.8), Cm(4.4), ANCHO - Cm(3.6), ALTO - Cm(6))
    tf = caja.text_frame
    tf.word_wrap = True
    for i, linea in enumerate(lineas[:MAX_VINETAS]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"·  {linea}"
        r.font.size = Pt(16)
        r.font.color.rgb = TINTA
        p.space_after = Pt(10)


def _resumir(texto: str, limite: int = 160) -> str:
    """Una diapositiva no admite un párrafo entero: se queda con la primera
    idea completa y avisa con puntos suspensivos si se cortó."""
    texto = " ".join(texto.split())
    if len(texto) <= limite:
        return texto
    corte = texto.rfind(".", 0, limite)
    if corte > 40:
        return texto[:corte + 1]
    return texto[:limite].rsplit(" ", 1)[0] + "…"


def generar(guion: dict, destino: str, bibliografia: dict[str, str],
            en_texto: dict[str, str], formato: dict | None = None) -> dict:
    G.validar(guion, set(bibliografia) if bibliografia else None)

    pres = Presentation()
    pres.slide_width, pres.slide_height = ANCHO, ALTO

    portada = _lamina(pres, guion["titulo"])
    if guion.get("autor"):
        caja = portada.shapes.add_textbox(Cm(1.8), Cm(4.2), ANCHO - Cm(3.6), Cm(2))
        r = caja.text_frame.paragraphs[0].add_run()
        r.text = guion["autor"]
        r.font.size = Pt(16)
        r.font.color.rgb = SUAVE

    actual, pendientes = None, []

    def cerrar():
        if actual is not None:
            _vinetas(actual, pendientes)

    for b in guion["bloques"]:
        clase = b["clase"]
        if clase == "titulo" and b.get("nivel", 1) <= 2:
            cerrar()
            actual, pendientes = _lamina(pres, b["texto"]), []
        elif clase == "titulo":
            pendientes.append(b["texto"])
        elif clase in ("parrafo", "cita"):
            if actual is None:
                actual, pendientes = _lamina(pres, guion["titulo"]), []
            pendientes.append(_resumir(G.texto_con_citas(b, en_texto)))
        elif clase == "lista":
            if actual is None:
                actual, pendientes = _lamina(pres, guion["titulo"]), []
            pendientes.extend(_resumir(str(i), 90) for i in b["items"])
        elif clase == "tabla":
            cerrar()
            actual = _lamina(pres, b.get("leyenda") or "Datos")
            pendientes = [" · ".join(str(c) for c in f) for f in b["filas"]]
        elif clase == "figura":
            cerrar()
            actual = _lamina(pres, b.get("leyenda") or "Figura")
            pendientes = []
            ruta = Path(b["ruta"])
            if ruta.exists():
                actual.shapes.add_picture(str(ruta), Cm(4), Cm(4.6),
                                          height=ALTO - Cm(7))
        elif clase == "bibliografia":
            cerrar()
            actual = _lamina(pres, "Referencias")
            pendientes = sorted(bibliografia.values())

    cerrar()
    Path(destino).parent.mkdir(parents=True, exist_ok=True)
    pres.save(destino)
    return {"ruta": destino, "unidades": len(pres.slides.__iter__.__self__._sldIdLst)}
