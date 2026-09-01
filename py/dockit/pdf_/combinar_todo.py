#!/usr/bin/env python3
"""
Script para combinar PDFs y PPTXs en un solo PDF optimizado para impresión.
- PDFs: quita carátula, 9 por hoja (3x3), crop suave
- PPTXs: extrae todo el texto y lo coloca como texto legible (títulos en negrita)
"""

import os
import re
import glob
import tempfile
import shutil
from pypdf import PdfReader, PdfWriter, Transformation, PageObject
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch, cm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, HRFlowable
from reportlab.lib import colors
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

FOLDER = '/home/leonidas/SSD500/Develoment/tareas/matmil'
OUTPUT = os.path.join(FOLDER, 'TODO_PARA_IMPRIMIR.pdf')
TEMP_DIR = tempfile.mkdtemp(prefix='matmil_')

# Letter size in points
PAGE_W, PAGE_H = letter  # 612 x 792 points
MARGIN = 4  # points - minimal edge margin for n-up
GAP = 1  # points - minimal gap between cells


def crop_page_whitespace(page, margin_pts=30):
    """Crop white margins from a PDF page."""
    x0 = float(page.mediabox.left)
    y0 = float(page.mediabox.bottom)
    x1 = float(page.mediabox.right)
    y1 = float(page.mediabox.top)
    
    new_x0 = x0 + margin_pts
    new_y0 = y0 + margin_pts
    new_x1 = x1 - margin_pts
    new_y1 = y1 - margin_pts
    
    if new_x1 - new_x0 < 100 or new_y1 - new_y0 < 100:
        return
    
    page.mediabox.left = new_x0
    page.mediabox.bottom = new_y0
    page.mediabox.right = new_x1
    page.mediabox.top = new_y1


def place_page_in_cell(new_page, src_page, col, row, cols, rows):
    """Place a source page into a specific cell position on the output page."""
    src_w = float(src_page.mediabox.width)
    src_h = float(src_page.mediabox.height)
    src_x0 = float(src_page.mediabox.left)
    src_y0 = float(src_page.mediabox.bottom)
    
    usable_w = PAGE_W - 2 * MARGIN
    usable_h = PAGE_H - 2 * MARGIN
    cell_w = usable_w / cols
    cell_h = usable_h / rows
    
    inner_w = cell_w - GAP
    inner_h = cell_h - GAP
    scale = min(inner_w / src_w, inner_h / src_h)
    
    scaled_w = src_w * scale
    scaled_h = src_h * scale
    x_offset = MARGIN + col * cell_w + (cell_w - scaled_w) / 2
    y_offset = PAGE_H - MARGIN - (row + 1) * cell_h + (cell_h - scaled_h) / 2
    
    x_offset -= src_x0 * scale
    y_offset -= src_y0 * scale
    
    new_page.merge_transformed_page(
        src_page,
        Transformation().scale(scale).translate(x_offset, y_offset)
    )


def escape_html(text):
    """Escape HTML special characters for reportlab Paragraph."""
    text = text.replace('&', '&amp;')
    text = text.replace('<', '&lt;')
    text = text.replace('>', '&gt;')
    return text


def is_timestamp(text):
    """Check if text is just a timestamp like '11:43' or '14:59'."""
    return bool(re.match(r'^\d{1,2}:\d{2}$', text.strip()))


def is_slide_number(text):
    """Check if text is just a slide number."""
    return bool(re.match(r'^\d{1,3}$', text.strip()))


def extract_text_from_pptx(pptx_path):
    """
    Extract all text from a PPTX file, organized by slides.
    Returns list of dicts: [{title, texts: [str], tables: [[[str]]]}]
    """
    prs = Presentation(pptx_path)
    slides_data = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            'title': None,
            'texts': [],
            'tables': [],
            'slide_num': slide_idx + 1
        }
        
        # Get title
        if slide.shapes.title and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()
            if not is_timestamp(title) and not is_slide_number(title):
                slide_data['title'] = title
        
        # Process all shapes
        for shape in slide.shapes:
            # Tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.strip().replace('\n', ' ').replace('\x0b', ' ')
                        row_data.append(cell_text)
                    table_data.append(row_data)
                slide_data['tables'].append(table_data)
                continue
            
            # Text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    # Skip empty, timestamps, slide numbers
                    if not text or is_timestamp(text) or is_slide_number(text):
                        continue
                    # Skip if it's the same as the title
                    if slide_data['title'] and text == slide_data['title']:
                        continue
                    
                    # Check if bold
                    is_bold = any(run.font.bold for run in para.runs if run.font.bold is not None)
                    
                    slide_data['texts'].append({
                        'text': text,
                        'bold': is_bold
                    })
        
        # Only add slide if it has content
        if slide_data['title'] or slide_data['texts'] or slide_data['tables']:
            slides_data.append(slide_data)
    
    return slides_data


def create_slides_text_pdf(all_presentations, output_path):
    """
    Create a PDF with all slide text content, formatted for readability.
    Uses small margins and compact font to fit as much as possible.
    """
    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        leftMargin=1.2*cm,
        rightMargin=1.2*cm,
        topMargin=1*cm,
        bottomMargin=1*cm
    )
    
    # Custom styles for compact printing
    styles = getSampleStyleSheet()
    
    style_pres_title = ParagraphStyle(
        'PresTitle',
        parent=styles['Heading1'],
        fontSize=11,
        leading=13,
        spaceAfter=4,
        spaceBefore=8,
        textColor=colors.black,
        alignment=TA_CENTER,
    )
    
    style_slide_title = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading2'],
        fontSize=8,
        leading=10,
        spaceAfter=2,
        spaceBefore=4,
        textColor=colors.black,
    )
    
    style_body = ParagraphStyle(
        'SlideBody',
        parent=styles['Normal'],
        fontSize=7,
        leading=8.5,
        spaceAfter=1,
        spaceBefore=0,
    )
    
    style_bold = ParagraphStyle(
        'SlideBold',
        parent=style_body,
        fontName='Helvetica-Bold',
    )
    
    style_table_cell = ParagraphStyle(
        'TableCell',
        parent=styles['Normal'],
        fontSize=6,
        leading=7,
    )
    
    story = []
    
    for pres_name, slides_data in all_presentations:
        # Presentation header
        clean_name = os.path.splitext(pres_name)[0]
        story.append(Paragraph(f"<b>{escape_html(clean_name)}</b>", style_pres_title))
        story.append(HRFlowable(width="100%", thickness=1, color=colors.black))
        story.append(Spacer(1, 4))
        
        for slide in slides_data:
            # Slide title
            if slide['title']:
                title_text = escape_html(slide['title'])
                story.append(Paragraph(f"<b>{title_text}</b>", style_slide_title))
            
            # Slide texts
            for item in slide['texts']:
                text = escape_html(item['text'])
                if item['bold']:
                    story.append(Paragraph(f"<b>{text}</b>", style_bold))
                else:
                    story.append(Paragraph(text, style_body))
            
            # Tables
            for table_data in slide['tables']:
                if not table_data:
                    continue
                
                # Build table for reportlab
                num_cols = max(len(row) for row in table_data)
                rl_data = []
                for row in table_data:
                    rl_row = []
                    for j in range(num_cols):
                        cell_text = row[j] if j < len(row) else ''
                        rl_row.append(Paragraph(escape_html(cell_text), style_table_cell))
                    rl_data.append(rl_row)
                
                if rl_data:
                    # Calculate column widths
                    avail_width = doc.width
                    col_width = avail_width / num_cols
                    
                    t = Table(rl_data, colWidths=[col_width] * num_cols)
                    t.setStyle(TableStyle([
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                        ('FONTSIZE', (0, 0), (-1, -1), 6),
                        ('TOPPADDING', (0, 0), (-1, -1), 1),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 1),
                        ('LEFTPADDING', (0, 0), (-1, -1), 2),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 2),
                        ('BACKGROUND', (0, 0), (-1, 0), colors.Color(0.9, 0.9, 0.9)),
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ]))
                    story.append(Spacer(1, 3))
                    story.append(t)
                    story.append(Spacer(1, 3))
            
            # Small separator between slides
            story.append(Spacer(1, 2))
        
        # Separator between presentations
        story.append(Spacer(1, 6))
    
    doc.build(story)
    return output_path


def main():
    print("=" * 60)
    print("COMBINADOR DE PDFs Y PPTXs PARA IMPRESIÓN")
    print("=" * 60)
    
    # ============================================================
    # STEP 1: Collect all PDF content pages (remove covers)
    # ============================================================
    pdf_files = sorted(glob.glob(os.path.join(FOLDER, '*.pdf')))
    pdf_files = [f for f in pdf_files if os.path.basename(f) != 'TODO_PARA_IMPRIMIR.pdf']
    
    all_pdf_pages = []
    
    if pdf_files:
        print(f"\n📄 Recolectando páginas de {len(pdf_files)} PDFs (quitando carátulas)...")
        for pdf_file in pdf_files:
            name = os.path.basename(pdf_file)
            reader = PdfReader(pdf_file)
            original = len(reader.pages)
            
            for i in range(1, len(reader.pages)):
                page = reader.pages[i]
                crop_page_whitespace(page, margin_pts=30)
                all_pdf_pages.append(page)
            
            print(f"  📋 {name}: {original} págs → {original - 1} de contenido")
        
        print(f"\n  📦 Total páginas de contenido: {len(all_pdf_pages)}")
    
    # ============================================================
    # STEP 2: Extract text from all presentations
    # ============================================================
    pptx_files = sorted(glob.glob(os.path.join(FOLDER, '*.pptx')))
    
    all_presentations = []
    
    if pptx_files:
        print(f"\n📊 Extrayendo texto de {len(pptx_files)} presentaciones...")
        for pptx_file in pptx_files:
            name = os.path.basename(pptx_file)
            print(f"\n  🎯 {name}")
            
            slides_data = extract_text_from_pptx(pptx_file)
            all_presentations.append((name, slides_data))
            
            total_texts = sum(len(s['texts']) for s in slides_data)
            total_tables = sum(len(s['tables']) for s in slides_data)
            print(f"     {len(slides_data)} slides con contenido, {total_texts} párrafos, {total_tables} tablas")
    
    # ============================================================
    # STEP 3: Build combined PDF
    # ============================================================
    print(f"\n📐 Armando layout optimizado...")
    
    writer = PdfWriter()
    total_sheets = 0
    
    PDF_COLS, PDF_ROWS = 3, 3
    PDF_PER_SHEET = PDF_COLS * PDF_ROWS  # 9
    
    # --- Phase 1: PDF pages in 3x3 grid ---
    if all_pdf_pages:
        for sheet_start in range(0, len(all_pdf_pages), PDF_PER_SHEET):
            new_page = PageObject.create_blank_page(width=PAGE_W, height=PAGE_H)
            
            end = min(sheet_start + PDF_PER_SHEET, len(all_pdf_pages))
            pdf_count_on_sheet = end - sheet_start
            
            for idx in range(pdf_count_on_sheet):
                src_page = all_pdf_pages[sheet_start + idx]
                col = idx % PDF_COLS
                row = idx // PDF_COLS
                place_page_in_cell(new_page, src_page, col, row, PDF_COLS, PDF_ROWS)
            
            writer.add_page(new_page)
            total_sheets += 1
    
    pdf_sheets = total_sheets
    print(f"  📄 PDFs: {pdf_sheets} hojas (9/hoja, {len(all_pdf_pages)} páginas)")
    
    # --- Phase 2: Slide text pages ---
    if all_presentations:
        slides_pdf_path = os.path.join(TEMP_DIR, 'slides_text.pdf')
        create_slides_text_pdf(all_presentations, slides_pdf_path)
        
        slides_reader = PdfReader(slides_pdf_path)
        slide_text_pages = len(slides_reader.pages)
        
        for page in slides_reader.pages:
            writer.add_page(page)
            total_sheets += 1
        
        print(f"  📊 Diapos (texto): {slide_text_pages} hojas de texto legible")
    
    # ============================================================
    # STEP 4: Save final PDF
    # ============================================================
    with open(OUTPUT, 'wb') as f:
        writer.write(f)
    
    print(f"\n{'=' * 60}")
    print(f"✅ ¡LISTO! PDF final guardado en:")
    print(f"   {OUTPUT}")
    print(f"   Total de hojas a imprimir: {total_sheets}")
    print(f"{'=' * 60}")
    
    # Cleanup
    shutil.rmtree(TEMP_DIR, ignore_errors=True)
    print("\n🧹 Archivos temporales limpiados.")


if __name__ == '__main__':
    main()
