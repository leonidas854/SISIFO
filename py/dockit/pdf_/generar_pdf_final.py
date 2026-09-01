from pathlib import Path
#!/usr/bin/env python3
"""
Genera UN SOLO PDF listo para imprimir, con:

  PARTE I   - Los PDFs de los trabajos de los equipos tal cual se ven,
              recortados a su contenido y apilados en 3 columnas por hoja.
              Descarta caratulas, indices y hojas en blanco.
  PARTE II  - Todo el contenido de las 3 presentaciones (.pptx) como texto.
              Las diapositivas que son FOTOS fueron transcritas a mano
              (archivo ocr.txt); las que tienen texto real se leen del pptx.
  PARTE III - Los mismos trabajos de los equipos pero como texto legible.
              DESACTIVADA (ver INCLUIR_TP_TEXTO): repetia lo de la PARTE I.

Nada de este PDF es contenido propio: la PARTE I y la III salen de los
TP-*.pdf que ya estaban en la carpeta, y la PARTE II de los .pptx del docente.

Diseno del texto: 3 columnas, letra chica pero legible, titulos en negrita,
margenes minimos, sin fondos de color -> minima tinta, minimas hojas.

Ajustes: N_COLS (columnas del texto), NUP_COLS (columnas de la PARTE I),
INCLUIR_TP_TEXTO (volver a agregar la PARTE III).
"""

import os
import re
import glob
import tempfile
import unicodedata

from pypdf import PdfReader, PdfWriter, PageObject, Transformation
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import cm
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER, TA_LEFT
from reportlab.lib import colors
from reportlab.platypus import (
    BaseDocTemplate, Frame, PageTemplate, Paragraph, Spacer,
    KeepTogether, HRFlowable,
)

def _carpeta_trabajo() -> str:
    """Carpeta con los archivos a procesar: TALLER_PROYECTO, o donde estés."""
    if _v := os.environ.get("TALLER_PROYECTO"):
        return str(Path(_v).resolve())
    return str(Path.cwd().resolve())


FOLDER = _carpeta_trabajo()
OUTPUT = os.path.join(FOLDER, 'MATERIA-COMPLETA-IMPRIMIR.pdf')

# ocr.txt = transcripcion manual de las diapositivas que son imagenes
OCR_FILE = os.path.join(FOLDER, 'ocr.txt')

PAGE_W, PAGE_H = letter
MARGIN_X = 0.8 * cm
MARGIN_TOP = 0.8 * cm
MARGIN_BOT = 0.7 * cm
COL_GAP = 0.42 * cm
N_COLS = 3            # columnas del texto (2 = letra mas grande, 3 = mas denso)

DOC_TITLE = 'LIDERAZGO E INTELIGENCIA ESTRATEGICA  -  EMI 2026'

# --- PARTE I: los PDFs de los equipos, acomodados en columnas -------------
# No es una grilla rigida: cada pagina se recorta a su contenido real y se
# van apilando en columnas, asi no quedan huecos blancos entre una y otra.
NUP_COLS = 3                  # columnas por hoja
NUP_MARGIN = 8                # borde de la hoja, en puntos
NUP_GAP = 6                   # separacion entre paginas, en puntos
NUP_PAD = 3                   # aire alrededor del contenido recortado
BBOX_DPI = 72                 # resolucion para detectar donde hay tinta
NUP_SHRINK = 0.86             # cuanto se puede achicar una pagina con tal de
                              # que entre en la columna y no quede hueco

# Repetir los trabajos de los equipos tambien como texto legible (PARTE III)
INCLUIR_TP_TEXTO = False

# Presentaciones en orden, con el nombre que se muestra en el PDF
PPTX_ORDER = [
    ('1', '1-LIDERAZGO-A-24-jul-26-EE-CC.pptx',
     'CLASE 1 (24-JUL) - PRESENTACION, PLAN DE TRABAJO Y EVALUACION DIAGNOSTICA'),
    ('2', '2-EXPOSICION-JUE-30-jul-26.pptx',
     'CLASE 2 (30-JUL) - DOCTRINA DE LIDERAZGO: CONCEPTO, PODER Y TEORIAS'),
    ('3', '3-LIDERAZGO-A-VIE-21-AGO-26-EE-CC.pptx',
     'CLASE 3 (21-AGO) - FACTORES, PRINCIPIOS, COMPETENCIAS Y VALORES MILITARES'),
]

# Rotulos sobrepuestos en las fotos, ya incorporados a la transcripcion
OVERLAY_LABELS = {'EL MAS IMPORTANTE', 'ES EL MAS IMPORTANTE'}

# Diapositivas que son solo caratula/decoracion: no aportan contenido
SKIP_SLIDES = {
    'p1_s01', 'p1_s37',
    'p2_s01', 'p2_s39', 'p2_s16', 'p2_s17',
    'p3_s01', 'p3_s44',
}


# --------------------------------------------------------------------------
# Estilos
# --------------------------------------------------------------------------

def build_styles():
    s = {}
    # A 3 columnas la columna es angosta: letra mas chica y sin justificar,
    # para que no queden huecos feos entre palabras.
    tight = N_COLS >= 3
    body_size = 6.5 if tight else 7.4
    body_lead = 7.7 if tight else 8.7
    align = TA_LEFT if tight else TA_JUSTIFY

    s['part'] = ParagraphStyle(
        'Part', fontName='Helvetica-Bold', fontSize=9.2, leading=11,
        spaceBefore=1, spaceAfter=2.5, alignment=TA_CENTER, textColor=colors.black,
    )
    s['section'] = ParagraphStyle(
        'Section', fontName='Helvetica-Bold', fontSize=7.6, leading=9.2,
        spaceBefore=5, spaceAfter=2.2, textColor=colors.black,
    )
    s['slide'] = ParagraphStyle(
        'Slide', fontName='Helvetica-Bold', fontSize=6.9, leading=8.3,
        spaceBefore=3.6, spaceAfter=1.2, textColor=colors.black,
    )
    s['sub'] = ParagraphStyle(
        'Sub', fontName='Helvetica-Bold', fontSize=6.4, leading=7.7,
        spaceBefore=1.8, spaceAfter=0.6, leftIndent=2, textColor=colors.black,
    )
    s['body'] = ParagraphStyle(
        'Body', fontName='Times-Roman', fontSize=body_size, leading=body_lead,
        spaceBefore=0, spaceAfter=1.0, alignment=align,
    )
    s['bullet'] = ParagraphStyle(
        'Bullet', parent=s['body'], leftIndent=6.5, firstLineIndent=-5,
        spaceAfter=0.7,
    )
    return s


# --------------------------------------------------------------------------
# Utilidades de texto
# --------------------------------------------------------------------------

def esc(t):
    return (t.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;'))


def clean(t):
    """Normaliza espacios y basura tipica de extraccion de PDF."""
    t = t.replace('\x0b', ' ').replace('\xa0', ' ')
    t = re.sub(r'\.{4,}', ' ', t)          # puntos guia del indice
    t = re.sub(r'(\s\.){3,}\s*', ' ', t)   # ". . . ." separados
    t = re.sub(r'[ \t]+', ' ', t)
    return t.strip()


def strip_accents(t):
    return ''.join(c for c in unicodedata.normalize('NFD', t)
                   if unicodedata.category(c) != 'Mn')


# --------------------------------------------------------------------------
# Fuente 1: texto nativo de los .pptx
# --------------------------------------------------------------------------

def read_pptx_text():
    """Devuelve {slide_key: [(tipo, texto), ...]} con el texto real del pptx."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(sh.shapes)
            else:
                yield sh

    junk = re.compile(r'^(\d{1,2}:\d{2}|\d{1,3}|[a-e]\)?)$')
    data = {}

    for tag, fname, _ in PPTX_ORDER:
        path = os.path.join(FOLDER, fname)
        if not os.path.exists(path):
            continue
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, 1):
            key = f'p{tag}_s{i:02d}'
            title = ''
            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title = slide.shapes.title.text.strip()
            except (AttributeError, ValueError):
                title = ''

            items = []
            for sh in walk(slide.shapes):
                if getattr(sh, 'has_table', False) and sh.has_table:
                    for row in sh.table.rows:
                        cells = [clean(c.text) for c in row.cells]
                        cells = [c for c in cells if c]
                        if cells:
                            items.append(('bullet', ' | '.join(cells)))
                    continue
                if not sh.has_text_frame:
                    continue
                for para in sh.text_frame.paragraphs:
                    txt = clean(para.text)
                    if not txt or junk.match(txt):
                        continue
                    kind = 'title' if txt == title else 'bullet'
                    items.append((kind, txt))

            if items:
                data[key] = items
    return data


# --------------------------------------------------------------------------
# Fuente 2: transcripcion de las diapositivas que son fotos (ocr.txt)
# --------------------------------------------------------------------------

def read_ocr_file():
    """
    Formato de ocr.txt:
        @@@ p3_s07
        # Titulo o subtitulo (negrita)
        * vineta
        texto de parrafo normal
    """
    data = {}
    if not os.path.exists(OCR_FILE):
        return data

    key = None
    for raw in open(OCR_FILE, encoding='utf-8'):
        line = raw.rstrip('\n')
        if line.startswith('@@@ '):
            key = line[4:].strip()
            data[key] = []
            continue
        if key is None:
            continue
        line = clean(line)
        if not line:
            continue
        if line.startswith('# '):
            data[key].append(('title', line[2:].strip()))
        elif line.startswith('* '):
            data[key].append(('bullet', line[2:].strip()))
        else:
            data[key].append(('body', line))
    return {k: v for k, v in data.items() if v}


# --------------------------------------------------------------------------
# Fuente 3: texto de los PDFs de los equipos
# --------------------------------------------------------------------------

# Lineas de caratula / indice / nomina que no aportan contenido de estudio
NOISE_LINE = re.compile(
    r'^(comando general|escuela militar|mcal\.|bolivia$|cochabamba|'
    r'carrera\s*:|semestre\s*:|codigos y nombres|estudiantes\s*:|'
    r'asignatura\s*:|u\. academica|gestion\s*:|docente\s*:|fecha\s*:|'
    r'indice$|indice general|pagina \d+|\W*$)',
    re.IGNORECASE)

# Nombre de estudiante seguido (o precedido) de su codigo: C11424-3
ROSTER_LINE = re.compile(r'C\d{5}\s*-\s*\d')

# Linea del indice: "1.1. DESCRIPCION DEL CASO .......... 4"
TOC_LINE = re.compile(r'\.{4,}\s*\d+\s*$|\s\.(\s\.){3,}')

HEAD_LINE = re.compile(
    r'^(\d+(\.\d+)*\.?\s+\S|mision|trabajo pract|conclusion|recomendacion|'
    r'analisis|descripcion|ejemplo|caso \d|situacion|contexto|justificacion|'
    r'principio|valor militar|estado)', re.IGNORECASE)


def looks_like_index(page_text):
    """Detecta la pagina de INDICE (lineas que terminan en numero de pagina)."""
    lines = [l.strip() for l in page_text.split('\n') if len(l.strip()) > 10]
    if not lines or len(page_text) > 2500:
        return False
    hits = sum(1 for l in lines
               if re.search(r'\.{3,}', l) or re.search(r'\s\d{1,3}$', l))
    return hits >= max(3, len(lines) * 0.5)


def pdf_lines(path):
    """
    Texto del PDF via pdftotext (mucho mas limpio que pypdf).
    Descarta la caratula (pag. 1) y la pagina de indice.
    """
    import subprocess
    try:
        raw = subprocess.run(['pdftotext', '-q', path, '-'],
                             capture_output=True, text=True, timeout=120).stdout
        pages = raw.split('\f')
    except (OSError, subprocess.SubprocessError):
        pages = [(p.extract_text() or '') for p in PdfReader(path).pages]

    lines = []
    for n, page in enumerate(pages):
        if n == 0:                       # caratula
            continue
        if looks_like_index(page):       # indice
            continue
        lines.extend(page.split('\n'))
    return lines


def read_tp_pdfs():
    """Devuelve [(nombre_equipo, [(tipo, texto), ...]), ...]."""
    files = sorted(glob.glob(os.path.join(FOLDER, 'TP-*.pdf')))
    out = []

    for path in files:
        name = re.sub(r'\s+', ' ',
                      os.path.splitext(os.path.basename(path))[0]).upper()

        kept = []
        for line in pdf_lines(path):
            line = clean(line)
            if not line or len(line) < 3:
                continue
            plain = strip_accents(line)
            if NOISE_LINE.match(plain) or ROSTER_LINE.search(line):
                continue
            if TOC_LINE.search(line):
                continue
            if line in ('•', '-', '*', '', 'o'):
                continue
            kept.append(line)

        # Reune el parrafo partido en varias lineas: se corta solo cuando la
        # linea anterior termina en punto/dos puntos y la siguiente arranca
        # en mayuscula o vineta.
        merged = []
        for line in kept:
            starts_new = (
                re.match(r'^[•\-•▪■]', line)
                or re.match(r'^\d+(\.\d+)*[.)]\s', line)
                or (line[:1].isupper()
                    and merged and re.search(r'[.:;!?"”)]$', merged[-1]))
            )
            if merged and not starts_new:
                merged[-1] = merged[-1] + ' ' + line
            else:
                merged.append(line)

        blocks = []
        for line in merged:
            line = re.sub(r'^[•\-•▪■]\s*', '', line).strip()
            if not line:
                continue
            kind = 'sub' if (HEAD_LINE.match(line) and len(line) < 120) else 'body'
            blocks.append((kind, line))

        if blocks:
            out.append((name, blocks))
    return out


# --------------------------------------------------------------------------
# PARTE I: los PDFs de los equipos en grilla 3x3
# --------------------------------------------------------------------------

def render_pages(path, dpi=BBOX_DPI):
    """Renderiza el PDF a imagenes (una por pagina) en un directorio temporal."""
    import subprocess
    from PIL import Image

    tmp = tempfile.mkdtemp(prefix='matmil_bbox_')
    subprocess.run(['pdftoppm', '-r', str(dpi), '-gray', '-png',
                    path, os.path.join(tmp, 'p')],
                   check=True, capture_output=True, timeout=300)
    imgs = []
    for f in sorted(glob.glob(os.path.join(tmp, 'p-*.png'))):
        with Image.open(f) as im:
            imgs.append(im.convert('L').copy())
    return imgs, tmp


def content_box(img, page):
    """
    Devuelve el rectangulo (en puntos PDF) donde realmente hay contenido.
    None si la pagina esta en blanco.
    """
    # Todo lo que no sea casi blanco cuenta como tinta
    ink = img.point(lambda v: 255 if v < 245 else 0)
    bbox = ink.getbbox()
    if not bbox:
        return None

    left, upper, right, lower = bbox
    k = 72.0 / BBOX_DPI
    box = page.mediabox
    x0 = float(box.left) + left * k - NUP_PAD
    x1 = float(box.left) + right * k + NUP_PAD
    y1 = float(box.top) - upper * k + NUP_PAD
    y0 = float(box.top) - lower * k - NUP_PAD

    # Nunca salir de la pagina original
    x0 = max(x0, float(box.left));   x1 = min(x1, float(box.right))
    y0 = max(y0, float(box.bottom)); y1 = min(y1, float(box.top))
    if x1 - x0 < 20 or y1 - y0 < 20:
        return None
    return x0, y0, x1, y1


def collect_tp_pages():
    """
    Paginas utiles de los PDFs de los equipos, ya recortadas a su contenido.
    Descarta caratulas, indices y hojas en blanco.
    """
    import shutil
    import subprocess

    out = []
    descartadas = {'caratula': 0, 'indice': 0, 'vacia': 0}

    for path in sorted(glob.glob(os.path.join(FOLDER, 'TP-*.pdf'))):
        reader = PdfReader(path)
        raw = subprocess.run(['pdftotext', '-q', path, '-'],
                             capture_output=True, text=True).stdout
        texts = raw.split('\f')
        imgs, tmp = render_pages(path)
        try:
            for i, page in enumerate(reader.pages):
                if i == 0:
                    descartadas['caratula'] += 1
                    continue
                if i < len(texts) and looks_like_index(texts[i]):
                    descartadas['indice'] += 1
                    continue
                if i >= len(imgs):
                    continue
                box = content_box(imgs[i], page)
                if box is None:
                    descartadas['vacia'] += 1
                    continue
                x0, y0, x1, y1 = box
                page.mediabox.left, page.mediabox.bottom = x0, y0
                page.mediabox.right, page.mediabox.top = x1, y1
                out.append(page)
        finally:
            shutil.rmtree(tmp, ignore_errors=True)

    return out, descartadas


def build_nup_sheets():
    """
    Apila las paginas recortadas en columnas, una debajo de otra, saltando de
    columna (y de hoja) recien cuando no entra la siguiente. Asi no quedan
    celdas vacias ni margenes blancos heredados de los originales.
    """
    pages, descartadas = collect_tp_pages()

    col_w = (PAGE_W - 2 * NUP_MARGIN - NUP_GAP * (NUP_COLS - 1)) / NUP_COLS
    avail_h = PAGE_H - 2 * NUP_MARGIN
    dims = [(float(p.mediabox.width), float(p.mediabox.height)) for p in pages]

    def columnas_con(factor):
        """Cuantas columnas ocupa todo si cada pagina se achica ese factor."""
        n, used = 0, 0.0
        for w, h in dims:
            ph = h * min(col_w * factor / w, avail_h / h)
            gap = NUP_GAP if used > 0 else 0.0
            if used + gap + ph > avail_h:
                n, used, gap = n + 1, 0.0, 0.0
            used += gap + ph
        return n + 1

    # Busca el tamano mas grande que igual entre en la menor cantidad de hojas.
    # Achicar un 8% puede ahorrar una hoja entera, y casi no se nota.
    opciones = [round(1.0 - i * 0.01, 2) for i in range(int(NUP_SHRINK * 100))]
    opciones = [f for f in opciones if f >= NUP_SHRINK]
    hojas_por_factor = {f: -(-columnas_con(f) // NUP_COLS) for f in opciones}
    minimo = min(hojas_por_factor.values())
    factor = max(f for f, n in hojas_por_factor.items() if n == minimo)
    col_full, col_w = col_w, col_w * factor

    sheets = []
    sheet = PageObject.create_blank_page(width=PAGE_W, height=PAGE_H)
    col, used = 0, 0.0

    for src in pages:
        box = src.mediabox
        w, h = float(box.width), float(box.height)
        scale = min(col_w / w, avail_h / h)
        ph = h * scale

        gap = NUP_GAP if used > 0 else 0.0
        if used + gap + ph > avail_h:        # no entra: siguiente columna
            col += 1
            used, gap = 0.0, 0.0
            if col >= NUP_COLS:
                sheets.append(sheet)
                sheet = PageObject.create_blank_page(width=PAGE_W, height=PAGE_H)
                col = 0
        used += gap

        x_left = (NUP_MARGIN + col * (col_full + NUP_GAP)
                  + (col_full - w * scale) / 2)
        y_bottom = PAGE_H - NUP_MARGIN - used - ph

        sheet.merge_transformed_page(
            src,
            Transformation().scale(scale).translate(
                x_left - float(box.left) * scale,
                y_bottom - float(box.bottom) * scale))
        used += ph

    sheets.append(sheet)
    return sheets, len(pages), descartadas


# --------------------------------------------------------------------------
# Armado del documento
# --------------------------------------------------------------------------

def render_items(story, items, st):
    """Convierte [(tipo, texto)] en flowables, agrupando titulo+1ra linea."""
    pending = []
    for kind, txt in items:
        txt = esc(txt)
        if kind == 'title':
            pending.append(Paragraph(f'<b>{txt}</b>', st['sub']))
        elif kind == 'sub':
            pending.append(Paragraph(f'<b>{txt}</b>', st['sub']))
        elif kind == 'bullet':
            pending.append(Paragraph(f'&bull; {txt}', st['bullet']))
        else:
            pending.append(Paragraph(txt, st['body']))

        # Evita que un titulo quede solo al pie de la columna
        if kind in ('title', 'sub'):
            continue
        if pending:
            story.append(KeepTogether(pending) if len(pending) > 1 else pending[0])
            pending = []
    if pending:
        story.append(KeepTogether(pending) if len(pending) > 1 else pending[0])


def build_story(st):
    pptx_text = read_pptx_text()
    ocr_text = read_ocr_file()
    story = []

    # ---------------- PARTE II: las 3 presentaciones ----------------
    story.append(Paragraph('PARTE II - CONTENIDO DE LAS 3 PRESENTACIONES', st['part']))
    story.append(HRFlowable(width='100%', thickness=1.1, color=colors.black,
                            spaceBefore=1, spaceAfter=4))

    n_slides = 0
    for tag, fname, label in PPTX_ORDER:
        story.append(Paragraph(esc(label), st['section']))
        story.append(HRFlowable(width='100%', thickness=0.6, color=colors.black,
                                spaceBefore=0, spaceAfter=3))

        keys = sorted(k for k in set(pptx_text) | set(ocr_text)
                      if k.startswith(f'p{tag}_s'))
        for key in keys:
            if key in SKIP_SLIDES:
                continue

            ocr_items = ocr_text.get(key, [])
            nat_items = pptx_text.get(key, [])

            # Si la diapositiva es una foto, manda la transcripcion; el texto
            # nativo (rotulos sobrepuestos) se agrega solo si aporta algo.
            if ocr_items:
                seen = {strip_accents(t).upper() for _, t in ocr_items}
                seen |= OVERLAY_LABELS
                extra = [(k, t) for k, t in nat_items
                         if strip_accents(t).upper() not in seen and len(t) > 4]
                items = ocr_items + extra
            else:
                items = nat_items

            if not items:
                continue

            # La primera linea de la diapositiva es su titulo
            head_kind, head_txt = items[0]
            body = items[1:]
            n_slides += 1

            head = Paragraph(f'<b>{esc(head_txt)}</b>', st['slide'])
            if body:
                first_kind, first_txt = body[0]
                first_style = st['bullet'] if first_kind == 'bullet' else st['body']
                first_txt = (f'&bull; {esc(first_txt)}' if first_kind == 'bullet'
                             else esc(first_txt))
                story.append(KeepTogether([head, Paragraph(first_txt, first_style)]))
                render_items(story, body[1:], st)
            else:
                story.append(head)

    # ---------------- PARTE III: trabajos de los equipos, en texto ----------
    tps = read_tp_pdfs() if INCLUIR_TP_TEXTO else []
    if tps:
        story.append(Spacer(1, 6))
        story.append(Paragraph(
            'PARTE III - TRABAJOS DE LOS EQUIPOS (TEXTO LEGIBLE)', st['part']))
        story.append(HRFlowable(width='100%', thickness=1.1, color=colors.black,
                                spaceBefore=1, spaceAfter=4))
        for name, blocks in tps:
            story.append(Paragraph(esc(name), st['section']))
            story.append(HRFlowable(width='100%', thickness=0.6, color=colors.black,
                                    spaceBefore=0, spaceAfter=3))
            render_items(story, blocks, st)

    return story, n_slides, len(tps)


def make_page_decorator(offset=0):
    def draw(canvas, doc):
        canvas.saveState()
        canvas.setFont('Helvetica', 6)
        canvas.setFillGray(0.35)
        canvas.drawString(MARGIN_X, PAGE_H - MARGIN_TOP + 4, DOC_TITLE)
        canvas.drawRightString(PAGE_W - MARGIN_X, PAGE_H - MARGIN_TOP + 4,
                               f'pag. {doc.page + offset}')
        canvas.setStrokeGray(0.6)
        canvas.setLineWidth(0.4)
        canvas.line(MARGIN_X, PAGE_H - MARGIN_TOP + 1.5,
                    PAGE_W - MARGIN_X, PAGE_H - MARGIN_TOP + 1.5)
        canvas.restoreState()
    return draw


def main():
    # --- PARTE I: los PDFs de los equipos, recortados y apilados ---
    sheets, n_orig_pages, desc = build_nup_sheets()
    print(f'PARTE I  : {n_orig_pages} paginas utiles -> {len(sheets)} hojas '
          f'de {NUP_COLS} columnas')
    print(f'           descartadas: {desc["caratula"]} caratulas, '
          f'{desc["indice"]} indices, {desc["vacia"]} en blanco')

    # --- PARTES II y III: todo el texto ---
    st = build_styles()
    story, n_slides, n_tps = build_story(st)

    tmp_text = os.path.join(tempfile.gettempdir(), 'matmil_texto.pdf')
    doc = BaseDocTemplate(
        tmp_text, pagesize=letter,
        leftMargin=MARGIN_X, rightMargin=MARGIN_X,
        topMargin=MARGIN_TOP, bottomMargin=MARGIN_BOT,
        title='Liderazgo e Inteligencia Estrategica - Material completo',
        author='EMI 2026',
    )

    usable_w = PAGE_W - 2 * MARGIN_X
    usable_h = PAGE_H - MARGIN_TOP - MARGIN_BOT - 8   # 8 pt para el encabezado
    col_w = (usable_w - COL_GAP * (N_COLS - 1)) / N_COLS

    frames = [
        Frame(MARGIN_X + i * (col_w + COL_GAP), MARGIN_BOT, col_w, usable_h,
              leftPadding=0, rightPadding=0, topPadding=0, bottomPadding=0)
        for i in range(N_COLS)
    ]
    doc.addPageTemplates([
        PageTemplate(id='dos-columnas', frames=frames,
                     onPage=make_page_decorator(offset=len(sheets)))
    ])

    doc.build(story)

    # --- Union final: grilla 3x3 + texto, en un solo PDF ---
    writer = PdfWriter()
    for sheet in sheets:
        writer.add_page(sheet)
    text_pdf = PdfReader(tmp_text)
    for page in text_pdf.pages:
        writer.add_page(page)
    with open(OUTPUT, 'wb') as fh:
        writer.write(fh)
    os.remove(tmp_text)

    total = len(sheets) + len(text_pdf.pages)
    print(f'PARTE II : {n_slides} diapositivas de las 3 presentaciones')
    print(f'PARTE III: {n_tps} trabajos de equipos en texto')
    print('=' * 58)
    print('LISTO. PDF generado:')
    print(f'  {OUTPUT}')
    print(f'  Hojas a imprimir: {total} '
          f'({len(sheets)} de grilla + {len(text_pdf.pages)} de texto)')
    print('=' * 58)


if __name__ == '__main__':
    main()
