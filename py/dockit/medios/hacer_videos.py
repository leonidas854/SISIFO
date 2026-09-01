#!/usr/bin/env python3
"""Arma un video por tema: cada diapositiva de `diapos_terminados` mostrada
mientras suena su narracion de `audios_diapos_original`.

Pasos (todos reanudables; se saltan los productos ya hechos salvo --rehacer):

  1. mapear  - empareja diapositiva <-> audio por similitud de texto y deja
               salida/mapeo.json. Los decks terminados perdieron duplicados
               (temas 4 y 9) y el tema 14 tiene dos diapos intercambiadas,
               asi que el emparejamiento NO puede ser 1 a 1 por posicion.
  2. render  - pptx -> pdf (LibreOffice) -> png 1920x1080 (pdftoppm).
               El tema 9 ya es pdf y se usa tal cual.
  3. audio   - concatena los mp3 del tema con una pausa exacta entre diapos
               y los codifica a AAC (una sola pista .m4a por tema).
  4. video   - concat demuxer de imagenes con la duracion exacta de cada
               narracion + la pista de audio -> mp4 H.264.

Nada se escribe fuera de videos_diapos/: los originales no se tocan.

Uso:
    python3 hacer_videos.py                  # los 14 temas
    python3 hacer_videos.py --temas 4 9 14   # solo esos
    python3 hacer_videos.py --paso mapear    # solo el emparejamiento
    python3 hacer_videos.py --crf 18 --fps 30 --pausa 0.8
"""

from __future__ import annotations

import argparse
import difflib
import json
import re
import shutil
import subprocess
import sys
import time
from pathlib import Path

RAIZ = Path(__file__).resolve().parents[2]
DIAPOS = RAIZ / "diapos_terminados"
ORIGINALES = RAIZ / "diapos_original"
AUDIOS = RAIZ / "audios_diapos_original"
BASE = RAIZ / "videos_diapos"
TRABAJO = BASE / "trabajo"
SALIDA = BASE / "salida"

TEMAS = list(range(1, 15))
ANCHO, ALTO = 1920, 1080


# --------------------------------------------------------------------------- util


def log(msg: str) -> None:
    print(msg, flush=True)


def correr(cmd: list[str], **kw) -> subprocess.CompletedProcess:
    r = subprocess.run(cmd, capture_output=True, text=True, **kw)
    if r.returncode != 0:
        cola = (r.stderr or r.stdout or "").strip().splitlines()[-15:]
        raise RuntimeError(f"fallo: {' '.join(cmd[:4])} ...\n" + "\n".join(cola))
    return r


def normalizar(t: str) -> str:
    return re.sub(r"\s+", " ", t).strip().lower()


def archivo_deck(tema: int) -> Path:
    for p in sorted(DIAPOS.iterdir()):
        if p.name.split("_")[0] == str(tema) and p.suffix.lower() in (".pptx", ".pdf"):
            return p
    raise FileNotFoundError(f"no hay deck para el tema {tema} en {DIAPOS}")


def carpeta_audio(tema: int) -> Path:
    return AUDIOS / f"presentacion_{tema:02d}"


def audios_del_tema(tema: int) -> list[Path]:
    return sorted(carpeta_audio(tema).glob("diapositiva_*.mp3"))


# ----------------------------------------------------------------------- 1. mapear


def textos_diapositivas(tema: int) -> list[str]:
    """Texto de cada diapositiva del deck terminado, en orden."""
    deck = archivo_deck(tema)
    if deck.suffix.lower() == ".pdf":
        crudo = correr(["pdftotext", "-layout", str(deck), "-"]).stdout
        paginas = crudo.split("\f")
        if paginas and not paginas[-1].strip():
            paginas.pop()
        return [normalizar(p) for p in paginas]
    from pptx import Presentation  # solo se necesita aqui

    pres = Presentation(str(deck))
    salida = []
    for slide in pres.slides:
        partes = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
        salida.append(normalizar(" ".join(partes)))
    return salida


def textos_originales(tema: int) -> list[str]:
    """Texto de cada diapositiva del deck ORIGINAL, en orden.

    El original es el puente fiable: tiene exactamente tantas diapositivas como
    audios (los audios se generaron de ahi, uno por diapositiva), mientras que
    el terminado perdio duplicados y reordeno alguna. Emparejar terminado con
    original es facil porque el texto es casi identico; emparejarlo con las
    transcripciones no lo es, porque los TXT anaden descripciones de imagenes.
    """
    from pptx import Presentation

    for p in sorted(ORIGINALES.iterdir()):
        if p.name.split("_")[0] == str(tema) and p.suffix.lower() == ".pptx":
            pres = Presentation(str(p))
            return [
                normalizar(" ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame))
                for s in pres.slides
            ]
    raise FileNotFoundError(f"no hay deck original para el tema {tema}")


def mapear_tema(tema: int) -> dict:
    diapos = textos_diapositivas(tema)
    originales = textos_originales(tema)
    mp3s = audios_del_tema(tema)

    if len(originales) != len(mp3s):
        raise RuntimeError(
            f"tema {tema}: el deck original tiene {len(originales)} diapositivas "
            f"y hay {len(mp3s)} audios; el puente deja de valer"
        )

    n_d, n_o = len(diapos), len(originales)
    filas, usados = [], []
    for i, td in enumerate(diapos):
        candidatos = []
        for j, to in enumerate(originales):
            sim = difflib.SequenceMatcher(None, td, to).ratio()
            # entre diapositivas duplicadas del original gana la de posicion
            # relativa mas cercana.
            castigo = 0.10 * abs(i / max(n_d - 1, 1) - j / max(n_o - 1, 1))
            candidatos.append((sim - castigo, sim, j))
        _, sim, j = max(candidatos)
        usados.append(j)
        filas.append(
            {
                "diapositiva": i + 1,
                "audio": j + 1,
                "mp3": str(mp3s[j].relative_to(AUDIOS)),
                "similitud": round(sim, 3),
            }
        )

    sobrantes = [k + 1 for k in range(n_o) if k not in usados]
    repetidos = sorted({j + 1 for j in usados if usados.count(j) > 1})
    dudosos = [f["diapositiva"] for f in filas if f["similitud"] < 0.75]
    return {
        "tema": tema,
        "deck": archivo_deck(tema).name,
        "diapositivas": n_d,
        "audios": len(mp3s),
        "audios_sin_usar": sobrantes,
        "audios_repetidos": repetidos,
        "diapositivas_dudosas": dudosos,
        "similitud_minima": round(min(f["similitud"] for f in filas), 3),
        "pares": filas,
    }


def paso_mapear(temas: list[int], rehacer: bool) -> dict:
    destino = SALIDA / "mapeo.json"
    mapeo = json.loads(destino.read_text()) if destino.exists() and not rehacer else {}
    for tema in temas:
        clave = str(tema)
        if clave in mapeo and not rehacer:
            continue
        log(f"  mapeando tema {tema:2d} ...")
        mapeo[clave] = mapear_tema(tema)
        m = mapeo[clave]
        aviso = []
        if m["audios_sin_usar"]:
            aviso.append(f"audios sin usar: {m['audios_sin_usar']}")
        if m["audios_repetidos"]:
            aviso.append(f"audios repetidos: {m['audios_repetidos']}")
        if m["diapositivas_dudosas"]:
            aviso.append(f"emparejamiento dudoso en diapos {m['diapositivas_dudosas']}")
        detalle = "; ".join(aviso) or "1 a 1 limpio"
        log(f"     {m['diapositivas']} diapos / {m['audios']} audios - {detalle}")
    destino.parent.mkdir(parents=True, exist_ok=True)
    destino.write_text(json.dumps(mapeo, ensure_ascii=False, indent=1), encoding="utf-8")
    return mapeo


# ----------------------------------------------------------------------- 2. render


def paso_render(tema: int, rehacer: bool) -> list[Path]:
    deck = archivo_deck(tema)
    dir_png = TRABAJO / "png" / f"tema_{tema:02d}"
    esperadas = len(textos_diapositivas(tema))

    hechas = sorted(dir_png.glob("diapo_*.png"))
    if len(hechas) == esperadas and not rehacer:
        return hechas

    if rehacer and dir_png.exists():
        shutil.rmtree(dir_png)
    dir_png.mkdir(parents=True, exist_ok=True)

    if deck.suffix.lower() == ".pdf":
        pdf = deck
    else:
        dir_pdf = TRABAJO / "pdf"
        dir_pdf.mkdir(parents=True, exist_ok=True)
        pdf = dir_pdf / f"tema_{tema:02d}.pdf"
        if not pdf.exists() or rehacer:
            perfil = (TRABAJO / "perfil_libreoffice").as_uri()
            correr(
                [
                    "libreoffice",
                    f"-env:UserInstallation={perfil}",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(dir_pdf),
                    str(deck),
                ]
            )
            generado = dir_pdf / (deck.stem + ".pdf")
            generado.replace(pdf)

    correr(
        [
            "pdftoppm",
            "-png",
            "-r",
            "144",
            "-scale-to-x",
            str(ANCHO),
            "-scale-to-y",
            str(ALTO),
            "-aa",
            "yes",
            "-aaVector",
            "yes",
            str(pdf),
            str(dir_png / "diapo"),
        ]
    )
    hechas = sorted(dir_png.glob("diapo-*.png")) or sorted(dir_png.glob("diapo_*.png"))
    for p in list(hechas):
        num = int(re.search(r"(\d+)", p.stem.split("-")[-1]).group(1))
        p.rename(dir_png / f"diapo_{num:03d}.png")
    hechas = sorted(dir_png.glob("diapo_*.png"))
    if len(hechas) != esperadas:
        raise RuntimeError(
            f"tema {tema}: se renderizaron {len(hechas)} imagenes y el deck tiene {esperadas} diapositivas"
        )
    return hechas


# ------------------------------------------------------------------------ 3. audio


def muestras_exactas(mp3: Path) -> tuple[int, int]:
    """(muestras decodificadas, frecuencia). Se decodifica de verdad: la
    duracion de la cabecera mp3 se desvia unos milisegundos por archivo y ese
    error se acumularia hasta descuadrar la imagen del audio al final del tema.
    """
    hz = int(
        correr(
            [
                "ffprobe",
                "-v",
                "error",
                "-select_streams",
                "a:0",
                "-show_entries",
                "stream=sample_rate",
                "-of",
                "csv=p=0",
                str(mp3),
            ]
        ).stdout.strip()
    )
    crudo = subprocess.run(
        ["ffmpeg", "-v", "error", "-i", str(mp3), "-f", "s16le", "-ac", "1", "-ar", str(hz), "-"],
        capture_output=True,
        check=True,
    ).stdout
    return len(crudo) // 2, hz


def paso_audio(tema: int, pares: list[dict], pausa: float, kbps: int, rehacer: bool) -> tuple[Path, list[float]]:
    destino = TRABAJO / "audio" / f"tema_{tema:02d}.m4a"
    destino.parent.mkdir(parents=True, exist_ok=True)

    duraciones = []
    entradas = []
    for par in pares:
        mp3 = AUDIOS / par["mp3"]
        muestras, hz = muestras_exactas(mp3)
        duraciones.append(muestras / hz + pausa)
        entradas.append(mp3)

    if destino.exists() and not rehacer:
        return destino, duraciones

    cmd = ["ffmpeg", "-y", "-hide_banner", "-loglevel", "error"]
    for mp3 in entradas:
        cmd += ["-i", str(mp3)]
    partes = [f"[{k}:a]apad=pad_dur={pausa}[a{k}];" for k in range(len(entradas))]
    cadena = "".join(partes) + "".join(f"[a{k}]" for k in range(len(entradas)))
    cadena += f"concat=n={len(entradas)}:v=0:a=1[out]"
    cmd += [
        "-filter_complex",
        cadena,
        "-map",
        "[out]",
        "-c:a",
        "aac",
        "-b:a",
        f"{kbps}k",
        "-ar",
        "44100",
        "-ac",
        "1",
        str(destino),
    ]
    correr(cmd)
    return destino, duraciones


# ------------------------------------------------------------------------ 4. video


def duracion(path: Path) -> float:
    r = correr(
        ["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "csv=p=0", str(path)]
    )
    return float(r.stdout.strip())


def paso_video(
    tema: int,
    imagenes: list[Path],
    duraciones: list[float],
    audio: Path,
    crf: int,
    fps: int,
    preset: str,
    rehacer: bool,
) -> Path:
    destino = SALIDA / f"tema_{tema:02d}.mp4"
    destino.parent.mkdir(parents=True, exist_ok=True)
    if destino.exists() and not rehacer:
        return destino

    lista = TRABAJO / "listas" / f"tema_{tema:02d}_imagenes.txt"
    lista.parent.mkdir(parents=True, exist_ok=True)
    lineas = []
    for img, dur in zip(imagenes, duraciones):
        lineas.append(f"file '{img}'")
        lineas.append(f"duration {dur:.6f}")
    # el concat demuxer ignora la duracion de la ultima entrada: se repite el
    # archivo y se deja cola de sobra; -shortest corta al terminar el audio.
    lineas.append(f"file '{imagenes[-1]}'")
    lineas.append("duration 3.0")
    lineas.append(f"file '{imagenes[-1]}'")
    lista.write_text("\n".join(lineas) + "\n", encoding="utf-8")

    correr(
        [
            "ffmpeg",
            "-y",
            "-hide_banner",
            "-loglevel",
            "error",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            str(lista),
            "-i",
            str(audio),
            "-map",
            "0:v:0",
            "-map",
            "1:a:0",
            "-c:v",
            "libx264",
            "-preset",
            preset,
            "-crf",
            str(crf),
            "-tune",
            "stillimage",
            "-pix_fmt",
            "yuv420p",
            "-r",
            str(fps),
            # un fotograma clave por minuto: el cambio de diapositiva ya
            # fuerza uno, y bajar de aqui duplicaba el peso del archivo sobre
            # contenido que no se mueve.
            "-g",
            str(fps * 60),
            "-vf",
            f"scale={ANCHO}:{ALTO}:force_original_aspect_ratio=decrease,"
            f"pad={ANCHO}:{ALTO}:(ow-iw)/2:(oh-ih)/2:white,setsar=1",
            "-c:a",
            "copy",
            "-movflags",
            "+faststart",
            "-shortest",
            str(destino),
        ]
    )
    return destino


# -------------------------------------------------------------------- 5. verificar


def paso_verificar(tema: int, imagenes: list[Path], duraciones: list[float], mp4: Path) -> dict:
    """Comprueba que el video dura lo que suma la narracion y que en mitad de
    cada diapositiva se ve esa diapositiva y no la vecina."""
    import numpy as np
    from PIL import Image

    def miniatura(origen) -> "np.ndarray":
        im = Image.open(origen).convert("L").resize((160, 90))
        return np.asarray(im, dtype=float)

    esperado = sum(duraciones)
    real = duracion(mp4)
    tmp = TRABAJO / "verificacion"
    tmp.mkdir(parents=True, exist_ok=True)

    referencias = [miniatura(p) for p in imagenes]
    inicios, acumulado = [], 0.0
    for d in duraciones:
        inicios.append(acumulado)
        acumulado += d

    muestras = sorted({0, len(imagenes) // 3, len(imagenes) // 2, 2 * len(imagenes) // 3, len(imagenes) - 1})
    fallos = []
    for i in muestras:
        t = inicios[i] + max(0.4, min(1.5, duraciones[i] / 2))
        foto = tmp / f"t{tema:02d}_d{i + 1:03d}.png"
        correr(
            ["ffmpeg", "-y", "-v", "error", "-ss", f"{t:.3f}", "-i", str(mp4),
             "-frames:v", "1", str(foto)]
        )
        actual = miniatura(foto)
        distancias = [float(np.abs(actual - r).mean()) for r in referencias]
        mejor = int(np.argmin(distancias))
        if mejor != i and distancias[i] - distancias[mejor] > 0.5:
            fallos.append({"diapositiva": i + 1, "segundo": round(t, 2), "se_ve": mejor + 1})
        foto.unlink(missing_ok=True)

    return {
        "duracion_esperada_s": round(esperado, 2),
        "duracion_real_s": round(real, 2),
        "desfase_s": round(real - esperado, 2),
        "diapositivas_comprobadas": [i + 1 for i in muestras],
        "fallos": fallos,
    }


# -------------------------------------------------------------------------- main


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--temas", nargs="*", type=int, default=TEMAS, help="temas a procesar (por defecto todos)")
    ap.add_argument(
        "--paso",
        choices=["mapear", "render", "audio", "video", "todo"],
        default="todo",
    )
    ap.add_argument("--pausa", type=float, default=0.6, help="silencio entre diapositivas, en segundos")
    ap.add_argument("--crf", type=int, default=21, help="calidad H.264: menor = mejor y mas pesado")
    ap.add_argument("--audio-kbps", type=int, default=96, help="bitrate AAC (el origen es mp3 de 48k)")
    ap.add_argument("--sin-verificar", action="store_true", help="no comprobar la sincronia al final")
    ap.add_argument("--fps", type=int, default=25)
    ap.add_argument("--preset", default="slow")
    ap.add_argument("--rehacer", action="store_true", help="no reutilizar productos previos")
    args = ap.parse_args()

    SALIDA.mkdir(parents=True, exist_ok=True)
    TRABAJO.mkdir(parents=True, exist_ok=True)

    log("1) emparejando diapositivas con audios")
    mapeo = paso_mapear(args.temas, args.rehacer and args.paso in ("mapear", "todo"))
    if args.paso == "mapear":
        log(f"\nmapeo en {SALIDA / 'mapeo.json'}")
        return 0

    informe = []
    for tema in args.temas:
        t0 = time.time()
        pares = mapeo[str(tema)]["pares"]
        log(f"\n2) tema {tema:2d}: renderizando {len(pares)} diapositivas")
        imagenes = paso_render(tema, args.rehacer)
        if args.paso == "render":
            continue

        log(f"3) tema {tema:2d}: montando la pista de audio")
        audio, duraciones = paso_audio(tema, pares, args.pausa, args.audio_kbps, args.rehacer)
        if args.paso == "audio":
            continue

        log(f"4) tema {tema:2d}: codificando video")
        mp4 = paso_video(
            tema, imagenes, duraciones, audio, args.crf, args.fps, args.preset, args.rehacer
        )
        control = None
        if not args.sin_verificar:
            control = paso_verificar(tema, imagenes, duraciones, mp4)
            if control["fallos"] or abs(control["desfase_s"]) > 0.5:
                log(f"   AVISO tema {tema}: {control}")
        seg = duracion(mp4)
        mb = mp4.stat().st_size / 1e6
        informe.append(
            {
                "tema": tema,
                "archivo": mp4.name,
                "diapositivas": len(imagenes),
                "duracion_s": round(seg, 1),
                "duracion": f"{int(seg // 60)}:{int(seg % 60):02d}",
                "mb": round(mb, 1),
                "mbits_por_min": round(mb * 8 / (seg / 60), 2),
                "verificacion": control,
            }
        )
        log(
            f"   -> {mp4.name}  {informe[-1]['duracion']}  {mb:.1f} MB"
            f"  ({time.time() - t0:.0f}s)"
        )

    if informe:
        (SALIDA / "informe.json").write_text(
            json.dumps(informe, ensure_ascii=False, indent=1), encoding="utf-8"
        )
        total_mb = sum(i["mb"] for i in informe)
        total_min = sum(i["duracion_s"] for i in informe) / 60
        log(f"\n{len(informe)} videos - {total_min:.0f} min - {total_mb:.0f} MB en {SALIDA}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
