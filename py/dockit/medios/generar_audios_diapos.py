#!/usr/bin/env python3
"""Extrae el contenido de cada diapositiva y genera un MP3 por diapositiva.

La fuente se mantiene intacta.  Los resultados se escriben en
``audios_diapos_original`` con una transcripcion TXT junto a cada audio.
"""

from __future__ import annotations

import argparse
import asyncio
import csv
import hashlib
import json
import re
import shutil
import subprocess
import sys
import tempfile
from collections import Counter
from dataclasses import asdict, dataclass
from pathlib import Path
from statistics import median
from typing import Iterable

import edge_tts
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


EMU_PER_INCH = 914_400
VOICE = "es-BO-MarceloNeural"
RATE = "-5%"
LOGO_TEXT = (
    "Texto visible en el distintivo institucional: 200. Policía Boliviana. "
    "1826–2026. 200 años de historia, vocación y servicio a la patria."
)
CREST_DESCRIPTION = (
    "Descripción de la imagen: escudo institucional de la Policía Boliviana, "
    "con cóndor, escudo central, rifles cruzados y ramas."
)
ARROWS_DESCRIPTION = (
    "Descripción de la imagen: dos flechas apuntan hacia abajo e indican la "
    "secuencia entre los enunciados."
)
SPECIAL_IMAGE_DESCRIPTION = (
    "Descripción de la imagen: un oficial de la Policía Boliviana señala un "
    "mapa digital de zonas de patrullaje; al fondo aparecen otros dos policías, "
    "vehículos policiales y la bandera de Bolivia. Texto visible dentro de la "
    "imagen: «PLANIFRULLAJE - Policía Boliviana». «Consiste en diseñar y "
    "estructurar el despliegue del personal policial en un territorio "
    "determinado y durante un tiempo específico. Su fin principal es disuadir "
    "la comisión de delitos, garantizar la seguridad ciudadana y mantener el "
    "orden público». «Safety Shield». «Shift Timelines». Números 1, 2, 3 y 4. "
    "«Tactical Districo». La placa del uniforme y el quinto marcador contienen "
    "caracteres que no son legibles con certeza."
)


@dataclass
class TextBlock:
    text: str
    source_text: str
    x: float
    y: float
    w: float
    h: float
    bold: bool
    max_font: float
    shape_name: str
    z_index: int
    auto_number: int | None = None

    @property
    def right(self) -> float:
        return self.x + self.w

    @property
    def bottom(self) -> float:
        return self.y + self.h

    @property
    def cx(self) -> float:
        return self.x + self.w / 2

    @property
    def cy(self) -> float:
        return self.y + self.h / 2

    @property
    def word_count(self) -> int:
        return len(re.findall(r"\b[\wÁÉÍÓÚÜÑáéíóúüñ]+\b", self.text))


@dataclass
class SlideRecord:
    deck_number: int
    deck_file: str
    slide_number: int
    transcript_relpath: str
    audio_relpath: str
    transcript_sha256: str
    characters: int
    words: int
    engine: str = "pending"
    duration_seconds: float | None = None


def clean_text(text: str) -> str:
    text = text.replace("\v", "\n").replace("\xa0", " ")
    lines = []
    for line in text.splitlines():
        line = re.sub(r"[ \t]+", " ", line).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def comparable_text(text: str) -> str:
    return "".join(character for character in text.casefold() if character.isalnum())


def paragraph_prefix(paragraph, running_number: int | None) -> tuple[str, int | None, int | None]:
    """Return visible bullet/number prefix, next counter and explicit number."""
    ppr = paragraph._p.find(qn("a:pPr"))
    if ppr is None:
        return "", None, None

    auto = ppr.find(qn("a:buAutoNum"))
    if auto is not None:
        start = auto.get("startAt")
        number = int(start) if start else (1 if running_number is None else running_number + 1)
        num_type = auto.get("type", "arabicPeriod")
        if num_type.startswith("arabic"):
            return f"{number}. ", number, number
        return f"Elemento {number}. ", number, number

    if ppr.find(qn("a:buChar")) is not None:
        return "• ", None, None
    return "", None, None


def extract_shape_text(shape) -> tuple[str, int | None]:
    if not getattr(shape, "has_text_frame", False):
        return "", None

    lines: list[str] = []
    running_number: int | None = None
    first_auto_number: int | None = None
    for paragraph in shape.text_frame.paragraphs:
        value = clean_text(paragraph.text)
        prefix, running_number, explicit_number = paragraph_prefix(paragraph, running_number)
        if explicit_number is not None and first_auto_number is None:
            first_auto_number = explicit_number
        if value:
            lines.append(prefix + value)
    return "\n".join(lines), first_auto_number


def source_shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(
        value
        for value in (clean_text(paragraph.text) for paragraph in shape.text_frame.paragraphs)
        if value
    )


def font_metadata(shape) -> tuple[bool, float]:
    bold_values: list[bool] = []
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return False, 0.0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.bold is not None:
                bold_values.append(bool(run.font.bold))
            if run.font.size is not None:
                sizes.append(float(run.font.size.pt))
    return any(bold_values), max(sizes, default=0.0)


def get_text_blocks(slide) -> list[TextBlock]:
    blocks: list[TextBlock] = []
    for z_index, shape in enumerate(slide.shapes):
        text, auto_number = extract_shape_text(shape)
        if not text:
            continue
        bold, max_font = font_metadata(shape)
        blocks.append(
            TextBlock(
                text=text,
                source_text=source_shape_text(shape),
                x=shape.left / EMU_PER_INCH,
                y=shape.top / EMU_PER_INCH,
                w=shape.width / EMU_PER_INCH,
                h=shape.height / EMU_PER_INCH,
                bold=bold,
                max_font=max_font,
                shape_name=shape.name,
                z_index=z_index,
                auto_number=auto_number,
            )
        )
    return blocks


def row_order(blocks: Iterable[TextBlock], tolerance: float = 0.38) -> list[TextBlock]:
    remaining = sorted(blocks, key=lambda b: (b.y, b.x, b.z_index))
    ordered: list[TextBlock] = []
    while remaining:
        anchor = remaining[0]
        row = [b for b in remaining if abs(b.y - anchor.y) <= tolerance]
        row.sort(key=lambda b: (b.x, b.y, b.z_index))
        ordered.extend(row)
        selected = {id(b) for b in row}
        remaining = [b for b in remaining if id(b) not in selected]
    return ordered


def is_top_block(block: TextBlock) -> bool:
    # Wide title/subtitle/intro blocks at the top of the slide.
    if block.y < 1.78 and block.w >= 7.0:
        return True
    if block.y < 1.25 and block.w >= 5.5:
        return True
    return False


def is_heading(block: TextBlock) -> bool:
    if block.word_count > 14:
        return False
    if block.h > 1.35:
        return False
    if block.bold:
        return True
    if block.max_font >= 27:
        return True
    if block.text.rstrip().endswith(":") and block.word_count <= 10:
        return True
    return False


def horizontal_overlap(a: TextBlock, b: TextBlock) -> float:
    overlap = max(0.0, min(a.right, b.right) - max(a.x, b.x))
    return overlap / max(0.01, min(a.w, b.w))


def body_heading_score(body: TextBlock, heading: TextBlock) -> float | None:
    # A description immediately below a heading in the same column.
    vgap = body.y - heading.bottom
    overlap = horizontal_overlap(body, heading)
    below_score: float | None = None
    if -0.28 <= vgap <= 3.2 and overlap >= 0.22:
        below_score = max(0.0, vgap) + abs(body.cx - heading.cx) * 0.08

    # A description to the right of a label, frequently joined by an arrow.
    hgap = body.x - heading.right
    right_score: float | None = None
    if -0.30 <= hgap <= 2.6 and abs(body.cy - heading.cy) <= 1.25:
        right_score = max(0.0, hgap) + abs(body.cy - heading.cy) * 0.12 + 0.04

    candidates = [s for s in (below_score, right_score) if s is not None]
    return min(candidates) if candidates else None


def order_blocks(blocks: list[TextBlock]) -> list[TextBlock]:
    if not blocks:
        return []

    top = row_order([b for b in blocks if is_top_block(b)], tolerance=0.28)
    top_ids = {id(b) for b in top}
    rest = [b for b in blocks if id(b) not in top_ids]

    # Wide centered prompts immediately below the title are introductions.
    intro = row_order(
        [
            b
            for b in rest
            if (b.w >= 7.0 and b.y < 2.75)
            or (
                b.w >= 5.0
                and b.y < 2.35
                and abs(b.cx - 6.665) <= 1.05
                and not is_heading(b)
            )
        ],
        tolerance=0.28,
    )
    intro_ids = {id(b) for b in intro}
    rest = [b for b in rest if id(b) not in intro_ids]

    headings = [b for b in rest if is_heading(b)]
    heading_ids = {id(b) for b in headings}
    bodies = [b for b in rest if id(b) not in heading_ids]

    assigned: dict[int, list[TextBlock]] = {id(h): [] for h in headings}
    leftovers: list[TextBlock] = []
    for body in bodies:
        scored = []
        for heading in headings:
            score = body_heading_score(body, heading)
            if score is not None:
                scored.append((score, abs(body.y - heading.y), heading.x, heading))
        if scored:
            scored.sort(key=lambda item: item[:3])
            assigned[id(scored[0][3])].append(body)
        else:
            leftovers.append(body)

    ordered: list[TextBlock] = [*top, *intro]

    # Numbered PowerPoint items have a definitive semantic order.
    numbered = [b for b in rest if b.auto_number is not None]
    if len(numbered) >= 2:
        numbered_ids = {id(b) for b in numbered}
        non_numbered_headings = [h for h in headings if id(h) not in numbered_ids]
        for heading in row_order(non_numbered_headings):
            ordered.append(heading)
            ordered.extend(row_order(assigned[id(heading)], tolerance=0.28))
        ordered.extend(sorted(numbered, key=lambda b: (b.auto_number or 0, b.y, b.x)))
        leftovers = [b for b in leftovers if id(b) not in numbered_ids]
    else:
        for heading in row_order(headings):
            ordered.append(heading)
            ordered.extend(row_order(assigned[id(heading)], tolerance=0.28))

    ordered.extend(row_order(leftovers))

    # Safety: preserve every block exactly once.
    seen: set[int] = set()
    unique: list[TextBlock] = []
    for block in ordered:
        if id(block) not in seen:
            seen.add(id(block))
            unique.append(block)
    for block in blocks:
        if id(block) not in seen:
            unique.append(block)
    return unique


def image_targets(slide) -> Counter[str]:
    targets: Counter[str] = Counter()
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        blip = shape._element.blipFill.blip
        rel_id = blip.rEmbed
        if not rel_id:
            continue
        rel = slide.part.rels[rel_id]
        targets[Path(str(rel.target_ref)).name] += 1
    return targets


def image_notes(slide) -> list[str]:
    targets = image_targets(slide)
    notes: list[str] = []
    if targets["image2.png"] or targets["image3.png"]:
        notes.append(CREST_DESCRIPTION)
    if targets["image5.png"]:
        notes.append(ARROWS_DESCRIPTION)
    if targets["image6.png"]:
        notes.append(SPECIAL_IMAGE_DESCRIPTION)
    if targets["image4.png"]:
        notes.append(LOGO_TEXT)
    return notes


def apply_audited_order(
    blocks: list[TextBlock], audited_blocks: list[str], deck_number: int, slide_number: int
) -> list[TextBlock]:
    remaining = list(blocks)
    ordered: list[TextBlock] = []
    for audited_text in audited_blocks:
        wanted = comparable_text(audited_text)
        matches = [
            (index, block)
            for index, block in enumerate(remaining)
            if comparable_text(block.source_text) == wanted
            or comparable_text(block.text) == wanted
        ]
        if not matches:
            raise RuntimeError(
                "El orden auditado no coincide con el PPTX en "
                f"presentación {deck_number}, diapositiva {slide_number}: {audited_text!r}"
            )
        index, block = matches[0]
        ordered.append(block)
        remaining.pop(index)
    if remaining:
        raise RuntimeError(
            "El orden auditado omitió bloques en "
            f"presentación {deck_number}, diapositiva {slide_number}: "
            f"{[block.source_text for block in remaining]!r}"
        )
    return ordered


def extract_narration(
    slide,
    deck_number: int,
    slide_number: int,
    audited_blocks: list[str] | None = None,
) -> tuple[str, list[TextBlock]]:
    blocks = get_text_blocks(slide)
    ordered = (
        apply_audited_order(blocks, audited_blocks, deck_number, slide_number)
        if audited_blocks is not None
        else order_blocks(blocks)
    )
    parts = [b.text for b in ordered]
    parts.extend(image_notes(slide))
    narration = "\n\n".join(part.strip() for part in parts if part.strip()).strip() + "\n"
    return narration, ordered


def xml_text_counter(slide) -> Counter[str]:
    values = slide._element.xpath(".//a:t/text()")
    return Counter(
        character
        for character in "".join(values).casefold()
        if character.isalnum()
    )


def block_text_counter(blocks: list[TextBlock]) -> Counter[str]:
    return Counter(
        character
        for block in blocks
        for character in block.source_text.casefold()
        if character.isalnum()
    )


def natural_deck_number(path: Path) -> int:
    match = re.match(r"\s*(\d+)", path.name)
    if not match:
        raise ValueError(f"No se pudo determinar el número de {path.name}")
    return int(match.group(1))


def load_reading_orders(path: Path | None) -> dict[tuple[int, int], list[str]]:
    if path is None:
        return {}
    data = json.loads(path.read_text(encoding="utf-8"))
    result: dict[tuple[int, int], list[str]] = {}
    for deck, slides in data["decks"].items():
        for slide, blocks in slides.items():
            result[(int(deck), int(slide))] = list(blocks)
    return result


def sha256_text(value: str) -> str:
    return hashlib.sha256(value.encode("utf-8")).hexdigest()


def prepare_transcripts(
    source_dir: Path,
    output_dir: Path,
    reading_orders: dict[tuple[int, int], list[str]],
) -> tuple[list[SlideRecord], list[dict]]:
    output_dir.mkdir(parents=True, exist_ok=True)
    records: list[SlideRecord] = []
    audits: list[dict] = []
    deck_paths = sorted(source_dir.glob("*.pptx"), key=natural_deck_number)
    if not deck_paths:
        raise FileNotFoundError(f"No se encontraron PPTX en {source_dir}")

    for deck_path in deck_paths:
        deck_number = natural_deck_number(deck_path)
        deck_dir = output_dir / f"presentacion_{deck_number:02d}"
        deck_dir.mkdir(parents=True, exist_ok=True)
        prs = Presentation(deck_path)
        combined: list[str] = []
        for slide_number, slide in enumerate(prs.slides, start=1):
            audited = reading_orders.get((deck_number, slide_number))
            narration, ordered_blocks = extract_narration(
                slide,
                deck_number,
                slide_number,
                audited,
            )
            native_xml = xml_text_counter(slide)
            native_blocks = block_text_counter(ordered_blocks)
            missing = native_xml - native_blocks
            extra = native_blocks - native_xml
            audit = {
                "presentacion": deck_number,
                "archivo": deck_path.name,
                "diapositiva": slide_number,
                "bloques_texto": len(ordered_blocks),
                "imagenes": dict(image_targets(slide)),
                "palabras_xml_no_extraidas": dict(missing),
                "palabras_extra_no_xml": dict(extra),
                "orden_lectura": "auditado" if audited is not None else "heurístico",
            }
            audits.append(audit)

            stem = f"diapositiva_{slide_number:03d}"
            transcript_path = deck_dir / f"{stem}.txt"
            audio_path = deck_dir / f"{stem}.mp3"
            transcript_path.write_text(narration, encoding="utf-8")
            combined.append(
                f"PRESENTACIÓN {deck_number:02d} — DIAPOSITIVA {slide_number:03d}\n\n"
                f"{narration.rstrip()}\n"
            )
            records.append(
                SlideRecord(
                    deck_number=deck_number,
                    deck_file=deck_path.name,
                    slide_number=slide_number,
                    transcript_relpath=str(transcript_path.relative_to(output_dir)),
                    audio_relpath=str(audio_path.relative_to(output_dir)),
                    transcript_sha256=sha256_text(narration),
                    characters=len(narration),
                    words=len(re.findall(r"\b\w+\b", narration)),
                )
            )
        (deck_dir / "transcripcion_completa.txt").write_text(
            "\n\n".join(combined).rstrip() + "\n", encoding="utf-8"
        )
    return records, audits


async def neural_tts(text: str, output_path: Path) -> None:
    communicator = edge_tts.Communicate(text=text, voice=VOICE, rate=RATE)
    await communicator.save(str(output_path))


def local_tts(text: str, output_path: Path) -> None:
    with tempfile.TemporaryDirectory(prefix="pptx_audio_fallback_") as temp_dir:
        wav_path = Path(temp_dir) / "fallback.wav"
        subprocess.run(
            ["espeak-ng", "-v", "es-la", "-s", "150", "-w", str(wav_path), text],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.PIPE,
        )
        subprocess.run(
            [
                "ffmpeg",
                "-v",
                "error",
                "-y",
                "-i",
                str(wav_path),
                "-codec:a",
                "libmp3lame",
                "-b:a",
                "64k",
                str(output_path),
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.PIPE,
        )


def valid_mp3(path: Path) -> bool:
    if not path.exists() or path.stat().st_size < 1_000:
        return False
    probe = subprocess.run(
        [
            "ffprobe",
            "-v",
            "error",
            "-select_streams",
            "a:0",
            "-show_entries",
            "stream=codec_name",
            "-of",
            "default=noprint_wrappers=1:nokey=1",
            str(path),
        ],
        text=True,
        capture_output=True,
    )
    return probe.returncode == 0 and bool(probe.stdout.strip())


def duration_seconds(path: Path) -> float | None:
    probe = subprocess.run(
        [
            "ffprobe",
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "default=noprint_wrappers=1:nokey=1",
            str(path),
        ],
        text=True,
        capture_output=True,
    )
    try:
        return round(float(probe.stdout.strip()), 3) if probe.returncode == 0 else None
    except ValueError:
        return None


async def generate_one(record: SlideRecord, output_dir: Path, semaphore: asyncio.Semaphore) -> None:
    transcript_path = output_dir / record.transcript_relpath
    audio_path = output_dir / record.audio_relpath
    state_path = (
        output_dir
        / ".estado"
        / f"presentacion_{record.deck_number:02d}_diapositiva_{record.slide_number:03d}.sha256"
    )
    state_path.parent.mkdir(parents=True, exist_ok=True)
    text = transcript_path.read_text(encoding="utf-8").strip()
    saved_hash = state_path.read_text(encoding="utf-8").strip() if state_path.exists() else ""
    if valid_mp3(audio_path) and saved_hash == record.transcript_sha256:
        record.engine = "edge-tts (existente)"
        record.duration_seconds = duration_seconds(audio_path)
        return

    part_path = audio_path.with_suffix(".part.mp3")
    part_path.unlink(missing_ok=True)
    async with semaphore:
        last_error: Exception | None = None
        for attempt in range(1, 6):
            try:
                await neural_tts(text, part_path)
                if not valid_mp3(part_path):
                    raise RuntimeError("edge-tts produjo un MP3 inválido")
                part_path.replace(audio_path)
                state_path.write_text(record.transcript_sha256 + "\n", encoding="utf-8")
                record.engine = "edge-tts"
                record.duration_seconds = duration_seconds(audio_path)
                return
            except Exception as exc:  # noqa: BLE001 - se registra y reintenta
                last_error = exc
                part_path.unlink(missing_ok=True)
                await asyncio.sleep(min(2 ** (attempt - 1), 16))

        try:
            await asyncio.to_thread(local_tts, text, part_path)
            if not valid_mp3(part_path):
                raise RuntimeError("el motor local produjo un MP3 inválido")
            part_path.replace(audio_path)
            state_path.write_text(record.transcript_sha256 + "\n", encoding="utf-8")
            record.engine = "espeak-ng (respaldo local)"
            record.duration_seconds = duration_seconds(audio_path)
        except Exception as fallback_error:  # noqa: BLE001
            part_path.unlink(missing_ok=True)
            raise RuntimeError(
                f"Falló edge-tts ({last_error}) y el respaldo local ({fallback_error})"
            ) from fallback_error


async def generate_audios(records: list[SlideRecord], output_dir: Path, concurrency: int) -> None:
    semaphore = asyncio.Semaphore(concurrency)
    completed = 0
    lock = asyncio.Lock()

    async def wrapped(record: SlideRecord) -> None:
        nonlocal completed
        await generate_one(record, output_dir, semaphore)
        async with lock:
            completed += 1
            if completed % 10 == 0 or completed == len(records):
                print(f"Audio: {completed}/{len(records)}", flush=True)

    await asyncio.gather(*(wrapped(record) for record in records))


def write_reports(output_dir: Path, records: list[SlideRecord], audits: list[dict]) -> None:
    manifest_path = output_dir / "manifiesto.json"
    manifest_path.write_text(
        json.dumps(
            {
                "voz": VOICE,
                "velocidad": RATE,
                "total_presentaciones": len({r.deck_number for r in records}),
                "total_diapositivas": len(records),
                "registros": [asdict(r) for r in records],
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    (output_dir / "auditoria_extraccion.json").write_text(
        json.dumps(audits, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )

    with (output_dir / "indice.csv").open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(asdict(records[0]).keys()))
        writer.writeheader()
        writer.writerows(asdict(r) for r in records)

    counts = Counter(r.deck_number for r in records)
    total_duration = sum(r.duration_seconds or 0.0 for r in records)
    fallbacks = [r for r in records if r.engine.startswith("espeak")]
    readme = [
        "# Audios por diapositiva",
        "",
        f"- Fuente: presentaciones de `{Path('diapos_original')}`.",
        f"- Total: {len(records)} diapositivas en {len(counts)} presentaciones.",
        f"- Formato: MP3, un archivo por diapositiva.",
        f"- Voz principal: `{VOICE}`, velocidad `{RATE}`.",
        f"- Duración total: {total_duration / 60:.1f} minutos.",
        f"- Audios generados con respaldo local: {len(fallbacks)}.",
        "- Cada MP3 tiene a su lado un TXT con el texto exacto narrado.",
        "- Las transcripciones incorporan el texto rasterizado del distintivo institucional.",
        "- Las imágenes informativas se describen; el fondo meramente decorativo no se verbaliza.",
        "- Se conservaron literalmente la ortografía y los duplicados de las diapositivas fuente.",
        "",
        "## Cantidad por presentación",
        "",
    ]
    readme.extend(f"- Presentación {deck:02d}: {counts[deck]} audios." for deck in sorted(counts))
    (output_dir / "README.md").write_text("\n".join(readme) + "\n", encoding="utf-8")


def validate(records: list[SlideRecord], audits: list[dict], output_dir: Path) -> None:
    bad_audio = [r.audio_relpath for r in records if not valid_mp3(output_dir / r.audio_relpath)]
    bad_duration = [
        r.audio_relpath
        for r in records
        if (duration_seconds(output_dir / r.audio_relpath) or 0.0) <= 0.5
    ]
    empty_text = [r.transcript_relpath for r in records if r.characters <= 1]
    missing_native = [a for a in audits if a["palabras_xml_no_extraidas"]]
    if bad_audio:
        raise RuntimeError(f"MP3 inválidos o ausentes: {bad_audio[:10]}")
    if bad_duration:
        raise RuntimeError(f"MP3 sin duración útil: {bad_duration[:10]}")
    if empty_text:
        raise RuntimeError(f"Transcripciones vacías: {empty_text[:10]}")
    if missing_native:
        sample = missing_native[:5]
        raise RuntimeError(f"Texto XML no extraído en {len(missing_native)} diapositivas: {sample}")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--concurrency", type=int, default=4)
    parser.add_argument("--reading-order", type=Path)
    parser.add_argument("--transcripts-only", action="store_true")
    args = parser.parse_args()

    source_dir = args.source.resolve()
    output_dir = args.output.resolve()
    reading_orders = load_reading_orders(args.reading_order.resolve() if args.reading_order else None)
    records, audits = prepare_transcripts(source_dir, output_dir, reading_orders)
    print(f"Transcripciones: {len(records)}", flush=True)
    if not args.transcripts_only:
        asyncio.run(generate_audios(records, output_dir, max(1, args.concurrency)))
        for record in records:
            record.duration_seconds = duration_seconds(output_dir / record.audio_relpath)
        validate(records, audits, output_dir)
    write_reports(output_dir, records, audits)
    print(f"Resultado: {output_dir}", flush=True)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
