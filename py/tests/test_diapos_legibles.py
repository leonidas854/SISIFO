"""Las diapositivas se proyectan: el texto tiene que leerse desde el fondo.

Referencia habitual en salas: cuerpo a 24 pt o más, títulos a 36 pt o más.
Con 16 pt, la última fila no lee nada.
"""
import pytest
from pptx import Presentation
from pptx.util import Pt

from dockit.generadores import pptx_node

CUERPO_MIN, TITULO_MIN = 24, 36

GUION = {
    "tipo": "pptx", "titulo": "El problema del oráculo", "autor": "Autor",
    "bloques": [
        {"clase": "titulo", "nivel": 1, "texto": "Introducción"},
        {"clase": "parrafo", "texto": "El oráculo media el acceso a datos externos."},
        {"clase": "titulo", "nivel": 1, "texto": "Riesgos"},
        {"clase": "lista", "items": ["Manipulación del precio", "Centralización",
                                     "Falta de auditoría"]},
    ],
}

pytestmark = pytest.mark.skipif(not pptx_node.disponible(),
                                reason="requiere Node con pptxgenjs")


@pytest.fixture
def pres(tmp_path):
    d = tmp_path / "d.pptx"
    pptx_node.generar(GUION, str(d), {}, {})
    return Presentation(str(d))


def tamanos(pres):
    """(tamaño en puntos, texto) de cada corrida con texto."""
    out = []
    for s in pres.slides:
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip() and r.font.size:
                        out.append((r.font.size.pt, r.text.strip()))
    return out


# Lo que NO es cuerpo: numeración, créditos de portada y el pie que repite el
# título del trabajo en cada lámina. El pie va pequeño a propósito: es una
# referencia de ubicación, no algo que nadie deba leer desde el fondo.
CHROME = ("Autor",)


def es_cuerpo(texto: str, titulo_trabajo: str = "") -> bool:
    if texto.isdigit() or len(texto) <= 18:
        return False
    if titulo_trabajo and texto.strip() == titulo_trabajo:
        return False
    return not any(texto.startswith(c) for c in CHROME)


def test_el_cuerpo_se_lee_desde_el_fondo(pres):
    cuerpo = [(t, x) for t, x in tamanos(pres)
              if es_cuerpo(x, GUION["titulo"])]
    assert cuerpo, "no encuentro texto de cuerpo"
    pequenos = [(t, x[:40]) for t, x in cuerpo if t < CUERPO_MIN]
    assert not pequenos, f"texto por debajo de {CUERPO_MIN} pt: {pequenos[:3]}"


def test_los_titulos_son_titulos(pres):
    grandes = [t for t, x in tamanos(pres) if t >= TITULO_MIN]
    assert grandes, f"ningún texto llega a {TITULO_MIN} pt"


def test_una_lamina_por_titulo_del_indice(pres):
    titulos = []
    for s in pres.slides:
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                titulos.append(sh.text_frame.text.strip())
                break
    junto = " | ".join(titulos)
    assert "Introducción" in junto and "Riesgos" in junto, \
        "las láminas deben seguir el índice del trabajo"


def test_no_se_amontona_el_texto(pres):
    """Más de seis viñetas en una lámina es un muro de texto."""
    for s in pres.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                lineas = [p for p in sh.text_frame.paragraphs if p.text.strip()]
                assert len(lineas) <= 7, f"{len(lineas)} líneas en una sola caja"


def test_las_referencias_no_secuestran_la_presentacion(tmp_path):
    """Caso real: 79 referencias generaron 16 láminas seguidas. Proyectadas
    no las lee nadie; su sitio es el informe."""
    bib = {f"clave{i}": f"Autor{i}, A. ({2000+i}). Un título largo de artículo "
                        f"académico número {i}. Revista." for i in range(40)}
    guion = {"tipo": "pptx", "titulo": "T", "bloques": [
        {"clase": "titulo", "nivel": 1, "texto": "Intro"},
        {"clase": "lista", "items": ["Uno", "Dos"]},
        {"clase": "bibliografia"}]}
    d = tmp_path / "r.pptx"
    pptx_node.generar(guion, str(d), bib, {})
    pres = Presentation(str(d))
    titulos = []
    for s in pres.slides:
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                titulos.append(sh.text_frame.text.strip())
                break
    refs = [t for t in titulos if t.startswith("Referencias")]
    assert len(refs) <= 2, f"{len(refs)} láminas de referencias es demasiado"
    assert len(pres.slides) <= 6, f"{len(pres.slides)} láminas para 2 ideas"


def test_la_imagen_va_en_la_misma_lamina_que_las_vinetas(tmp_path):
    """«cada diapo debe tener imágenes»: el diagrama acompaña a las viñetas,
    no ocupa una lámina aparte que rompe el ritmo de la exposición."""
    from dockit.imagen import ilustrar
    png = ilustrar.ilustrar_lamina("Riesgos", ["Manipulación", "Centralización"],
                                   tmp_path, "fig")
    assert png, "no se generó el diagrama de prueba"
    guion = {"tipo": "pptx", "titulo": "T", "bloques": [
        {"clase": "titulo", "nivel": 1, "texto": "Riesgos"},
        {"clase": "lista", "items": ["Manipulación", "Centralización"]},
        {"clase": "figura", "ruta": str(png), "leyenda": "Riesgos"}]}
    d = tmp_path / "j.pptx"
    pptx_node.generar(guion, str(d), {}, {})
    pres = Presentation(str(d))

    laminas_con_imagen = [
        s for s in pres.slides
        if any(sh.shape_type == 13 or sh.__class__.__name__ == "Picture"
               for sh in s.shapes)]
    assert laminas_con_imagen, "ninguna lámina lleva imagen"

    s = laminas_con_imagen[0]
    textos = " ".join(sh.text_frame.text for sh in s.shapes
                      if sh.has_text_frame and sh.text_frame.text)
    assert "Manipulación" in textos, \
        "la imagen quedó en una lámina aparte, sin sus viñetas"
