"""Una presentación académica tiene partes que el público espera encontrar.

Sin portada con los datos del trabajo, sin agenda, sin separadores y sin
cierre, el deck es una lista de láminas sueltas. Y la marca de la herramienta
no pinta nada en la portada de un trabajo del usuario.
"""
import pytest
from pptx import Presentation

from dockit.generadores import pptx_node

pytestmark = pytest.mark.skipif(not pptx_node.disponible(),
                                reason="requiere Node con pptxgenjs")

PORTADA = {
    "institucion": "Escuela Militar de Ingeniería",
    "carrera": "Ingeniería de Sistemas",
    "materia": "Seguridad de la Información",
    "docente": "Ing. Juan Pérez",
    "autor": "leonidas854",
    "fecha": "septiembre de 2026",
    "lugar": "Cochabamba, Bolivia",
}

GUION = {
    "tipo": "pptx", "titulo": "El problema del oráculo", "autor": "leonidas854",
    "portada": PORTADA,
    "bloques": [
        {"clase": "titulo", "nivel": 1, "texto": "Introducción"},
        {"clase": "lista", "items": ["Idea uno", "Idea dos"]},
        {"clase": "titulo", "nivel": 1, "texto": "Conclusiones"},
        {"clase": "lista", "items": ["Cierre"]},
        {"clase": "bibliografia"},
    ],
}


@pytest.fixture
def pres(tmp_path):
    d = tmp_path / "f.pptx"
    pptx_node.generar(GUION, str(d), {"k": "Autor, A. (2020). Obra."}, {})
    return Presentation(str(d))


def texto_de(slide) -> str:
    return "\n".join(sh.text_frame.text for sh in slide.shapes
                     if sh.has_text_frame and sh.text_frame.text)


def todo(pres) -> str:
    return "\n".join(texto_de(s) for s in pres.slides)


def test_la_portada_lleva_los_datos_academicos(pres):
    # la institución va en versalitas por diseño, así que se compara sin caso
    portada = texto_de(pres.slides[0]).lower()
    for campo in ("institucion", "materia", "docente", "autor", "fecha"):
        assert PORTADA[campo].lower() in portada, f"la portada no muestra «{campo}»"


def test_la_marca_de_la_herramienta_no_aparece(pres):
    assert "SÍSIFO" not in todo(pres) and "SISIFO" not in todo(pres), \
        "la herramienta no debe firmar el trabajo del usuario"


def test_hay_agenda_al_principio(pres):
    """El público necesita saber la estructura antes de empezar."""
    primeras = "\n".join(texto_de(s) for s in list(pres.slides)[:3])
    assert "Contenido" in primeras or "Agenda" in primeras
    assert "Introducción" in primeras and "Conclusiones" in primeras


def test_hay_lamina_de_cierre(pres):
    ultima = texto_de(list(pres.slides)[-1])
    assert "Gracias" in ultima or "Preguntas" in ultima, \
        "el deck termina de golpe en las referencias"


def test_las_laminas_llevan_pie_con_el_titulo(pres):
    """En una defensa, quien entra a media exposición debe ubicarse."""
    intermedia = texto_de(list(pres.slides)[3])
    assert "El problema del oráculo" in intermedia


def test_la_portada_funciona_sin_datos_academicos(tmp_path):
    """No todos los trabajos declaran institución: no puede reventar."""
    guion = dict(GUION)
    guion.pop("portada")
    d = tmp_path / "g.pptx"
    pptx_node.generar(guion, str(d), {}, {})
    assert Presentation(str(d)).slides


def test_las_etiquetas_de_la_portada_se_ven(pres):
    """«PRESENTADO POR» iba en gris apagado sobre fondo oscuro: ilegible."""
    portada = pres.slides[0]
    for sh in portada.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip().upper() in ("PRESENTADO POR", "DOCENTE"):
                    color = r.font.color
                    assert color and color.rgb, "la etiqueta no declara color"
                    # sobre fondo oscuro hace falta un tono claro
                    rgb = str(color.rgb)
                    claro = sum(int(rgb[i:i+2], 16) for i in (0, 2, 4)) / 3
                    assert claro > 110, \
                        f"«{r.text}» va a {rgb}, demasiado oscuro para el fondo"
