"""Tests del guion y los generadores. Se escriben antes que el código.

El guion es la estructura que describe un documento sin comprometerse con el
formato: el mismo guion produce el .docx, el .pptx y el .xlsx.
"""
import json
from pathlib import Path

import pytest

from dockit.generadores import guion as g


# ── validación del guion ─────────────────────────────────────────────────

def test_bloque_desconocido_falla():
    with pytest.raises(g.GuionInvalido, match="desconocida"):
        g.validar({"tipo": "docx", "titulo": "T",
                   "bloques": [{"clase": "chirimoya"}]})


def test_titulo_sin_texto_falla():
    with pytest.raises(g.GuionInvalido):
        g.validar({"tipo": "docx", "titulo": "T",
                   "bloques": [{"clase": "titulo", "nivel": 1}]})


def test_tabla_descuadrada_falla():
    with pytest.raises(g.GuionInvalido, match="celdas"):
        g.validar({"tipo": "docx", "titulo": "T", "bloques": [
            {"clase": "tabla", "cabecera": ["a", "b"], "filas": [["1"]]}]})


def test_cita_fuera_de_la_bibliografia_falla():
    """La barrera contra la referencia inventada, también al producir."""
    with pytest.raises(g.GuionInvalido, match="fantasma"):
        g.validar({"tipo": "docx", "titulo": "T", "bloques": [
            {"clase": "parrafo", "texto": "x", "citas": ["fantasma"]}]},
            disponibles={"real"})


def test_guion_valido_pasa():
    g.validar({"tipo": "docx", "titulo": "T", "bloques": [
        {"clase": "titulo", "nivel": 1, "texto": "Intro"},
        {"clase": "parrafo", "texto": "Texto", "citas": ["k"]},
        {"clase": "bibliografia"}]}, disponibles={"k"})


# ── generación real ──────────────────────────────────────────────────────

@pytest.fixture
def guion_completo():
    return {
        "tipo": "docx", "titulo": "Informe de prueba", "autor": "Autor",
        "bloques": [
            {"clase": "titulo", "nivel": 1, "texto": "Introducción"},
            {"clase": "parrafo", "texto": "Primer párrafo del informe.",
             "citas": ["nath2024"]},
            {"clase": "lista", "items": ["uno", "dos", "tres"]},
            {"clase": "tabla", "cabecera": ["Concepto", "Valor"],
             "filas": [["alfa", "1"], ["beta", "2"]],
             "leyenda": "Tabla 1. Datos de prueba"},
            {"clase": "cita", "texto": "Una cita en bloque."},
            {"clase": "titulo", "nivel": 2, "texto": "Conclusión"},
            {"clase": "parrafo", "texto": "Cierre."},
            {"clase": "bibliografia"},
        ],
    }


BIB = {"nath2024": "Nath, S. G. (2024). Digital Evidence Chain of Custody. IEEE."}
EN_TEXTO = {"nath2024": "(Nath, 2024)"}


def test_docx_se_genera_y_abre(tmp_path, guion_completo):
    from dockit.generadores import docx as gd
    destino = tmp_path / "informe.docx"
    doc = gd.generar(guion_completo, str(destino), BIB, EN_TEXTO)
    assert destino.exists() and destino.stat().st_size > 0
    from docx import Document
    d = Document(str(destino))
    texto = "\n".join(p.text for p in d.paragraphs)
    assert "Introducción" in texto
    assert "Primer párrafo" in texto
    assert "(Nath, 2024)" in texto, "la cita en el texto debe insertarse"
    assert "Nath, S. G. (2024)" in texto, "la bibliografía debe aparecer"
    assert len(d.tables) == 1
    assert doc["unidades"] >= 1


def test_docx_inserta_la_cita_donde_toca(tmp_path):
    from dockit.generadores import docx as gd
    guion = {"tipo": "docx", "titulo": "T", "bloques": [
        {"clase": "parrafo", "texto": "La cadena documenta cada traspaso.",
         "citas": ["nath2024"]}]}
    destino = tmp_path / "x.docx"
    gd.generar(guion, str(destino), BIB, EN_TEXTO)
    from docx import Document
    p = Document(str(destino)).paragraphs
    cuerpo = [x.text for x in p if "cadena documenta" in x.text]
    assert cuerpo, "no encuentro el párrafo"
    assert cuerpo[0].rstrip().endswith("(Nath, 2024)."), \
        f"la cita debe cerrar la frase, salió: {cuerpo[0]!r}"


def test_pptx_una_diapositiva_por_titulo(tmp_path, guion_completo):
    from dockit.generadores import pptx as gp
    guion_completo["tipo"] = "pptx"
    destino = tmp_path / "diapos.pptx"
    doc = gp.generar(guion_completo, str(destino), BIB, EN_TEXTO)
    from pptx import Presentation
    pres = Presentation(str(destino))
    # portada + una por cada título de nivel 1-2 + bibliografía
    assert len(pres.slides) >= 3
    assert doc["unidades"] == len(pres.slides)
    titulos = []
    for s in pres.slides:
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text:
                titulos.append(sh.text_frame.text)
    junto = "\n".join(titulos)
    assert "Introducción" in junto and "Conclusión" in junto


def test_xlsx_una_hoja_por_tabla(tmp_path, guion_completo):
    from dockit.generadores import xlsx as gx
    guion_completo["tipo"] = "xlsx"
    destino = tmp_path / "datos.xlsx"
    doc = gx.generar(guion_completo, str(destino), BIB, EN_TEXTO)
    import openpyxl
    wb = openpyxl.load_workbook(str(destino))
    assert "Tabla 1" in wb.sheetnames or len(wb.sheetnames) >= 1
    hoja = wb[wb.sheetnames[0]]
    assert hoja.cell(1, 1).value == "Concepto"
    assert hoja.cell(2, 1).value == "alfa"
    assert doc["unidades"] >= 1


def test_tipo_no_soportado_falla(tmp_path, guion_completo):
    from dockit.generadores import generar_desde_guion
    guion_completo["tipo"] = "jeroglifico"
    with pytest.raises(g.GuionInvalido):
        generar_desde_guion(guion_completo, str(tmp_path / "x"), BIB, EN_TEXTO)


# ── regresiones halladas produciendo un trabajo real ─────────────────────

def test_no_se_cita_una_clave_sin_forma_en_texto(tmp_path):
    """Bug real: una referencia no verificada quedaba fuera de citas_en_texto,
    y el generador escribía «(cosic2011ontological)» en el documento en vez de
    negarse. Citar algo no verificado es justo lo que el sistema debe impedir."""
    from dockit.generadores import generar_desde_guion
    guion = {"tipo": "docx", "titulo": "T", "bloques": [
        {"clase": "parrafo", "texto": "Algo.", "citas": ["sin_verificar"]}]}
    with pytest.raises(g.GuionInvalido, match="sin_verificar"):
        generar_desde_guion(guion, str(tmp_path / "x.docx"),
                            bibliografia={"sin_verificar": "Entrada suelta"},
                            en_texto={})   # no tiene forma en texto


def test_bibliografia_se_escribe_completa(tmp_path):
    """Bug real: la sección Referencias salía vacía porque el emparejamiento
    clave -> entrada APA se hacía adivinando por el apellido."""
    from dockit.generadores import docx as gd
    guion = {"tipo": "docx", "titulo": "T", "bloques": [
        {"clase": "parrafo", "texto": "Algo.", "citas": ["nath2024"]},
        {"clase": "bibliografia"}]}
    destino = tmp_path / "x.docx"
    gd.generar(guion, str(destino), BIB, EN_TEXTO)
    from docx import Document
    texto = "\n".join(p.text for p in Document(str(destino)).paragraphs)
    assert "Nath, S. G. (2024). Digital Evidence Chain of Custody. IEEE." in texto
