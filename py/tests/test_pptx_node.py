"""El generador PPTX con gráficos nativos (pptxgenjs sobre Node).

python-pptx no sabe hacer gráficos nativos ni notas del orador con soltura;
pptxgenjs sí. Se elige uno u otro según lo que haya instalado, pero el
contrato es el mismo: mismo guion dentro, mismo tipo de resultado fuera.
"""
import shutil

import pytest

from dockit.generadores import pptx_node

BIB = {"nath2024": "Nath, S. G. (2024). Digital Evidence Chain of Custody. IEEE."}
EN_TEXTO = {"nath2024": "(Nath, 2024)"}

GUION = {
    "tipo": "pptx", "titulo": "Cadena de custodia", "autor": "Autor",
    "bloques": [
        {"clase": "titulo", "nivel": 1, "texto": "Introducción"},
        {"clase": "parrafo", "texto": "La cadena documenta cada traspaso.",
         "citas": ["nath2024"]},
        {"clase": "tabla", "cabecera": ["Enfoque", "Valor"],
         "filas": [["Papel", "2"], ["Cadena de bloques", "5"]],
         "leyenda": "Tabla 1. Comparación"},
        {"clase": "bibliografia"},
    ],
}

hay_node = pptx_node.disponible()


def test_detecta_si_puede_usarse():
    """disponible() no debe reventar aunque falte Node: solo decir que no."""
    assert isinstance(pptx_node.disponible(), bool)


@pytest.mark.skipif(not hay_node, reason="requiere Node con pptxgenjs")
def test_genera_pptx_con_texto_nativo(tmp_path):
    destino = tmp_path / "d.pptx"
    r = pptx_node.generar(GUION, str(destino), BIB, EN_TEXTO)
    assert destino.exists() and destino.stat().st_size > 0
    assert r["unidades"] >= 3

    from pptx import Presentation
    pres = Presentation(str(destino))
    texto = "\n".join(
        sh.text_frame.text
        for s in pres.slides for sh in s.shapes
        if sh.has_text_frame and sh.text_frame.text)
    # el texto tiene que ser editable en PowerPoint, no estar dentro de una imagen
    assert "Introducción" in texto
    assert "(Nath, 2024)" in texto, "la cita debe ir como texto, no incrustada"


@pytest.mark.skipif(not hay_node, reason="requiere Node con pptxgenjs")
def test_la_tabla_llega_como_tabla_o_grafico(tmp_path):
    destino = tmp_path / "t.pptx"
    pptx_node.generar(GUION, str(destino), BIB, EN_TEXTO)
    from pptx import Presentation
    pres = Presentation(str(destino))
    tiene = any(sh.has_table or sh.has_chart
                for s in pres.slides for sh in s.shapes)
    assert tiene, "la tabla debe seguir siendo un objeto nativo, no una foto"


def test_error_de_node_se_reporta_claro(tmp_path):
    """Si el guion es inválido, el fallo tiene que decir qué pasó."""
    with pytest.raises(Exception) as e:
        pptx_node.generar({"tipo": "pptx"}, str(tmp_path / "x.pptx"), {}, {})
    assert str(e.value)
