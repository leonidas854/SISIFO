"""Regresiones del contrato que evita imágenes vacías o fuera de contexto."""

from __future__ import annotations

import json
import xml.etree.ElementTree as ET

import pytest

from dockit.visual.dominio import (
    PlanVisual, Procedencia, Visual, cargar_plan, construir_plan_desde_guion,
    validar_plan,
)
from dockit.visual.pptx_adapter import auditar_presentacion
from dockit.visual.semantica import SemanticaLexica
from dockit.visual.vector import generar_svg


def visual_valida(**cambios) -> Visual:
    datos = dict(
        diapositiva=2,
        opcion=1,
        titulo="Cadena de custodia digital",
        proposito="Explicar cómo el registro continuo protege la integridad de la evidencia.",
        tipo="proceso",
        motor="vector",
        concepto_visual="Cuatro etapas enlazadas desde la recolección hasta la entrega.",
        conceptos=("recolección", "registro", "traslado", "entrega"),
        texto_visible=("Recolección", "Registro", "Traslado", "Entrega"),
        texto_alternativo="Proceso de cuatro etapas que mantiene trazable la evidencia.",
        procedencia=Procedencia(referencia="nist2006", verificada=True),
    )
    datos.update(cambios)
    return Visual(**datos)


def test_plan_visual_completo_pasa():
    plan = PlanVisual(1, "T", (visual_valida(),))
    assert validar_plan(plan, {(2, 1): 0.91}) == []


def test_ley_sin_rotulo_ni_fuente_falla():
    ley = visual_valida(
        titulo="Ley N° 101 del régimen disciplinario",
        proposito="Identificar la norma que define las faltas disciplinarias policiales.",
        tipo="ley",
        conceptos=("faltas disciplinarias", "régimen policial"),
        texto_visible=(),
        procedencia=Procedencia(),
    )
    codigos = {h.codigo for h in validar_plan(PlanVisual(1, "T", (ley,)))}
    assert {"VIS-021", "VIS-022"} <= codigos


def test_ley_no_delega_texto_a_imagegen():
    ley = visual_valida(
        titulo="Ley N° 348",
        proposito="Mostrar el nombre oficial de la ley y su finalidad de protección.",
        tipo="ley", motor="imagegen",
        conceptos=("protección", "violencia"), texto_visible=("Ley N° 348",),
        prompt="A readable law book with title Ley 348 written on its cover",
    )
    codigos = {h.codigo for h in validar_plan(PlanVisual(1, "T", (ley,)))}
    assert "VIS-020" in codigos
    assert "VIS-031" in codigos


def test_baja_relacion_semantica_bloquea():
    v = visual_valida()
    hallazgos = validar_plan(PlanVisual(1, "T", (v,)), {(2, 1): 0.12})
    assert any(h.codigo == "VIS-050" and h.severidad == "error" for h in hallazgos)


def test_comparador_lexico_distingue_relacion():
    s = SemanticaLexica()
    relacionada, ajena = s.comparar([
        ("integridad y trazabilidad de la evidencia", "evidencia íntegra con registro trazable"),
        ("integridad y trazabilidad de la evidencia", "paisaje de montañas y comida tropical"),
    ])
    assert relacionada > 0.48
    assert ajena == 0


def test_carga_plan_legado_y_expone_sus_carencias(tmp_path):
    ruta = tmp_path / "legado.json"
    ruta.write_text(json.dumps({"titulo_tema": "Normas", "trabajos": [{
        "diapositiva": 3, "opcion": 1, "titulo": "Artículo 1: deber",
        "que_debe_leerse": "El funcionario cumple el deber que la ley le asigna.",
        "concepto_visual": "Un reloj junto a un escudo.",
        "motor": "vector", "spec": {"tipo": "fila", "iconos": ["reloj", "escudo"]},
    }]}), encoding="utf-8")
    plan = cargar_plan(ruta)
    assert plan.formato_legado and plan.visuales[0].tipo == "diagrama"
    codigos = {h.codigo for h in validar_plan(plan)}
    assert "VIS-002" in codigos
    assert "VIS-012" in codigos
    assert "VIS-021" in codigos, "un artículo sin texto visible reproduce el defecto real"


def test_borrador_sale_del_guion_sin_inventar_fuente():
    plan = construir_plan_desde_guion({
        "titulo": "Normativa",
        "bloques": [
            {"clase": "titulo", "nivel": 1, "texto": "Ley N° 101"},
            {"clase": "parrafo", "texto": "La norma regula el régimen disciplinario."},
        ],
    })
    visual = plan["visuales"][1]
    assert visual["tipo"] == "ley"
    assert visual["texto_visible"] == ["Ley N° 101"]
    assert visual["procedencia"] == {}, "el sistema no debe inventar una referencia"


def test_svg_normativo_contiene_texto_real_y_accesible():
    ley = visual_valida(
        titulo="Ley N° 101",
        proposito="Identificar la norma disciplinaria y sus dos alcances principales.",
        tipo="ley",
        conceptos=("faltas leves", "faltas graves"),
        texto_visible=("Ley N° 101", "Faltas leves", "Faltas graves"),
        texto_alternativo="Libro abierto con el título y los tipos de falta.",
    )
    svg = generar_svg(ley)
    ET.fromstring(svg)  # SVG bien formado
    assert "Ley N° 101" in svg
    assert "Faltas leves" in svg and "Faltas graves" in svg
    assert "<title" in svg and "<desc" in svg


def test_pptx_textual_sin_titulo_ni_visual_falla(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    ruta = tmp_path / "mal.pptx"
    p = Presentation()
    portada = p.slides.add_slide(p.slide_layouts[6])
    portada.shapes.add_textbox(Inches(1), Inches(0.4), Inches(5), Inches(1)).text = "Portada"
    slide = p.slides.add_slide(p.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1)).text = "Solo cuerpo"
    p.save(ruta)
    codigos = {h.codigo for h in auditar_presentacion(ruta)}
    assert "PPT-013" in codigos


def test_pptx_exige_texto_visible_de_ley(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    ruta = tmp_path / "ley.pptx"
    p = Presentation()
    slide = p.slides.add_slide(p.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(8), Inches(0.8)).text = "Ley N° 101"
    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(5), Inches(2))
    p.save(ruta)
    visual = visual_valida(
        diapositiva=1, titulo="Ley N° 101", tipo="ley",
        proposito="Identificar la ley y el régimen disciplinario que establece.",
        conceptos=("régimen disciplinario", "faltas"),
        texto_visible=("Ley N° 101", "Régimen disciplinario"),
    )
    hallazgos = auditar_presentacion(ruta, PlanVisual(1, "Ley", (visual,)))
    assert any(h.codigo == "PPT-021" and "Régimen" in h.mensaje for h in hallazgos)
